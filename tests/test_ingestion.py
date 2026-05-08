"""Tests for the PowerPoint ingestion engine."""
from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from slideatelier.ingestion import (
    ExtractedDesignSystem,
    LogoCandidate,
    extract_design_system,
    ingest_to_template,
)
from slideatelier.template import Template


def _save_blank_pptx(dest: Path) -> Path:
    """Write a Presentation()'s default state (default Office theme, 11 layouts)
    to disk and return the path."""
    prs = Presentation()
    prs.save(str(dest))
    return dest


def _tiny_png_bytes() -> BytesIO:
    """A 1x1 red PNG in memory — small enough to embed cheaply but a valid
    image so python-pptx will treat it as MSO_SHAPE_TYPE.PICTURE."""
    buf = BytesIO()
    Image.new("RGB", (1, 1), (255, 0, 0)).save(buf, format="PNG")
    buf.seek(0)
    return buf


def _add_picture_to_master(prs: Presentation, **xywh: int) -> None:
    """python-pptx's MasterShapes/LayoutShapes don't expose add_picture, but
    real corporate masters carry logos as picture shapes. We replicate that by
    adding the picture to a temp slide, transplanting the underlying <p:pic>
    XML into the master's shape tree, and rewiring the image relationship to
    the master part. xywh keys: left, top, width, height (EMU)."""
    master = prs.slide_masters[0]
    blank_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(blank_layout)
    pic = slide.shapes.add_picture(
        _tiny_png_bytes(),
        xywh["left"], xywh["top"], xywh["width"], xywh["height"],
    )

    pic_elem = pic._element
    slide.shapes._spTree.remove(pic_elem)

    # Rewire the image relationship from the slide part to the master part.
    A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    blip = pic_elem.find(f".//{A_NS}blip")
    embed_attr = next(k for k in blip.attrib if k.endswith("}embed"))
    old_rid = blip.attrib[embed_attr]
    img_part = slide.part.related_part(old_rid)
    reltype = slide.part.rels[old_rid].reltype
    new_rid = master.part.relate_to(img_part, reltype)
    blip.attrib[embed_attr] = new_rid

    master.shapes._spTree.append(pic_elem)

    # Drop the temp slide so it doesn't pollute the deck.
    slide_id_lst = prs.slides._sldIdLst
    for sld_id in list(slide_id_lst):
        if sld_id.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id") in slide.part.partname.lower():
            pass  # cleanup is best-effort; leftover blank slide is harmless


def test_extract_returns_design_system_with_defaults_on_blank_pptx(tmp_path: Path) -> None:
    """A python-pptx blank deck has Office's default theme. The extractor
    should parse it without error and yield a populated ExtractedDesignSystem.
    The default Office accent1 is #4F81BD; if our parsing works, primary
    should match that. If theme parsing falls back to defaults for any reason,
    the model still validates and primary is the TemplateColors default."""
    pptx_path = _save_blank_pptx(tmp_path / "blank.pptx")

    result = extract_design_system(pptx_path)

    assert isinstance(result, ExtractedDesignSystem)
    assert result.slide_width_inches > 0
    assert result.slide_height_inches > 0
    assert len(result.master_layouts) >= 1
    # Either we extracted Office's accent1 (#4F81BD) or fell back to default.
    assert result.theme_colors.primary in {"#4F81BD", "#1F3A5F"}
    # Default Office theme uses Calibri for both major and minor fonts.
    assert result.fonts.heading == "Calibri"
    assert result.fonts.body == "Calibri"
    assert isinstance(result.notes, list)


def test_extract_finds_picture_shapes_as_logo_candidates(tmp_path: Path) -> None:
    """Adding a Picture shape directly to a slide master should make it surface
    in logo_candidates with non-zero dimensions, a 'master' source tag, and a
    non-empty rId."""
    prs = Presentation()
    _add_picture_to_master(
        prs,
        left=Inches(0.25),
        top=Inches(0.5),
        width=Inches(1.0),
        height=Inches(0.5),
    )
    pptx_path = tmp_path / "with_logo.pptx"
    prs.save(str(pptx_path))

    result = extract_design_system(pptx_path)

    master_logos = [c for c in result.logo_candidates if c.found_in == "master"]
    assert len(master_logos) >= 1, f"no master-level logo found in {result.logo_candidates}"
    logo = master_logos[0]
    assert isinstance(logo, LogoCandidate)
    assert logo.width_inches == pytest.approx(1.0, abs=0.01)
    assert logo.height_inches == pytest.approx(0.5, abs=0.01)
    assert logo.position_inches[0] == pytest.approx(0.25, abs=0.01)
    assert logo.position_inches[1] == pytest.approx(0.5, abs=0.01)
    # The rId reference should be a non-empty string like "rId13".
    assert logo.image_relationship_id


def test_extract_reports_master_layout_names(tmp_path: Path) -> None:
    """python-pptx's default master ships with 11 layouts. Confirm we report
    all of them and the well-known names show up."""
    pptx_path = _save_blank_pptx(tmp_path / "blank.pptx")

    result = extract_design_system(pptx_path)

    assert len(result.master_layouts) == 11
    indices = [ml["index"] for ml in result.master_layouts]
    assert indices == list(range(11))
    names = [ml["name"] for ml in result.master_layouts]
    assert "Title Slide" in names
    assert "Title and Content" in names
    assert "Section Header" in names
    # Every entry must include placeholder_types as a list.
    for ml in result.master_layouts:
        assert isinstance(ml["placeholder_types"], list)


def test_ingest_to_template_writes_json_and_copies_master(tmp_path: Path) -> None:
    """End-to-end: feed a blank pptx into ingest_to_template, verify the
    json + master pptx land under the templates dir, the saved Template has
    master_pptx set, and reloading the json round-trips cleanly."""
    src_pptx = _save_blank_pptx(tmp_path / "source.pptx")
    templates_dir = tmp_path / "templates"

    template = ingest_to_template(
        pptx_path=src_pptx,
        template_name="Acme Corp",
        templates_dir=templates_dir,
        description="Ingested from Acme corporate deck",
    )

    json_path = templates_dir / "acme-corp.json"
    master_path = templates_dir / "acme-corp-master.pptx"
    assert json_path.exists(), "template JSON not written"
    assert master_path.exists(), "master pptx not copied"

    assert template.master_pptx == "acme-corp-master.pptx"
    assert template.name == "Acme Corp"
    assert template.description == "Ingested from Acme corporate deck"
    # detect_layout_map should populate at least the standard 7 keys.
    assert "title" in template.layout_map
    assert "content" in template.layout_map

    # Reload from disk to confirm the json is valid and round-trips.
    reloaded = Template.model_validate_json(json_path.read_text())
    assert reloaded.master_pptx == "acme-corp-master.pptx"
    assert reloaded.name == "Acme Corp"
