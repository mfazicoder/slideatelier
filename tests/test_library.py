"""Tests for the asset library system."""
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from slideatelier.asset_copier import copy_slide_into
from slideatelier.library import (
    LibraryAsset,
    LibraryCatalog,
    load_catalog,
    save_catalog,
    scan_library,
    slugify,
)


def test_slugify():
    assert slugify("Business Diagrams PowerPoint") == "business-diagrams-powerpoint"
    assert slugify("Funnels & Charts") == "funnels-charts"
    assert slugify("a/b\\c") == "abc"
    assert slugify("  spaces   here  ") == "spaces-here"


def _build_fake_library(root: Path) -> Path:
    """Create a minimal fake library with two categories and a slide-with-shapes."""
    cat_a = root / "Business Diagrams PowerPoint"
    cat_a.mkdir(parents=True)
    cat_b = root / "Funnels PowerPoint"
    cat_b.mkdir(parents=True)

    # First file: 2 slides, one with shapes
    prs1 = Presentation()
    prs1.slides.add_slide(prs1.slide_layouts[5])  # title only
    slide2 = prs1.slides.add_slide(prs1.slide_layouts[6])  # blank
    slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
    slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(1), Inches(2), Inches(1))
    box = slide2.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    box.text_frame.text = "Sample diagram label"
    prs1.save(cat_a / "Business-Diagrams-01.pptx")

    # Second file: 1 slide
    prs2 = Presentation()
    s = prs2.slides.add_slide(prs2.slide_layouts[6])
    s.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(1), Inches(1), Inches(2), Inches(1))
    prs2.save(cat_b / "Funnel-Infographic-01.pptx")

    return root


def test_scan_library_smoke(tmp_path: Path):
    root = _build_fake_library(tmp_path / "lib")
    catalog = scan_library(root)
    assert catalog.asset_count == 3  # 2 + 1
    assert len(catalog.categories) == 2
    ids = [a.id for a in catalog.assets]
    assert any("business-diagrams-powerpoint" in i for i in ids)
    assert any("funnels-powerpoint" in i for i in ids)


def test_scan_library_category_filter(tmp_path: Path):
    root = _build_fake_library(tmp_path / "lib")
    catalog = scan_library(root, categories=["funnel"])
    assert catalog.asset_count == 1
    assert "funnels-powerpoint" in catalog.assets[0].id


def test_scan_extracts_text_snippets(tmp_path: Path):
    root = _build_fake_library(tmp_path / "lib")
    catalog = scan_library(root)
    # Find the slide that has text
    diagrams = [a for a in catalog.assets if a.text_snippets]
    assert len(diagrams) >= 1
    assert "Sample diagram label" in diagrams[0].text_snippets[0]


def test_save_load_catalog_roundtrip(tmp_path: Path):
    root = _build_fake_library(tmp_path / "lib")
    catalog = scan_library(root)
    out = tmp_path / "catalog.json"
    save_catalog(catalog, out)

    loaded = load_catalog(out)
    assert loaded.asset_count == catalog.asset_count
    assert loaded.assets[0].id == catalog.assets[0].id


def test_catalog_find_by_id(tmp_path: Path):
    root = _build_fake_library(tmp_path / "lib")
    catalog = scan_library(root)
    asset_id = catalog.assets[0].id
    found = catalog.find(asset_id)
    assert found is not None
    assert found.id == asset_id
    assert catalog.find("nonexistent/x/1") is None


def test_copy_slide_preserves_native_shapes(tmp_path: Path):
    """Critical: copying a library slide must preserve shape primitives.
    A rectangle in the source must remain a rectangle in the destination —
    not a freeform path approximating one."""
    # Build a source file with known native shapes
    source_path = tmp_path / "source.pptx"
    src = Presentation()
    s = src.slides.add_slide(src.slide_layouts[6])
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
    oval = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(1), Inches(2), Inches(1))
    pentagon = s.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(1), Inches(3), Inches(2), Inches(1))
    src.save(source_path)
    # Capture source shape types for verification
    src_shape_types = [(sh.shape_type, sh.auto_shape_type if hasattr(sh, "auto_shape_type") and sh.has_text_frame else None)
                       for sh in src.slides[0].shapes]

    # Copy into a fresh dest
    dest = Presentation()
    copy_slide_into(source_path, 1, dest)

    # Inspect the new slide (last one in dest)
    new_slide = dest.slides[-1]
    shapes = list(new_slide.shapes)
    assert len(shapes) == 3

    # Each shape should still be its native auto_shape type, not a freeform path.
    # We verify by checking shape_type is AUTO_SHAPE (not FREEFORM=5).
    for sh in shapes:
        assert sh.shape_type is not None
        # AUTO_SHAPE type code is 1 in MSO_SHAPE_TYPE
        assert int(sh.shape_type) == 1, (
            f"Shape type {sh.shape_type} is not AUTO_SHAPE — native primitive lost during copy"
        )


def test_copy_slide_preserves_text(tmp_path: Path):
    source_path = tmp_path / "source.pptx"
    src = Presentation()
    s = src.slides.add_slide(src.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    box.text_frame.text = "Hello library"
    src.save(source_path)

    dest = Presentation()
    copy_slide_into(source_path, 1, dest)

    new_slide = dest.slides[-1]
    found_text = any(
        sh.has_text_frame and "Hello library" in sh.text_frame.text
        for sh in new_slide.shapes
    )
    assert found_text


def test_copy_slide_match_dimensions(tmp_path: Path):
    """When match_dimensions=True, dest dimensions update to source's."""
    source_path = tmp_path / "source.pptx"
    src = Presentation()
    src.slide_width = Inches(20)  # large
    src.slide_height = Inches(11.25)
    src.slides.add_slide(src.slide_layouts[6])
    src.save(source_path)

    dest = Presentation()
    assert dest.slide_width != src.slide_width
    copy_slide_into(source_path, 1, dest, match_dimensions=True)
    assert dest.slide_width == src.slide_width
    assert dest.slide_height == src.slide_height


def test_copy_slide_index_validation(tmp_path: Path):
    """Out-of-range slide index should raise IndexError clearly."""
    source_path = tmp_path / "source.pptx"
    Presentation().save(source_path)
    dest = Presentation()
    try:
        copy_slide_into(source_path, 99, dest)
        assert False, "should raise"
    except IndexError as e:
        assert "out of range" in str(e)


def test_copy_slide_scales_to_dest_dimensions(tmp_path: Path):
    """A 26.67"x15" library slide stamped onto a 13.33"x7.5" deck should be
    proportionally scaled so it fits the smaller canvas."""
    source_path = tmp_path / "source_big.pptx"
    src = Presentation()
    src.slide_width = Inches(26.67)
    src.slide_height = Inches(15)
    s = src.slides.add_slide(src.slide_layouts[6])
    s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(20), Inches(10), Inches(4), Inches(2)
    )
    src.save(source_path)

    dest = Presentation()
    dest.slide_width = Inches(13.33)
    dest.slide_height = Inches(7.5)

    copy_slide_into(source_path, 1, dest, scale_to_dest=True)

    new_slide = dest.slides[-1]
    shape = list(new_slide.shapes)[0]

    # scale_x = 13.33/26.67 ≈ 0.4998..., scale_y = 7.5/15 = 0.5
    expected_left = int(Inches(20) * (Inches(13.33) / Inches(26.67)))
    expected_top = int(Inches(10) * 0.5)
    expected_width = int(Inches(4) * (Inches(13.33) / Inches(26.67)))
    expected_height = int(Inches(2) * 0.5)

    # Allow ±1 EMU tolerance for floating-point rounding
    assert abs(shape.left - expected_left) <= 1
    assert abs(shape.top - expected_top) <= 1
    assert abs(shape.width - expected_width) <= 1
    assert abs(shape.height - expected_height) <= 1


def test_copy_slide_scale_disabled_preserves_coords(tmp_path: Path):
    """With scale_to_dest=False and match_dimensions=False, original
    absolute coordinates from the source slide are preserved verbatim."""
    source_path = tmp_path / "source_big.pptx"
    src = Presentation()
    src.slide_width = Inches(26.67)
    src.slide_height = Inches(15)
    s = src.slides.add_slide(src.slide_layouts[6])
    s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(20), Inches(10), Inches(4), Inches(2)
    )
    src.save(source_path)

    dest = Presentation()
    dest.slide_width = Inches(13.33)
    dest.slide_height = Inches(7.5)

    copy_slide_into(
        source_path, 1, dest, scale_to_dest=False, match_dimensions=False
    )

    new_slide = dest.slides[-1]
    shape = list(new_slide.shapes)[0]
    assert shape.left == Inches(20)
    assert shape.top == Inches(10)
    assert shape.width == Inches(4)
    assert shape.height == Inches(2)


def test_extra_with_bbox_and_color_override(tmp_path: Path):
    """Sprint D: SlideExtra config supports bbox (normalized rect) and
    color_override (#RRGGBB). The renderer must:
      - Position copied shapes within the bbox.
      - Recolor the dominant fill of the copied shapes to the override.
    """
    from pptx.dml.color import RGBColor
    from slideatelier.library import LibraryAsset, LibraryCatalog
    from slideatelier.models import Slide, SlideDeck, SlideExtra
    from slideatelier.renderer import DeckRenderer
    from slideatelier.template import Template

    # 1. Build a tiny library with a "blue diagram" — three blue rectangles
    #    plus one orange accent rectangle. Blue should dominate (3 vs 1).
    library_root = tmp_path / "lib"
    cat_dir = library_root / "Business Diagrams PowerPoint"
    cat_dir.mkdir(parents=True)
    BLUE = RGBColor(0x10, 0x20, 0xC0)
    ORANGE = RGBColor(0xC8, 0x6E, 0x3C)
    src_pptx = cat_dir / "Business-Diagrams-Recolor.pptx"
    src_pres = Presentation()
    src_slide = src_pres.slides.add_slide(src_pres.slide_layouts[6])
    for i in range(3):
        sh = src_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1 + 3 * i), Inches(1), Inches(2), Inches(1)
        )
        sh.fill.solid()
        sh.fill.fore_color.rgb = BLUE
    accent = src_slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1), Inches(3), Inches(2), Inches(1)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    src_pres.save(src_pptx)

    catalog = scan_library(library_root)
    catalog_path = tmp_path / "catalog.json"
    save_catalog(catalog, catalog_path)
    asset_id = catalog.assets[0].id

    # 2. Build a deck where one slide has an extra with bbox + color_override.
    bbox = {"left": 0.10, "top": 0.20, "width": 0.30, "height": 0.40}
    target_color_hex = "#ff0000"
    deck = SlideDeck(
        title="Recolor Test",
        core_message="Verify bbox + color_override flow into the renderer.",
        narrative_arc="Open with X. Defend with Y. Close with Z.",
        slides=[
            Slide(
                layout="content",
                title="Has a recolored extra",
                body=["one"],
                extras=[
                    SlideExtra(
                        type="library_asset",
                        position="right",
                        config={
                            "asset_ref": asset_id,
                            "bbox": bbox,
                            "color_override": target_color_hex,
                        },
                    )
                ],
            )
        ],
    )

    out = tmp_path / "deck.pptx"
    DeckRenderer(Template(), library_catalog_path=catalog_path).render(deck, out)
    assert out.exists() and out.stat().st_size > 0

    # 3. Read back the deck and verify position + recolor.
    rendered = Presentation(str(out))
    rendered_slide = rendered.slides[0]
    slide_w = rendered.slide_width
    slide_h = rendered.slide_height

    # The extra's shapes should land within the bbox rectangle (with a small
    # tolerance for scaling rounding).
    bbox_left = int(slide_w * bbox["left"])
    bbox_top = int(slide_h * bbox["top"])
    bbox_right = int(slide_w * (bbox["left"] + bbox["width"]))
    bbox_bottom = int(slide_h * (bbox["top"] + bbox["height"]))

    extra_shape_count = 0
    red_count = 0
    target_rgb = (0xFF, 0x00, 0x00)
    for shp in rendered_slide.shapes:
        try:
            sl, st, sw, sh_ = shp.left, shp.top, shp.width, shp.height
        except (AttributeError, TypeError):
            continue
        if sl is None or st is None or sw is None or sh_ is None:
            continue
        # Skip the body textbox (much wider than bbox)
        if sw > bbox_right - bbox_left + Inches(1):
            continue
        # Shapes for the extra should fit within the bbox (with small slack).
        slack = Inches(0.2)
        if sl >= bbox_left - slack and st >= bbox_top - slack and \
           sl + sw <= bbox_right + slack and st + sh_ <= bbox_bottom + slack:
            extra_shape_count += 1
            try:
                rgb = shp.fill.fore_color.rgb
                if rgb is not None and (rgb[0], rgb[1], rgb[2]) == target_rgb:
                    red_count += 1
            except Exception:  # noqa: BLE001
                pass

    assert extra_shape_count >= 3, \
        f"expected ≥3 extra shapes within bbox, got {extra_shape_count}"
    # The dominant blue fill should have been recolored to red. The orange
    # accent should remain (we only recolor the dominant color).
    assert red_count >= 1, \
        f"expected at least one shape recolored to {target_rgb}, got {red_count}"
