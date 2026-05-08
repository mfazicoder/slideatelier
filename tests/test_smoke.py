from pathlib import Path

from slideatelier.models import Slide, SlideDeck
from slideatelier.renderer import DeckRenderer, render_deck
from slideatelier.template import (
    Template,
    TemplateColors,
    is_valid_hex,
    parse_hex,
)

# Minimal valid deck values for tests that don't care about content quality.
# The real deck production fills these via Claude.
_CORE = "Test core message stating the answer in one sentence."
_ARC = "Open with X. Defend with Y. Close with Z."


def _deck(**overrides) -> SlideDeck:
    """Build a SlideDeck with sensible test defaults; override any field."""
    base = dict(title="Test", core_message=_CORE, narrative_arc=_ARC,
                slides=[Slide(layout="title", title="Test")])
    base.update(overrides)
    return SlideDeck(**base)


def test_models_validate():
    deck = _deck(slides=[Slide(layout="title", title="Hello")])
    assert deck.slides[0].title == "Hello"
    assert deck.core_message == _CORE


def test_models_require_core_message():
    """core_message and narrative_arc must be provided — they're not optional.
    If a test starts failing here, the schema added a required field; update _deck()."""
    try:
        SlideDeck(title="X", slides=[Slide(layout="title", title="X")])
        assert False, "SlideDeck should require core_message and narrative_arc"
    except Exception:
        pass


def test_hex_validation():
    assert is_valid_hex("#1F3A5F")
    assert is_valid_hex("1f3a5f")
    assert not is_valid_hex("blue")
    assert not is_valid_hex("#12345")
    assert parse_hex("1f3a5f") == "#1F3A5F"
    assert parse_hex("#1F3A5F") == "#1F3A5F"


def test_template_defaults():
    tpl = Template()
    assert tpl.colors.primary == "#1F3A5F"
    assert tpl.fonts.heading == "Calibri"
    assert tpl.fonts.sizes.slide_title == 24


def test_template_custom_colors():
    tpl = Template(colors=TemplateColors(primary="ff0000", accent="00ff00"))
    assert tpl.colors.primary == "#FF0000"
    assert tpl.colors.accent == "#00FF00"


def test_renderer_smoke(tmp_path: Path):
    deck = _deck(
        title="Smoke Test",
        subtitle="A trivial deck for verifying the renderer",
        slides=[
            Slide(layout="title", title="Smoke Test"),
            Slide(
                layout="content",
                title="Three things to know",
                body=["First insight", "Second insight", "Third insight"],
            ),
            Slide(
                layout="two_column",
                title="Before vs After",
                body_left=["Old approach", "Slow"],
                body_right=["New approach", "Fast"],
            ),
            Slide(
                layout="key_takeaway",
                title="The headline result",
                body=["Backed by Q3 data"],
            ),
            Slide(layout="closing", title="Next steps", body=["Decide", "Execute"]),
        ],
    )
    out = tmp_path / "deck.pptx"
    render_deck(deck, out)
    assert out.exists()
    assert out.stat().st_size > 0


def test_renderer_honors_block_bbox(tmp_path: Path):
    """Sprint Y.2: when a slide has block_bbox set, the renderer must place
    that block's text frame at the bbox-derived absolute rect (in EMU),
    not the layout's hardcoded default. Round-trip: set bbox → render →
    inspect resulting .pptx → confirm at least one shape is positioned
    inside the bbox-derived rect.

    The wireframe stage stores bboxes as 0..1 normalised over the slide
    area; the renderer projects onto slide_width × slide_height in EMU.
    """
    from pptx import Presentation as _Presentation
    deck = _deck(
        title="BBox Test",
        slides=[
            # Auto-promote scenario: a `title` layout with body content lifted
            # into a freeform overlay. Default `title` rendering would drop
            # body entirely; with bbox set, body MUST appear at bbox position.
            Slide(
                layout="title",
                title="Cover with body",
                body=["Bullet A", "Bullet B", "Bullet C"],
                block_bbox={
                    "body": {"left": 0.07, "top": 0.72, "width": 0.86, "height": 0.22},
                },
            ),
            # Content slide with title moved to bottom-right corner.
            Slide(
                layout="content",
                title="Title in unusual spot",
                body=["regular body item"],
                block_bbox={
                    "title": {"left": 0.55, "top": 0.80, "width": 0.40, "height": 0.10},
                },
            ),
        ],
    )
    out = tmp_path / "deck.pptx"
    render_deck(deck, out)
    assert out.exists() and out.stat().st_size > 0

    # Re-open and inspect.
    prs = _Presentation(str(out))
    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height

    # --- Slide 0: body should appear at the bbox-derived rect ---
    s0 = prs.slides[0]
    expected_left = int(0.07 * slide_w_emu)
    expected_top  = int(0.72 * slide_h_emu)
    body_shapes = [
        sh for sh in s0.shapes
        if sh.has_text_frame and "Bullet A" in sh.text_frame.text
    ]
    assert body_shapes, "body bullets must be rendered on title-layout slide when block_bbox['body'] is set"
    placed = body_shapes[0]
    # Allow rounding slop of ±1% of slide dim.
    slop_w = int(0.01 * slide_w_emu)
    slop_h = int(0.01 * slide_h_emu)
    assert abs(placed.left - expected_left) <= slop_w, (
        f"body left was {placed.left} EMU, expected ~{expected_left} EMU "
        f"(bbox.left=0.07 × slide_w={slide_w_emu})"
    )
    assert abs(placed.top - expected_top) <= slop_h, (
        f"body top was {placed.top} EMU, expected ~{expected_top} EMU "
        f"(bbox.top=0.72 × slide_h={slide_h_emu})"
    )

    # --- Slide 1: title should be at its bbox, NOT the default Inches(0.6, 0.5) ---
    s1 = prs.slides[1]
    expected_left = int(0.55 * slide_w_emu)
    expected_top  = int(0.80 * slide_h_emu)
    title_shapes = [
        sh for sh in s1.shapes
        if sh.has_text_frame and "Title in unusual spot" in sh.text_frame.text
    ]
    assert title_shapes, "title text must be rendered"
    placed = title_shapes[0]
    assert abs(placed.left - expected_left) <= int(0.01 * slide_w_emu)
    assert abs(placed.top - expected_top) <= int(0.01 * slide_h_emu)


def test_renderer_honors_block_style(tmp_path: Path):
    """Sprint Y.3: when slide.block_style[name] is set, the rendered shape's
    font properties (size/family/color/bold/italic) and paragraph alignment
    must reflect the overrides — not just the layout defaults.
    """
    from pptx import Presentation as _Presentation
    from pptx.util import Pt as _Pt
    from pptx.enum.text import PP_ALIGN as _PP_ALIGN

    deck = _deck(
        title="Style Test",
        slides=[
            Slide(
                layout="content",
                title="Bold red 32-pt centered title",
                body=["italic body"],
                block_style={
                    "title": {
                        "font_family": "Georgia",
                        "font_size": 32,
                        "color": "#cc0000",
                        "bold": True,
                        "italic": False,
                        "align": "center",
                    },
                    "body": {
                        "italic": True,
                        "color": "#333333",
                        "align": "right",
                    },
                },
            ),
        ],
    )
    out = tmp_path / "style.pptx"
    render_deck(deck, out)
    prs = _Presentation(str(out))
    s = prs.slides[0]

    # Find the title shape.
    title_shapes = [sh for sh in s.shapes
                    if sh.has_text_frame and "Bold red" in sh.text_frame.text]
    assert title_shapes
    tp = title_shapes[0].text_frame.paragraphs[0]
    title_run = tp.runs[0]
    assert tp.alignment == _PP_ALIGN.CENTER, "title align should be CENTER"
    assert title_run.font.size == _Pt(32), f"title size should be 32pt, got {title_run.font.size}"
    assert title_run.font.bold is True
    assert title_run.font.name == "Georgia"
    # Color check: python-pptx stores RGBColor; compare hex.
    assert title_run.font.color.rgb is not None
    assert str(title_run.font.color.rgb).lower() == "cc0000"

    # Find the body shape (bullets — paragraph styling, not run styling).
    body_shapes = [sh for sh in s.shapes
                   if sh.has_text_frame and "italic body" in sh.text_frame.text]
    assert body_shapes
    bp = body_shapes[0].text_frame.paragraphs[0]
    assert bp.alignment == _PP_ALIGN.RIGHT, "body align should be RIGHT"
    assert bp.font.italic is True
    assert str(bp.font.color.rgb).lower() == "333333"


def test_renderer_with_custom_template(tmp_path: Path):
    """Custom template with bright red — verifies template flows through to renderer."""
    tpl = Template(colors=TemplateColors(primary="#FF0000", accent="#00FF00"))
    deck = _deck(
        title="Custom Brand",
        slides=[Slide(layout="title", title="Custom"), Slide(layout="content", title="Test", body=["x"])],
    )
    out = tmp_path / "deck.pptx"
    DeckRenderer(tpl).render(deck, out)
    assert out.exists() and out.stat().st_size > 0


def test_master_mode_render(tmp_path: Path):
    """Save a blank Presentation as a master, then render through it."""
    from pptx import Presentation as PptxPresentation
    master_path = tmp_path / "master.pptx"
    PptxPresentation().save(str(master_path))

    tpl = Template(name="MasterTest", master_pptx=str(master_path))
    deck = _deck(
        title="Master Test",
        subtitle="Inheriting layouts",
        slides=[
            Slide(layout="title", title="Master Test"),
            Slide(layout="section_divider", title="Section One"),
            Slide(layout="content", title="Some content", body=["Point A", "Point B"]),
            Slide(layout="two_column", title="Compare", body_left=["L1", "L2"], body_right=["R1", "R2"]),
            Slide(layout="key_takeaway", title="The big idea"),
            Slide(layout="closing", title="Next steps", body=["Do thing"]),
        ],
    )
    out = tmp_path / "deck.pptx"
    DeckRenderer(tpl, templates_dir=tmp_path).render(deck, out)
    assert out.exists() and out.stat().st_size > 0


def test_master_mode_missing_file_raises(tmp_path: Path):
    tpl = Template(name="X", master_pptx=str(tmp_path / "does-not-exist.pptx"))
    deck = _deck(title="X", slides=[Slide(layout="title", title="X")])
    out = tmp_path / "deck.pptx"
    try:
        DeckRenderer(tpl, templates_dir=tmp_path).render(deck, out)
        assert False, "should have raised"
    except FileNotFoundError:
        pass


def test_cache_input_hash_is_deterministic_and_sensitive():
    from slideatelier.cache import compute_input_hash

    base = dict(
        system_prompt="prompt",
        brief="brief",
        requirements="reqs",
        model_id="claude-opus-4-7",
        prompt_version="v1.0.0",
    )
    h1 = compute_input_hash(**base)
    h2 = compute_input_hash(**base)
    assert h1 == h2  # deterministic

    # Each input field affects the hash
    for field in ("system_prompt", "brief", "requirements", "model_id", "prompt_version"):
        modified = dict(base)
        modified[field] = base[field] + "X"
        assert compute_input_hash(**modified) != h1, f"hash insensitive to {field}"


def test_cache_input_hash_no_collision_on_concatenation():
    """('ab', 'cd') and ('a', 'bcd') must produce different hashes — that's
    why we use a delimiter byte in the hash."""
    from slideatelier.cache import compute_input_hash

    h1 = compute_input_hash(system_prompt="ab", brief="cd", requirements="", model_id="m", prompt_version="v")
    h2 = compute_input_hash(system_prompt="a", brief="bcd", requirements="", model_id="m", prompt_version="v")
    assert h1 != h2


def test_cache_get_set_clear(tmp_path: Path):
    from slideatelier.cache import DeckCache

    cache = DeckCache(tmp_path / "cache", enabled=True)
    deck = _deck(title="Cached", slides=[Slide(layout="title", title="Cached")])
    cache.set("abc123", deck)

    retrieved = cache.get("abc123")
    assert retrieved is not None
    assert retrieved.title == "Cached"

    assert cache.get("nonexistent") is None

    entries = cache.list_entries()
    assert len(entries) == 1
    assert entries[0]["key"] == "abc123"
    assert entries[0]["title"] == "Cached"

    n = cache.clear()
    assert n == 1
    assert cache.get("abc123") is None


def test_cache_disabled_is_noop(tmp_path: Path):
    from slideatelier.cache import DeckCache

    cache = DeckCache(tmp_path / "cache", enabled=False)
    deck = _deck(title="X", slides=[Slide(layout="title", title="X")])
    cache.set("k", deck)
    assert cache.get("k") is None


def test_generated_deck_wrapper_serializes():
    from datetime import datetime, timezone

    from slideatelier.metadata import GeneratedDeck, GenerationMetadata

    wrapped = GeneratedDeck(
        metadata=GenerationMetadata(
            model_id="claude-opus-4-7",
            prompt_version="v1.0.0",
            generated_at=datetime(2026, 5, 6, tzinfo=timezone.utc),
            cache_hit=False,
            input_hash="abc",
            duration_seconds=12.3,
        ),
        deck=_deck(title="T", slides=[Slide(layout="title", title="T")]),
    )
    json_str = wrapped.model_dump_json()
    parsed = GeneratedDeck.model_validate_json(json_str)
    assert parsed.metadata.input_hash == "abc"
    assert parsed.deck.title == "T"


def test_prompt_version_current():
    """Sanity check — verify the current prompt version. Update this when bumping
    PROMPT_VERSION so it stays in sync with the actual code."""
    from slideatelier.metadata import PROMPT_VERSION
    assert PROMPT_VERSION == "v1.5.0", f"Expected v1.5.0, got {PROMPT_VERSION}"


def test_slide_extras_schema():
    """SlideExtra is the foundation for growable add-ons (charts, library assets, etc)."""
    from slideatelier.models import Slide, SlideExtra
    extra = SlideExtra(type="library_asset", position="right",
                       config={"asset_ref": "funnels-powerpoint/funnel-infographic-01/2"})
    s = Slide(layout="content", title="Test", extras=[extra])
    assert len(s.extras) == 1
    assert s.extras[0].config["asset_ref"] == "funnels-powerpoint/funnel-infographic-01/2"
    assert s.extras[0].position == "right"

    # Default empty extras
    s2 = Slide(layout="content", title="Plain")
    assert s2.extras == []


def test_slide_has_optional_strap_field():
    """Strap field defaults to '' so existing decks/specs validate without it."""
    from slideatelier.models import Slide
    s = Slide(layout="content", title="Test")
    assert s.strap == ""

    s2 = Slide(layout="content", title="Test", strap="Subtitle line")
    assert s2.strap == "Subtitle line"


def test_storyboard_schema_validates():
    """Storyboard is the Stage 1 primitive — must require core_message and narrative_arc."""
    from slideatelier.storyboard import Storyboard, StoryboardSlide

    sb = Storyboard(
        title="Test deck",
        subtitle="board",
        core_message="The single sentence answer.",
        narrative_arc="Open with X. Defend with Y. Close with Z.",
        slides=[
            StoryboardSlide(layout="title", title="", purpose="Opens the deck"),
            StoryboardSlide(layout="content", title="Insight-led title", purpose="Defends the thesis"),
        ],
    )
    assert len(sb.slides) == 2
    assert sb.slides[1].purpose == "Defends the thesis"


def test_storyboard_requires_core_message():
    from slideatelier.storyboard import Storyboard
    try:
        Storyboard(title="X", narrative_arc="X", slides=[])
        assert False, "core_message must be required"
    except Exception:
        pass


def test_critic_models_validate():
    """The critique stage's schemas must validate."""
    from slideatelier.critic import CritiqueResult, DeckCritique, SlideCritique

    dc = DeckCritique(
        core_message_assessment="Good",
        narrative_arc_assessment="Holds",
        weakest_slide_index=2,
        missing_slide=None,
        overall_verdict="Client-ready.",
    )
    sc = SlideCritique(
        slide_index=0,
        layout="title",
        current_title="Test",
        severity="ok",
        verdict="Works fine.",
    )
    cr = CritiqueResult(
        critic_version="v1.0.0",
        model_id="claude-opus-4-7",
        duration_seconds=5.0,
        slide_critiques=[sc],
        deck_critique=dc,
    )
    assert cr.fatal_count == 0
    assert cr.is_client_ready is True


def test_critic_severity_aggregates():
    from slideatelier.critic import CritiqueResult, DeckCritique, SlideCritique

    sc1 = SlideCritique(slide_index=0, layout="title", current_title="A", severity="ok", verdict="OK")
    sc2 = SlideCritique(slide_index=1, layout="content", current_title="B", severity="major",
                        verdict="bad", suggestion="rewrite")
    sc3 = SlideCritique(slide_index=2, layout="content", current_title="C", severity="fatal",
                        verdict="terrible", suggestion="redo")
    cr = CritiqueResult(
        critic_version="v1.0.0",
        model_id="claude-opus-4-7",
        duration_seconds=1.0,
        slide_critiques=[sc1, sc2, sc3],
        deck_critique=DeckCritique(
            core_message_assessment="weak", narrative_arc_assessment="broken",
            weakest_slide_index=2, missing_slide=None, overall_verdict="not ready",
        ),
    )
    assert cr.fatal_count == 1
    assert cr.major_count == 1
    assert cr.minor_count == 0
    assert cr.ok_count == 1
    assert cr.is_client_ready is False


def test_apply_critique_swaps_titles():
    """apply_critique_to_deck should swap proposed_title for fatal/major slides."""
    from slideatelier.critic import (
        CritiqueResult,
        DeckCritique,
        SlideCritique,
        apply_critique_to_deck,
    )

    deck = _deck(
        title="Deck",
        slides=[
            Slide(layout="title", title="Test"),
            Slide(layout="content", title="Boring topic", body=["a"]),
            Slide(layout="content", title="Another boring one", body=["b"]),
        ],
    )
    critique = CritiqueResult(
        critic_version="v1.0.0",
        model_id="claude-opus-4-7",
        duration_seconds=1.0,
        slide_critiques=[
            SlideCritique(slide_index=0, layout="title", current_title="Test", severity="ok", verdict="ok"),
            SlideCritique(
                slide_index=1, layout="content", current_title="Boring topic",
                severity="major", verdict="weak title", suggestion="rewrite",
                proposed_title="Insight-led replacement",
            ),
            SlideCritique(
                slide_index=2, layout="content", current_title="Another boring one",
                severity="minor", verdict="meh", suggestion="polish",
                proposed_title="Slightly better",
            ),
        ],
        deck_critique=DeckCritique(
            core_message_assessment="ok", narrative_arc_assessment="ok",
            weakest_slide_index=1, missing_slide=None, overall_verdict="ok with fixes",
        ),
    )
    revised, changed = apply_critique_to_deck(deck, critique, accept_severities=("fatal", "major"))
    assert changed == [1]
    assert revised.slides[0].title == "Test"  # untouched
    assert revised.slides[1].title == "Insight-led replacement"  # swapped (major)
    assert revised.slides[2].title == "Another boring one"  # untouched (minor not in accept list)
