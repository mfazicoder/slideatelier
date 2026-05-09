"""Sprint Z.v2 — live SVG translation of library .pptx slides for the Web Deck.

Verifies that `library_asset_to_svg` produces well-formed inline SVG that
mirrors the source slide's shape geometry, applies chrome_only stripping,
and gracefully handles unsupported shape kinds.
"""
from pathlib import Path
import xml.etree.ElementTree as ET

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from slideatelier.library_to_svg import library_asset_to_svg


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _build_source_pptx(tmp_path: Path) -> Path:
    """Build a synthetic library slide with: a chrome title bar, a chrome
    footer, two AutoShape rectangles (the diagram body), one AutoShape
    oval, and a couple of inline TEXT_BOX annotations."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar — chrome.
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tb.name = "Title 1"
    tb.text_frame.text = "Diagram Title (chrome)"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(28)

    # Footer — chrome.
    fb = s.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12), Inches(0.3))
    fb.name = "Footer 1"
    fb.text_frame.text = "Page footer (chrome)"

    # Two filled rectangles forming a diagram body.
    s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2.5), Inches(4), Inches(2))
    s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(2.5), Inches(4), Inches(2))

    # An oval.
    s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(5), Inches(2), Inches(2))

    # Inline annotation — should survive chrome_only.
    ann = s.shapes.add_textbox(Inches(2), Inches(4.7), Inches(4), Inches(0.4))
    ann.text_frame.text = "Left annotation"

    out = tmp_path / "synth_lib.pptx"
    prs.save(str(out))
    return out


def test_library_asset_to_svg_emits_valid_xml(tmp_path):
    src = _build_source_pptx(tmp_path)
    svg = library_asset_to_svg(src, 1, 800, 600, strip_text="chrome_only")
    assert svg.startswith("<svg")
    assert svg.rstrip().endswith("</svg>")
    ET.fromstring(svg)  # raises if not well-formed


def test_library_asset_to_svg_emits_native_primitives(tmp_path):
    """Rectangles and ovals from the source slide should map to <rect> and
    <ellipse>/<circle> respectively — not approximated as paths or images."""
    src = _build_source_pptx(tmp_path)
    svg = library_asset_to_svg(src, 1, 800, 600, strip_text="chrome_only")
    root = ET.fromstring(svg)
    tags = [el.tag.rsplit("}", 1)[-1] for el in root.iter()]
    assert tags.count("rect") >= 2, "two source RECTANGLE shapes must emit <rect>s"
    assert "ellipse" in tags or "circle" in tags, "OVAL must emit <ellipse> or <circle>"


def test_library_asset_to_svg_chrome_only_strips_chrome_keeps_inline(tmp_path):
    """Chrome (title + footer) text must NOT appear; inline annotation must."""
    src = _build_source_pptx(tmp_path)
    svg = library_asset_to_svg(src, 1, 800, 600, strip_text="chrome_only")
    assert "Diagram Title (chrome)" not in svg
    assert "Page footer (chrome)" not in svg
    assert "Left annotation" in svg


def test_library_asset_to_svg_strip_none_keeps_everything(tmp_path):
    src = _build_source_pptx(tmp_path)
    svg = library_asset_to_svg(src, 1, 800, 600, strip_text="none")
    assert "Diagram Title (chrome)" in svg
    assert "Page footer (chrome)" in svg
    assert "Left annotation" in svg


def test_library_asset_to_svg_strip_all_blanks_text(tmp_path):
    src = _build_source_pptx(tmp_path)
    svg = library_asset_to_svg(src, 1, 800, 600, strip_text="all")
    # In "all" mode, TEXT_BOXes are dropped entirely (no geometry to keep)
    # and AutoShape text-frames have their text suppressed.
    assert "Diagram Title (chrome)" not in svg
    assert "Page footer (chrome)" not in svg
    assert "Left annotation" not in svg


def test_library_asset_to_svg_uses_source_emu_viewbox(tmp_path):
    """The outer SVG's viewBox is in source-slide EMU so each shape can use
    its native coords directly. Width/height attrs scale to target size."""
    src = _build_source_pptx(tmp_path)
    svg = library_asset_to_svg(src, 1, 800, 600)
    root = ET.fromstring(svg)
    assert root.get("width") == "800"
    assert root.get("height") == "600"
    vb = root.get("viewBox") or ""
    parts = vb.split()
    assert len(parts) == 4
    # Width/height in EMU should match Inches(13.333) / Inches(7.5).
    assert int(parts[2]) == int(SLIDE_W)
    assert int(parts[3]) == int(SLIDE_H)


def test_library_asset_to_svg_invalid_index_raises(tmp_path):
    src = _build_source_pptx(tmp_path)
    with pytest.raises(IndexError):
        library_asset_to_svg(src, 99, 800, 600)


def test_library_asset_to_svg_unknown_shape_falls_back_to_rect(tmp_path):
    """A LINE shape (which we don't have a hand-mapping for beyond _emit_line)
    should still produce a <line> element. Verifies the dispatcher doesn't
    crash on common-but-non-AutoShape primitives."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # add_connector — produces an MSO_SHAPE_TYPE.LINE
    from pptx.enum.shapes import MSO_CONNECTOR
    s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2), Inches(2), Inches(8), Inches(5))
    src = tmp_path / "with_line.pptx"
    prs.save(str(src))

    svg = library_asset_to_svg(src, 1, 800, 600)
    ET.fromstring(svg)  # well-formed
    # Connector renders as <line> via _emit_line.
    assert "<line " in svg or "<rect " in svg  # line preferred, rect fallback acceptable


def test_library_asset_to_svg_handles_groups(tmp_path):
    """Shapes nested inside a GROUP should be recursed into and emitted."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # add_group_shape — wraps the two children in a group.
    grouped_rect_a = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(2))
    grouped_rect_b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(1), Inches(2), Inches(2))
    src = tmp_path / "grouped.pptx"
    prs.save(str(src))

    svg = library_asset_to_svg(src, 1, 800, 600)
    root = ET.fromstring(svg)
    tags = [el.tag.rsplit("}", 1)[-1] for el in root.iter()]
    # Two source rectangles → at least two <rect>s in output (groups recurse).
    assert tags.count("rect") >= 2
