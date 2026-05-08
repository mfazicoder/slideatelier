"""PowerPoint ingestion engine.

Extracts a corporate `.pptx`'s design system (theme colors, fonts, master
layouts, logos, footer, sample palette) into slideAtelier's Template format
so the user can render new decks against an inherited corporate look.

Pure parsing — no LLM calls.
"""
from __future__ import annotations

import re
import shutil
import xml.etree.ElementTree as ET
import zipfile
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pydantic import BaseModel

from .template import (
    Template,
    TemplateColors,
    TemplateFonts,
    TemplateFontSizes,
    save_template,
)
from .template_inspector import detect_layout_map

# DrawingML namespace used throughout theme/slide XML
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS = {"a": A_NS}

# theme slot -> our TemplateColors field
THEME_COLOR_MAP: dict[str, str] = {
    "dk1": "text",
    "lt1": "background",
    "accent1": "primary",
    "accent2": "accent",
    "accent3": "success",
    "accent4": "warning",
    "accent5": "danger",
    "dk2": "muted",
}

HEX6_RE = re.compile(r"^[0-9A-Fa-f]{6}$")


class LogoCandidate(BaseModel):
    image_relationship_id: str
    width_inches: float
    height_inches: float
    position_inches: tuple[float, float]
    found_in: str  # "master" | "layout:0" | "layout:1" etc.


class ExtractedDesignSystem(BaseModel):
    theme_colors: TemplateColors
    fonts: TemplateFonts
    master_layouts: list[dict]
    logo_candidates: list[LogoCandidate]
    footer_text: str | None
    sample_palette: list[str]
    slide_width_inches: float
    slide_height_inches: float
    notes: list[str]


# ---------- helpers --------------------------------------------------------

def _slugify(value: str) -> str:
    s = re.sub(r"[^a-zA-Z0-9]+", "-", value.strip().lower()).strip("-")
    return s or "template"


def _emu_to_inches(emu: int | None) -> float:
    if emu is None:
        return 0.0
    return float(emu) / 914400.0


def _hex_or_none(value: str | None) -> str | None:
    """Accept a 6-char hex string and return '#RRGGBB' uppercase, else None."""
    if value is None:
        return None
    v = value.strip().lstrip("#")
    if HEX6_RE.match(v):
        return "#" + v.upper()
    return None


def _read_theme_xml(pptx_path: Path) -> bytes | None:
    """Locate and read the first theme XML from the .pptx zip. Returns None if
    no theme file exists or the file is not a valid zip."""
    try:
        with zipfile.ZipFile(pptx_path) as zf:
            theme_names = sorted(
                n for n in zf.namelist() if n.startswith("ppt/theme/") and n.endswith(".xml")
            )
            if not theme_names:
                return None
            return zf.read(theme_names[0])
    except (zipfile.BadZipFile, KeyError):
        return None


def _extract_color_from_slot(slot_elem: ET.Element) -> str | None:
    """A theme color slot like <a:dk1> can wrap either an <a:srgbClr> (direct
    hex) or an <a:sysClr> (system color, with optional lastClr hex fallback).
    Return '#RRGGBB' or None if neither shape is parseable."""
    srgb = slot_elem.find("a:srgbClr", NS)
    if srgb is not None:
        v = srgb.get("val")
        h = _hex_or_none(v)
        if h:
            return h
    sys_clr = slot_elem.find("a:sysClr", NS)
    if sys_clr is not None:
        # sysClr names like "windowText" don't give us a hex directly; the
        # `lastClr` attribute caches the resolved value for non-Office viewers.
        last = sys_clr.get("lastClr")
        h = _hex_or_none(last)
        if h:
            return h
    return None


def _parse_theme(
    theme_bytes: bytes, notes: list[str]
) -> tuple[TemplateColors, TemplateFonts]:
    """Parse a theme1.xml blob into TemplateColors + TemplateFonts. Any slot
    that fails to resolve is left at the model default and a note is recorded."""
    color_overrides: dict[str, str] = {}
    font_overrides: dict[str, str] = {}

    try:
        root = ET.fromstring(theme_bytes)
    except ET.ParseError as e:
        notes.append(f"theme XML parse error: {e}; using default colors and fonts")
        return TemplateColors(), TemplateFonts()

    elements = root.find(".//a:themeElements", NS)
    if elements is None:
        notes.append("theme XML missing <a:themeElements>; using defaults")
        return TemplateColors(), TemplateFonts()

    clr_scheme = elements.find("a:clrScheme", NS)
    if clr_scheme is None:
        notes.append("theme XML missing <a:clrScheme>; using default colors")
    else:
        for theme_slot, our_field in THEME_COLOR_MAP.items():
            slot_elem = clr_scheme.find(f"a:{theme_slot}", NS)
            if slot_elem is None:
                notes.append(f"theme slot <a:{theme_slot}> missing; defaulting {our_field}")
                continue
            hex_val = _extract_color_from_slot(slot_elem)
            if hex_val is None:
                notes.append(
                    f"theme slot <a:{theme_slot}> had no resolvable hex; defaulting {our_field}"
                )
                continue
            color_overrides[our_field] = hex_val

    font_scheme = elements.find("a:fontScheme", NS)
    if font_scheme is None:
        notes.append("theme XML missing <a:fontScheme>; using default fonts")
    else:
        major = font_scheme.find("a:majorFont/a:latin", NS)
        minor = font_scheme.find("a:minorFont/a:latin", NS)
        if major is not None and major.get("typeface"):
            font_overrides["heading"] = major.get("typeface", "")
        if minor is not None and minor.get("typeface"):
            font_overrides["body"] = minor.get("typeface", "")

    colors = TemplateColors(**color_overrides) if color_overrides else TemplateColors()
    fonts = TemplateFonts(
        sizes=TemplateFontSizes(),
        **{k: v for k, v in font_overrides.items() if v},
    )
    return colors, fonts


def _walk_master_pictures(prs: Presentation) -> list[LogoCandidate]:
    """Find Picture shapes on slide masters and master-level layouts. Likely
    logos. Each shape's rId is recorded so the orchestrator can extract the
    image bytes later if it wants to copy the logo into the output folder."""
    candidates: list[LogoCandidate] = []

    for master in prs.slide_masters:
        for shape in master.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                candidates.append(
                    LogoCandidate(
                        image_relationship_id=_picture_rid(shape),
                        width_inches=_emu_to_inches(shape.width),
                        height_inches=_emu_to_inches(shape.height),
                        position_inches=(
                            _emu_to_inches(shape.left),
                            _emu_to_inches(shape.top),
                        ),
                        found_in="master",
                    )
                )
        for i, layout in enumerate(master.slide_layouts):
            for shape in layout.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    candidates.append(
                        LogoCandidate(
                            image_relationship_id=_picture_rid(shape),
                            width_inches=_emu_to_inches(shape.width),
                            height_inches=_emu_to_inches(shape.height),
                            position_inches=(
                                _emu_to_inches(shape.left),
                                _emu_to_inches(shape.top),
                            ),
                            found_in=f"layout:{i}",
                        )
                    )

    return candidates


def _picture_rid(shape) -> str:
    """python-pptx exposes the embed rId via shape.image.blob, but the
    relationship id itself sits on the underlying <p:blipFill><a:blip r:embed=..>
    XML. Pull it directly so we keep the original reference."""
    try:
        blip = shape._element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
        )
        if blip is not None:
            for k, v in blip.attrib.items():
                if k.endswith("}embed") or k == "embed" or k.endswith("}link") or k == "link":
                    return v
    except Exception:
        pass
    return ""


def _extract_master_footer(prs: Presentation) -> str | None:
    """Walk masters and master-layout placeholders. Return the first non-empty
    FOOTER placeholder text, or None."""
    for master in prs.slide_masters:
        for ph in master.placeholders:
            if ph.placeholder_format.type == PP_PLACEHOLDER.FOOTER:
                txt = (ph.text_frame.text or "").strip()
                if txt:
                    return txt
        for layout in master.slide_layouts:
            for ph in layout.placeholders:
                if ph.placeholder_format.type == PP_PLACEHOLDER.FOOTER:
                    txt = (ph.text_frame.text or "").strip()
                    if txt:
                        return txt
    return None


def _master_layout_summary(prs: Presentation) -> list[dict]:
    """[{'index': 0, 'name': 'Title Slide', 'placeholder_types': [...]}, ...]
    over every layout in every master. Index resets per master, but for
    single-master decks (the common case) it lines up with prs.slide_layouts."""
    out: list[dict] = []
    for master in prs.slide_masters:
        for i, layout in enumerate(master.slide_layouts):
            ph_types: list[str] = []
            for ph in layout.placeholders:
                t = ph.placeholder_format.type
                ph_types.append(t.name if hasattr(t, "name") else str(t))
            out.append({"index": i, "name": layout.name, "placeholder_types": ph_types})
    return out


def _sample_palette_from_slides(pptx_path: Path, max_slides: int = 10, top_n: int = 8) -> list[str]:
    """Walk the first N content-slide XML files and count any <a:srgbClr val=..>
    occurrences (shape fills, line colors, run colors). Return up to top_n
    distinct hex values ranked by frequency."""
    counter: Counter[str] = Counter()
    try:
        with zipfile.ZipFile(pptx_path) as zf:
            slide_files = sorted(
                n for n in zf.namelist()
                if n.startswith("ppt/slides/slide") and n.endswith(".xml")
            )
            for name in slide_files[:max_slides]:
                blob = zf.read(name)
                try:
                    root = ET.fromstring(blob)
                except ET.ParseError:
                    continue
                for srgb in root.iter(f"{{{A_NS}}}srgbClr"):
                    v = srgb.get("val")
                    h = _hex_or_none(v)
                    if h:
                        counter[h] += 1
    except (zipfile.BadZipFile, KeyError):
        return []

    return [h for h, _ in counter.most_common(top_n)]


# ---------- public API -----------------------------------------------------

def extract_design_system(pptx_path: Path) -> ExtractedDesignSystem:
    """Parse a .pptx and return its extracted design system. Always returns —
    missing pieces fall back to defaults and are recorded in `notes`."""
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")

    notes: list[str] = []
    prs = Presentation(str(pptx_path))

    # Theme XML — parsed directly from the zip (python-pptx doesn't surface it)
    theme_bytes = _read_theme_xml(pptx_path)
    if theme_bytes is None:
        notes.append("no ppt/theme/*.xml found in archive; using default colors and fonts")
        theme_colors = TemplateColors()
        fonts = TemplateFonts()
    else:
        theme_colors, fonts = _parse_theme(theme_bytes, notes)

    master_layouts = _master_layout_summary(prs)
    logo_candidates = _walk_master_pictures(prs)
    footer_text = _extract_master_footer(prs)
    sample_palette = _sample_palette_from_slides(pptx_path)

    slide_width_inches = _emu_to_inches(prs.slide_width)
    slide_height_inches = _emu_to_inches(prs.slide_height)

    notes.append(f"found {len(master_layouts)} layouts across {len(prs.slide_masters)} master(s)")
    if logo_candidates:
        notes.append(f"found {len(logo_candidates)} picture shape(s) on masters/layouts (likely logos)")
    if sample_palette:
        notes.append(f"sampled {len(sample_palette)} distinct hex colors from content slides")

    return ExtractedDesignSystem(
        theme_colors=theme_colors,
        fonts=fonts,
        master_layouts=master_layouts,
        logo_candidates=logo_candidates,
        footer_text=footer_text,
        sample_palette=sample_palette,
        slide_width_inches=slide_width_inches,
        slide_height_inches=slide_height_inches,
        notes=notes,
    )


def ingest_to_template(
    pptx_path: Path,
    template_name: str,
    templates_dir: Path,
    description: str = "",
) -> Template:
    """Extract a .pptx's design system, save it as templates/<slug>.json, and
    copy the source pptx in as templates/<slug>-master.pptx. The returned
    Template has master_pptx set so rendering goes through MasterRenderer and
    inherits all the master's visual decisions."""
    pptx_path = Path(pptx_path)
    templates_dir = Path(templates_dir)
    templates_dir.mkdir(parents=True, exist_ok=True)

    design = extract_design_system(pptx_path)

    slug = _slugify(template_name)
    master_filename = f"{slug}-master.pptx"
    master_dest = templates_dir / master_filename
    shutil.copy2(pptx_path, master_dest)

    layout_map = detect_layout_map(master_dest)

    template = Template(
        name=template_name,
        description=description,
        master_pptx=master_filename,
        layout_map=layout_map,
        colors=design.theme_colors,
        fonts=design.fonts,
        slide_width_inches=design.slide_width_inches or 13.333,
        slide_height_inches=design.slide_height_inches or 7.5,
    )

    save_template(template, templates_dir / f"{slug}.json")
    return template
