"""Copy a single slide from one .pptx file into another, preserving all native
shape primitives (rectangles stay rectangles, ovals stay ovals, freeforms
stay freeforms) so the output remains fully editable in PowerPoint.

This is the foundational primitive for the asset library: when a generated
deck references `asset_ref = "<id>"`, the renderer uses this function to
stamp the library slide into the output deck.

Approach:
  - Open both source and destination presentations
  - Find a "blank" layout in the destination (we add to a layout that has no
    placeholders so the copied shapes don't conflict)
  - Deep-copy the source slide's spTree children (shape XML elements) into
    the new destination slide
  - Copy any embedded media (images) so references don't break

Known limitations:
  - Theme-referenced colors/fonts may shift if source and dest themes differ
  - Animations are not preserved
  - Slide dimensions of the dest deck may differ from the source — for
    layout-sensitive assets you may want to set dest dimensions to match
    the asset's native size before copying
"""
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


def copy_slide_into(
    source_pptx: Path,
    source_slide_index: int,  # 1-based to match library catalog convention
    dest_presentation,
    *,
    match_dimensions: bool = False,
    scale_to_dest: bool = True,
):
    """Copy a slide from source_pptx (at the given 1-based index) into the
    dest_presentation. Returns the newly-added slide.

    If match_dimensions=True, the destination presentation's slide_width and
    slide_height are updated to match the source slide's dimensions before
    copying. Use this when the asset is dimensionally sensitive (most
    diagrams from a 26.7"×15" library are).

    If scale_to_dest=True (the default) and match_dimensions=False, the
    top-level shapes of the copied slide are proportionally scaled so they
    fit the destination canvas. This fixes the common case where a library
    slide authored at 26.67"×15" is stamped into a 13.33"×7.5" deck and
    shapes spill off the right/bottom edges. NOTE: text font sizes are
    intentionally left untouched — PowerPoint's autofit handles text
    overflow more gracefully than a blind multiplicative scale would.
    """
    source_pres = Presentation(str(source_pptx))
    if source_slide_index < 1 or source_slide_index > len(source_pres.slides):
        raise IndexError(
            f"Slide index {source_slide_index} out of range for {source_pptx} "
            f"(has {len(source_pres.slides)} slides, 1-based)"
        )

    # Decide on a scaling factor BEFORE we mutate dest dimensions.
    if match_dimensions:
        dest_presentation.slide_width = source_pres.slide_width
        dest_presentation.slide_height = source_pres.slide_height
        scale_x = scale_y = 1.0
    elif scale_to_dest:
        scale_x = dest_presentation.slide_width / source_pres.slide_width
        scale_y = dest_presentation.slide_height / source_pres.slide_height
    else:
        scale_x = scale_y = 1.0

    source_slide = source_pres.slides[source_slide_index - 1]

    # Use the blank layout (idx 6 by convention) so we don't conflict with
    # any layout placeholders. If the destination master uses different
    # indices, fall back to the last layout (often a "Blank" too).
    layouts = dest_presentation.slide_layouts
    blank_layout = None
    for layout in layouts:
        if "blank" in layout.name.lower():
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = layouts[6] if len(layouts) > 6 else layouts[-1]

    new_slide = dest_presentation.slides.add_slide(blank_layout)

    # Copy each shape from source to dest by deep-copying its XML element.
    # This preserves all native shape types (auto_shape, freeform, picture,
    # group, etc.) and their formatting.
    src_spTree = source_slide.shapes._spTree
    dst_spTree = new_slide.shapes._spTree

    # Iterate XML children directly — covers shapes, group shapes, pictures,
    # connectors. Skip the non-shape elements that begin/end the spTree.
    from pptx.oxml.ns import qn
    skip_tags = {qn("p:nvGrpSpPr"), qn("p:grpSpPr")}

    for child in list(src_spTree):
        if child.tag in skip_tags:
            continue
        new_element = deepcopy(child)
        dst_spTree.append(new_element)

    # Copy speaker notes if present
    if source_slide.has_notes_slide:
        notes_text = source_slide.notes_slide.notes_text_frame.text
        if notes_text:
            new_slide.notes_slide.notes_text_frame.text = notes_text

    # Copy embedded images (relationships).
    # Each picture in the copied XML refers to a relationship ID (rId) that
    # only existed in the source presentation. We need to recreate those
    # relationships in the destination so the image data is preserved.
    _rebind_image_relationships(source_slide, new_slide)

    # Proportionally scale top-level shapes if the dest canvas differs.
    # We only scale geometry (left/top/width/height); font sizes are left
    # alone so PowerPoint's text autofit can handle text overflow.
    # For nested groups, python-pptx propagates the new width/height down
    # to children automatically, so a single pass over top-level shapes
    # suffices.
    if scale_x != 1.0 or scale_y != 1.0:
        for shape in new_slide.shapes:
            try:
                shape.left = int(shape.left * scale_x)
                shape.top = int(shape.top * scale_y)
                shape.width = int(shape.width * scale_x)
                shape.height = int(shape.height * scale_y)
            except (TypeError, AttributeError):
                # Some shape kinds (e.g. connectors with absent extents)
                # may not expose all four properties — skip them safely.
                continue

    return new_slide


def copy_slide_shapes_onto(
    source_pptx: Path,
    source_slide_index: int,  # 1-based
    dest_slide,
    *,
    target_left_emu: int,
    target_top_emu: int,
    target_width_emu: int,
    target_height_emu: int,
    strip_text: "bool | str" = "chrome_only",
    recolor_to=None,  # RGBColor | None — when set, recolor dominant solid fill
):
    """Copy SHAPES from a library slide onto an EXISTING destination slide,
    scaling/offsetting them to fit a target rectangle. Used for "extras":
    attaching a library asset to the right half of a slide that already has
    a title and body of its own.

    Unlike `copy_slide_into`, this does NOT add a new slide. It mutates the
    given dest_slide by appending the source's shapes (scaled and translated).
    Image relationships are rebound so embedded pictures travel with the copy.

    `strip_text` modes (Sprint Z):
        - "none" / False: keep all text as-is.
        - "all" / True:   blank every text frame (legacy blunt strip — kills
                          inline annotations along with chrome).
        - "chrome_only" (DEFAULT): blank only outer slide chrome — title bar,
                          subtitle, footer, page-number — but KEEP text that
                          looks inline-to-the-diagram (quadrant labels,
                          axis labels, callouts, step names). Heuristic in
                          `_is_chrome_text`. This is what users actually want
                          when attaching a library diagram as an extra.
    """
    from copy import deepcopy

    source_pres = Presentation(str(source_pptx))
    if source_slide_index < 1 or source_slide_index > len(source_pres.slides):
        raise IndexError(
            f"Slide index {source_slide_index} out of range for {source_pptx} "
            f"(has {len(source_pres.slides)} slides, 1-based)"
        )
    source_slide = source_pres.slides[source_slide_index - 1]

    src_w = source_pres.slide_width
    src_h = source_pres.slide_height
    if src_w <= 0 or src_h <= 0:
        return  # nothing sensible to scale to

    scale_x = target_width_emu / src_w
    scale_y = target_height_emu / src_h

    # Track which top-level XML elements we add so we only scale those.
    src_spTree = source_slide.shapes._spTree
    dst_spTree = dest_slide.shapes._spTree

    from pptx.oxml.ns import qn
    skip_tags = {qn("p:nvGrpSpPr"), qn("p:grpSpPr")}

    new_top_level_elements = []
    for child in list(src_spTree):
        if child.tag in skip_tags:
            continue
        new_element = deepcopy(child)
        dst_spTree.append(new_element)
        new_top_level_elements.append(new_element)

    # Rebind any image relationships so embedded pictures don't break
    _rebind_image_relationships(source_slide, dest_slide)

    # Resolve strip mode. Backward-compat: True → "all", False → "none".
    if isinstance(strip_text, bool):
        strip_mode = "all" if strip_text else "none"
    else:
        strip_mode = (strip_text or "none").lower()
        if strip_mode not in {"all", "chrome_only", "none"}:
            strip_mode = "chrome_only"  # safe default for unknown values

    # Now scale + translate each newly-added top-level shape. We need to
    # iterate dest_slide.shapes to get python-pptx's shape objects (which
    # expose left/top/width/height APIs) — but only the LAST N shapes are
    # the ones we just added. Walk dest_slide.shapes and pick the ones
    # whose underlying element is in our `new_top_level_elements` set.
    #
    # NOTE: we collect the new shapes BEFORE scaling, then apply text-strip
    # (which may need pre-scale positions to detect chrome by source-slide
    # percentage), THEN scale. Order matters: shape.left/top/width/height
    # below are still in source-slide EMU when the strip step runs.
    new_set = set(id(e) for e in new_top_level_elements)
    new_shapes: list = []
    for shape in dest_slide.shapes:
        try:
            if id(shape.element) not in new_set:
                continue
        except AttributeError:
            continue
        new_shapes.append(shape)

    # Strip placeholder text from text frames in the copied shapes.
    # - "all":          blunt — every text frame is blanked. Library slides
    #                   ship with lorem-ipsum content the user almost never
    #                   wants. Legacy behaviour.
    # - "chrome_only":  smart — only outer slide chrome (title, subtitle,
    #                   footer, page-number, peripheral title-bar text) is
    #                   blanked; inline diagram annotations (quadrant labels,
    #                   axis labels, callouts) survive. Sprint Z default.
    # - "none":         no stripping.
    if strip_mode == "all":
        from pptx.oxml.ns import qn
        for new_el in new_top_level_elements:
            for tf_xml in new_el.findall(".//" + qn("a:t")):
                tf_xml.text = ""
    elif strip_mode == "chrome_only":
        for shape in new_shapes:
            try:
                if _is_chrome_text(shape, src_w, src_h):
                    _blank_shape_text(shape)
            except Exception:  # noqa: BLE001 — chrome detection is best-effort
                continue

    # Apply scale + translate. AFTER strip so position-percent heuristics
    # operate on source-slide coordinates.
    for shape in new_shapes:
        try:
            shape.left = int(shape.left * scale_x) + target_left_emu
            shape.top = int(shape.top * scale_y) + target_top_emu
            shape.width = int(shape.width * scale_x)
            shape.height = int(shape.height * scale_y)
        except (TypeError, AttributeError):
            continue

    # Optional dominant-color recolor. We walk only the newly-added top-level
    # shapes, count solid-fill RGB colors on AUTO_SHAPE / FREEFORM kinds, and
    # replace the most common color with `recolor_to`. This preserves accent
    # colors while shifting the asset's primary color to the user's choice.
    if recolor_to is not None:
        _recolor_dominant_fill(new_shapes, recolor_to)


def _recolor_dominant_fill(shapes: list, recolor_to) -> None:
    """Find the dominant solid-fill color across `shapes` (and their group
    descendants) and replace it with `recolor_to`. Skips pictures/text-only
    shapes — only AUTO_SHAPE and FREEFORM are considered.

    "Dominant" = most common (r,g,b) tuple across all eligible shapes. When
    multiple colors tie, the first one seen wins (deterministic enough)."""
    from collections import Counter
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    eligible_kinds = {MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM}

    def _walk(shape_list):
        for shp in shape_list:
            try:
                if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                    _walk(shp.shapes)
                    continue
            except (AttributeError, ValueError):
                pass
            yield shp

    counter: Counter = Counter()
    candidates: list = []  # (shape, current_rgb_tuple)
    for shp in _walk(shapes):
        try:
            if shp.shape_type not in eligible_kinds:
                continue
        except (AttributeError, ValueError):
            continue
        try:
            fill = shp.fill
            fc = fill.fore_color
            rgb = fc.rgb  # raises if not solid / not RGB
        except (AttributeError, TypeError, ValueError, Exception):  # noqa: BLE001
            continue
        if rgb is None:
            continue
        try:
            rgb_tuple = (rgb[0], rgb[1], rgb[2])
        except (TypeError, IndexError):
            continue
        counter[rgb_tuple] += 1
        candidates.append((shp, rgb_tuple))

    if not counter:
        return
    dominant, _ = counter.most_common(1)[0]
    for shp, rgb_tuple in candidates:
        if rgb_tuple != dominant:
            continue
        try:
            shp.fill.solid()
            shp.fill.fore_color.rgb = recolor_to
        except Exception:  # noqa: BLE001
            continue


def _rebind_image_relationships(source_slide, dest_slide):
    """When we deep-copy XML, picture elements still reference rIds from
    the source slide's relationships. We need to add the corresponding
    image parts to the destination's relationships and rewrite the rIds.
    """
    from pptx.oxml.ns import qn

    # Find all picture references (a:blip elements) in the destination slide
    # and rebind them.
    blip_elements = dest_slide.shapes._spTree.findall(
        ".//" + qn("a:blip")
    )
    if not blip_elements:
        return

    # Build a map of source rIds → image parts
    source_rels = source_slide.part.rels
    embed_attr = qn("r:embed")
    link_attr = qn("r:link")

    for blip in blip_elements:
        src_rid = blip.get(embed_attr) or blip.get(link_attr)
        if not src_rid:
            continue
        if src_rid not in source_rels:
            continue
        rel = source_rels[src_rid]
        # Add the image part to dest and get a new rId.
        # Use relate_to which de-duplicates if the same image is already related.
        try:
            new_rid = dest_slide.part.relate_to(rel.target_part, rel.reltype)
        except Exception:  # noqa: BLE001
            continue
        # Rewrite the blip's reference
        if blip.get(embed_attr):
            blip.set(embed_attr, new_rid)
        elif blip.get(link_attr):
            blip.set(link_attr, new_rid)


# ---------------------------------------------------------------------------
# Sprint Z — selective text stripping helpers
# ---------------------------------------------------------------------------

# Placeholder types that always count as "outer chrome": title bar, subtitle,
# footer, page numbers, header, date — all the slide-level scaffolding the
# user does NOT want carried over when they attach a diagram as an extra.
def _chrome_placeholder_types():
    """Return the set of PP_PLACEHOLDER values to treat as chrome.
    Returns an empty set if pptx isn't importable for any reason — caller
    falls through to position/name heuristics."""
    try:
        from pptx.enum.text import PP_PLACEHOLDER
    except ImportError:
        return set()
    out = set()
    for name in (
        "TITLE", "CENTER_TITLE", "SUBTITLE",
        "FOOTER", "SLIDE_NUMBER", "DATE", "HEADER",
    ):
        v = getattr(PP_PLACEHOLDER, name, None)
        if v is not None:
            out.add(v)
    return out


def _is_chrome_text(shape, slide_w_emu: int, slide_h_emu: int) -> bool:
    """Decide whether `shape`'s text is outer slide chrome (strip) versus
    inline diagram annotation (keep). Sprint Z heuristic. Layered:

      1. If shape is a title/subtitle/footer/page-number placeholder →
         chrome (most reliable signal).
      2. If shape is positioned in the top 12% of the slide AND looks
         titlish (large font OR short text) → chrome.
      3. If shape's vertical centre is below 88% of slide height → chrome
         (footer / page-number zone — fonts are tiny but always peripheral).
      4. If shape name pattern matches title/subtitle/footer/page → chrome.
      5. Otherwise → keep (assume inline annotation).

    `slide_w_emu` / `slide_h_emu` are the SOURCE slide's dimensions, since
    the caller hasn't yet scaled positions to the destination rect.
    Conservative: when the heuristic can't decide, returns False (keep).
    """
    if not getattr(shape, "has_text_frame", False):
        return False
    try:
        if not shape.text_frame.text.strip():
            return False  # already empty — nothing to strip
    except (AttributeError, TypeError):
        return False

    # 1. Placeholder-based detection.
    try:
        ph = shape.placeholder_format
        if ph is not None and ph.idx is not None:
            ph_type = ph.type
            if ph_type in _chrome_placeholder_types():
                return True
    except (AttributeError, ValueError):
        pass

    # 2 + 3. Position-based.
    try:
        top = shape.top
        height = shape.height
        if top is not None and height is not None and slide_h_emu > 0:
            top_pct = top / slide_h_emu
            v_centre_pct = (top + height / 2) / slide_h_emu
            # Top 12%: title-bar zone — only counts as chrome if titlish.
            if top_pct < 0.12 and _looks_titlish(shape):
                return True
            # Bottom 8% (centre below 92%): footer / page-number zone.
            if v_centre_pct > 0.92:
                return True
    except (TypeError, AttributeError):
        pass

    # 4. Shape-name pattern (last resort — names like "Title 1", "Footer 1").
    try:
        name = (shape.name or "").lower()
        for pat in ("title", "subtitle", "footer", "page number",
                    "slide number", "header placeholder", "date placeholder"):
            if pat in name:
                return True
    except AttributeError:
        pass

    return False


def _looks_titlish(shape) -> bool:
    """For text in the top 12% of a slide: distinguish a real title bar
    from a small annotation that just happens to sit near the top.

    Titlish if any run has font-size ≥ 24pt, OR text is short (≤8 words)
    with no explicit sizing. Long text (>12 words) is never titlish even
    in the top zone — probably a description band."""
    try:
        words = shape.text_frame.text.split()
        if len(words) > 12:
            return False
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                size = run.font.size
                if size is not None and size.pt >= 24:
                    return True
        return len(words) <= 8
    except (AttributeError, TypeError, ValueError):
        return True  # be permissive — better to over-strip a peripheral
                     # title-zone shape than under-strip


def _blank_shape_text(shape) -> None:
    """Blank every text run in `shape` while leaving the shape itself
    intact (so the rectangle/box still renders, just empty). Same XML
    pathway as the legacy `strip_text=True` blunt path, but scoped to
    one shape rather than the whole copied tree."""
    if not getattr(shape, "has_text_frame", False):
        return
    try:
        from pptx.oxml.ns import qn
        txBody = shape.text_frame._txBody
        for tf_xml in txBody.findall(".//" + qn("a:t")):
            tf_xml.text = ""
    except (AttributeError, ImportError):
        pass
