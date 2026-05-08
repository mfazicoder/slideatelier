"""Asset library — indexes a directory of .pptx files into a queryable catalog.

Each individual slide inside a .pptx is treated as a discrete "asset" (e.g.,
one specific 5-stage funnel design). The catalog records stable IDs so we can
reference assets from generated deck specs (`Slide.asset_ref = "<asset_id>"`).
"""
import json
import re
import unicodedata
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pydantic import BaseModel, Field


def slugify(s: str) -> str:
    """File-system-safe slug: ascii lowercase, hyphens for separators."""
    s = unicodedata.normalize("NFKD", s).encode("ascii", "ignore").decode("ascii")
    s = re.sub(r"[^\w\s-]", "", s).strip().lower()
    s = re.sub(r"[-\s]+", "-", s)
    return s.strip("-")


class LibraryAsset(BaseModel):
    """One slide inside a library .pptx, indexed and addressable."""
    id: str = Field(description="Stable ID: <category-slug>/<file-stem-slug>/<1-based-slide-index>")
    file_path: str = Field(description="Absolute path to the source .pptx")
    slide_index: int = Field(description="1-based slide index within the file")
    category: str = Field(description="Human-readable category name (folder name)")
    category_slug: str
    file_stem: str = Field(description="Base filename without .pptx extension")
    shape_count: int
    text_snippets: list[str] = Field(
        default_factory=list,
        description="Up to 5 distinct non-empty text strings found on the slide — useful for search",
    )
    slide_width_inches: float
    slide_height_inches: float
    has_thumbnail: bool = False
    notes: str = Field(default="", description="Human-curated description (filled in during phase 3)")
    style_tags: list[str] = Field(
        default_factory=list,
        description="Tags describing the visual/structural style (e.g. 'consulting', 'minimal', 'process'). Used to filter and recommend assets.",
    )


class LibraryCatalog(BaseModel):
    library_root: str
    scanned_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    asset_count: int = 0
    categories: list[str] = Field(default_factory=list)
    assets: list[LibraryAsset] = Field(default_factory=list)

    def by_category(self) -> dict[str, list[LibraryAsset]]:
        out: dict[str, list[LibraryAsset]] = {}
        for a in self.assets:
            out.setdefault(a.category, []).append(a)
        return out

    def find(self, asset_id: str) -> LibraryAsset | None:
        for a in self.assets:
            if a.id == asset_id:
                return a
        return None


# Heuristic table mapping a folder-name substring to a list of style tags.
# Order matters: longer / more specific keys are tried before shorter ones so
# that "Digital Marketing" beats plain "Marketing".
_CATEGORY_TAG_RULES: list[tuple[str, list[str]]] = [
    ("digital marketing", ["marketing", "contemporary", "digital"]),
    ("business diagrams", ["structured", "framework", "consulting"]),
    ("comparison", ["comparison", "structured", "consulting"]),
    ("real estate", ["real-estate", "contemporary"]),
    ("watercolor", ["artistic", "creative"]),
    ("alphabet", ["typographic", "decorative"]),
    ("roadmap", ["process", "sequential", "planning"]),
    ("timeline", ["process", "sequential"]),
    ("scrum", ["process", "agile", "technical"]),
    ("pyramid", ["hierarchy", "structured", "consulting"]),
    ("funnel", ["process", "narrowing", "marketing", "consulting"]),
    ("cycle", ["process", "circular", "structured"]),
    ("venn", ["comparison", "overlap", "structured"]),
    ("marketing", ["marketing", "contemporary"]),
]

# Tags inferred from substrings in a file's stem (filename without .pptx).
_FILE_STEM_TAG_RULES: list[tuple[str, list[str]]] = [
    ("outline", ["minimal"]),
    ("infographic", ["graphics-heavy"]),
    ("flat", ["flat", "minimal"]),
    ("3d", ["graphics-heavy"]),
    ("hand-drawn", ["artistic", "creative"]),
    ("hand drawn", ["artistic", "creative"]),
    ("sketch", ["artistic", "creative"]),
    ("watercolor", ["artistic", "creative"]),
    ("isometric", ["graphics-heavy", "modern"]),
    ("gradient", ["modern", "vibrant"]),
    ("minimal", ["minimal"]),
]


def auto_tag_asset(asset: "LibraryAsset") -> list[str]:
    """Heuristically infer style_tags for a library asset based on its
    category folder name and file stem. Returns a de-duplicated list.

    Falls back to ``["medium-corporate"]`` so that nothing is left untagged
    (which would render the asset effectively invisible to style filters).
    """
    tags: list[str] = []
    seen: set[str] = set()

    def _add(t: str) -> None:
        if t and t not in seen:
            seen.add(t)
            tags.append(t)

    cat_lc = (asset.category or "").lower()
    for keyword, rule_tags in _CATEGORY_TAG_RULES:
        if keyword in cat_lc:
            for t in rule_tags:
                _add(t)
            break  # one category match is enough

    stem_lc = (asset.file_stem or "").lower()
    for keyword, rule_tags in _FILE_STEM_TAG_RULES:
        if keyword in stem_lc:
            for t in rule_tags:
                _add(t)

    if not tags:
        _add("medium-corporate")
    return tags


def _extract_text_snippets(slide, max_snippets: int = 5) -> list[str]:
    """Collect distinct, non-empty, sufficiently-meaningful text strings from a slide."""
    seen: set[str] = set()
    snippets: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        # Filter trivial / repeated artifacts
        if len(text) < 3:
            continue
        # Normalize whitespace
        text = re.sub(r"\s+", " ", text)
        if text in seen:
            continue
        seen.add(text)
        snippets.append(text[:120])  # cap length
        if len(snippets) >= max_snippets:
            break
    return snippets


def scan_library(
    library_root: Path,
    *,
    categories: list[str] | None = None,
    progress_callback=None,
) -> LibraryCatalog:
    """Walk the library_root, optionally filtered to specific categories
    (matched against folder names), and build a LibraryCatalog.

    `categories` is a list of substrings. A folder matches if any substring
    appears (case-insensitive) in the folder name. None = scan everything.
    """
    if not library_root.exists():
        raise FileNotFoundError(f"Library root not found: {library_root}")

    catalog = LibraryCatalog(library_root=str(library_root.resolve()))

    # Pre-filter by category folder
    category_dirs: list[Path] = []
    for child in sorted(library_root.iterdir()):
        if not child.is_dir():
            continue
        if categories is not None:
            if not any(c.lower() in child.name.lower() for c in categories):
                continue
        category_dirs.append(child)

    catalog.categories = [d.name for d in category_dirs]

    for cat_dir in category_dirs:
        cat_slug = slugify(cat_dir.name)
        for pptx_path in sorted(cat_dir.glob("*.pptx")):
            # Skip Office lock files (~$prefix)
            if pptx_path.name.startswith("~$"):
                continue
            try:
                prs = Presentation(str(pptx_path))
            except Exception as e:  # noqa: BLE001
                if progress_callback:
                    progress_callback(f"skip (load error): {pptx_path.name} — {e}")
                continue

            file_stem = pptx_path.stem
            file_stem_slug = slugify(file_stem)
            width_in = prs.slide_width / 914400
            height_in = prs.slide_height / 914400

            for slide_idx, slide in enumerate(prs.slides, start=1):
                shape_count = len(list(slide.shapes))
                snippets = _extract_text_snippets(slide)
                asset = LibraryAsset(
                    id=f"{cat_slug}/{file_stem_slug}/{slide_idx}",
                    file_path=str(pptx_path.resolve()),
                    slide_index=slide_idx,
                    category=cat_dir.name,
                    category_slug=cat_slug,
                    file_stem=file_stem,
                    shape_count=shape_count,
                    text_snippets=snippets,
                    slide_width_inches=width_in,
                    slide_height_inches=height_in,
                )
                asset.style_tags = auto_tag_asset(asset)
                catalog.assets.append(asset)

            if progress_callback:
                progress_callback(f"{cat_dir.name} / {pptx_path.name} — {len(prs.slides)} slides")

    catalog.asset_count = len(catalog.assets)
    return catalog


def save_catalog(catalog: LibraryCatalog, path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(catalog.model_dump_json(indent=2))


def load_catalog(path: Path) -> LibraryCatalog:
    if not path.exists():
        raise FileNotFoundError(f"Catalog not found: {path}. Run `atelier library scan` first.")
    return LibraryCatalog.model_validate_json(path.read_text())
