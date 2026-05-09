"""Translate a library .pptx slide into inline SVG for the Web Deck viewer.

Sprint Z.v2 — replaces the baked-thumbnail PNG fallback with native SVG so
text reflows, theme tokens hot-swap, and `chrome_only` stripping applies to
the published web view (matching what the .pptx export does).

Approach:
  - Walk the source slide's shapes recursively (descend into groups).
  - For each shape, dispatch on shape_type / auto_shape_type and emit the
    appropriate SVG primitive (`<rect>`, `<ellipse>`, `<circle>`, `<polygon>`,
    `<path>`, `<image>`, `<text>`).
  - Coordinate space: outer `<svg>` uses `viewBox="0 0 {src_w_emu} {src_h_emu}"`
    so each shape can use its native EMU coords directly; the outer width/
    height attributes scale to the target rendering size automatically.
  - Solid fills + line colours resolved to `#RRGGBB`. Theme/scheme colours
    that python-pptx can't surface as RGB fall back to a sensible default.
  - Chrome filtering reuses `asset_copier._is_chrome_text` so the heuristic
    matches the .pptx attach path bit-for-bit.

Anything we can't translate cleanly (charts, tables, complex custGeom) gets
a placeholder `<rect>` of the shape's bounding box so the visual layout
stays roughly right. Worst case is no worse than the pre-Sprint-Z.v2
thumbnail.
"""
from __future__ import annotations

import base64
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from .asset_copier import _is_chrome_text


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def library_asset_to_svg(
    pptx_path: Path,
    slide_index: int,  # 1-based
    target_w_px: int,
    target_h_px: int,
    *,
    strip_text: str = "chrome_only",
) -> str:
    """Render a single slide of `pptx_path` as inline SVG.

    `strip_text`: same modes as `copy_slide_shapes_onto`:
      - "none":          keep all text.
      - "chrome_only":   strip outer chrome (title/footer/etc.), keep inline.
      - "all":           strip every text frame.
    """
    prs = Presentation(str(pptx_path))
    if slide_index < 1 or slide_index > len(prs.slides):
        raise IndexError(
            f"slide index {slide_index} out of range for {pptx_path} "
            f"(1-based, has {len(prs.slides)} slides)"
        )
    slide = prs.slides[slide_index - 1]
    src_w = int(prs.slide_width)
    src_h = int(prs.slide_height)
    if src_w <= 0 or src_h <= 0:
        return _empty_svg(target_w_px, target_h_px)

    # rId → image bytes for the picture-shape pass. Built once per render so
    # we don't re-parse the package per shape.
    image_lookup = _build_image_lookup(slide)

    parts: list[str] = [
        f'<svg xmlns="http://www.w3.org/2000/svg" '
        f'width="{target_w_px}" height="{target_h_px}" '
        f'viewBox="0 0 {src_w} {src_h}" '
        f'preserveAspectRatio="xMidYMid meet">'
    ]

    for shape in slide.shapes:
        try:
            _emit_shape(shape, parts, src_w, src_h, strip_text, image_lookup)
        except Exception as e:  # noqa: BLE001 — never let one shape kill the page
            parts.append(
                f'  <!-- skipped {getattr(shape, "name", "?")} '
                f'({type(e).__name__}: {_xml_escape(str(e))[:80]}) -->'
            )

    parts.append("</svg>")
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Shape dispatch
# ---------------------------------------------------------------------------

def _emit_shape(shape, parts, slide_w, slide_h, strip_text, image_lookup):
    """Top-level dispatch on shape kind."""
    st = shape.shape_type

    if st == MSO_SHAPE_TYPE.GROUP:
        # Recurse into children. python-pptx wraps group children with
        # their own scaled coords, so we just descend.
        for child in shape.shapes:
            try:
                _emit_shape(child, parts, slide_w, slide_h, strip_text, image_lookup)
            except Exception:  # noqa: BLE001
                continue
        return

    # Decide whether this shape's text counts as chrome. If so:
    #  - For pure TEXT_BOX shapes, skip the whole shape (no geometry).
    #  - For AUTO_SHAPE / others, emit the geometry but suppress the text.
    is_chrome = False
    if shape.has_text_frame and strip_text != "none":
        if strip_text == "all":
            is_chrome = True
        elif strip_text == "chrome_only":
            try:
                is_chrome = _is_chrome_text(shape, slide_w, slide_h)
            except Exception:  # noqa: BLE001
                is_chrome = False

    if st == MSO_SHAPE_TYPE.TEXT_BOX:
        if is_chrome:
            return  # text-only chrome — drop entirely
        _emit_textbox(shape, parts)
        return

    if st == MSO_SHAPE_TYPE.PICTURE:
        _emit_picture(shape, parts, image_lookup)
        return

    if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
        _emit_auto_shape(shape, parts)
        if shape.has_text_frame and not is_chrome:
            _emit_text_inside(shape, parts)
        return

    if st == MSO_SHAPE_TYPE.FREEFORM:
        # Custom geometry — best-effort: emit a bounding-box rect so the
        # silhouette is visible. Custom-path translation is queued.
        _emit_placeholder_rect(shape, parts, hint="freeform")
        if shape.has_text_frame and not is_chrome:
            _emit_text_inside(shape, parts)
        return

    if st == MSO_SHAPE_TYPE.LINE:
        _emit_line(shape, parts)
        return

    # CHART, TABLE, CONNECTOR (other forms), MEDIA, etc.: emit a low-key
    # bounding-box placeholder so the layout stays correct.
    _emit_placeholder_rect(shape, parts, hint=str(st))


# ---------------------------------------------------------------------------
# Per-shape emitters
# ---------------------------------------------------------------------------

def _emit_auto_shape(shape, parts):
    """AutoShape (rectangle, oval, triangle, diamond, chevron, ...).
    Reads `auto_shape_type` and emits the matching native SVG primitive.
    Unknown types fall back to a `<rect>` of the bounding box."""
    left, top, w, h = _bounds(shape)
    if left is None:
        return
    fill = _resolve_fill(shape)
    stroke = _resolve_line(shape)

    try:
        ast = shape.auto_shape_type
    except Exception:  # noqa: BLE001
        ast = None

    style = f'fill="{fill}" stroke="{stroke[0]}" stroke-width="{stroke[1]}"'

    if ast == MSO_SHAPE.RECTANGLE:
        parts.append(f'  <rect x="{left}" y="{top}" width="{w}" height="{h}" {style}/>')
    elif ast == MSO_SHAPE.ROUNDED_RECTANGLE:
        rxy = min(w, h) // 8
        parts.append(
            f'  <rect x="{left}" y="{top}" width="{w}" height="{h}" '
            f'rx="{rxy}" ry="{rxy}" {style}/>'
        )
    elif ast == MSO_SHAPE.OVAL:
        cx = left + w // 2
        cy = top + h // 2
        if w == h:
            parts.append(f'  <circle cx="{cx}" cy="{cy}" r="{w // 2}" {style}/>')
        else:
            parts.append(
                f'  <ellipse cx="{cx}" cy="{cy}" rx="{w // 2}" ry="{h // 2}" {style}/>'
            )
    elif ast == MSO_SHAPE.RIGHT_TRIANGLE:
        pts = f"{left},{top + h} {left + w},{top + h} {left + w},{top}"
        parts.append(f'  <polygon points="{pts}" {style}/>')
    elif ast == MSO_SHAPE.ISOCELES_TRIANGLE:
        pts = f"{left + w // 2},{top} {left},{top + h} {left + w},{top + h}"
        parts.append(f'  <polygon points="{pts}" {style}/>')
    elif ast == MSO_SHAPE.DIAMOND:
        pts = (f"{left + w // 2},{top} {left + w},{top + h // 2} "
               f"{left + w // 2},{top + h} {left},{top + h // 2}")
        parts.append(f'  <polygon points="{pts}" {style}/>')
    elif ast == MSO_SHAPE.CHEVRON:
        notch = h // 2
        pts = (f"{left},{top} {left + w - notch},{top} {left + w},{top + h // 2} "
               f"{left + w - notch},{top + h} {left},{top + h} "
               f"{left + notch},{top + h // 2}")
        parts.append(f'  <polygon points="{pts}" {style}/>')
    elif ast == MSO_SHAPE.PENTAGON:
        pts = (f"{left + w // 2},{top} {left + w},{top + int(h * 0.4)} "
               f"{left + int(w * 0.8)},{top + h} {left + int(w * 0.2)},{top + h} "
               f"{left},{top + int(h * 0.4)}")
        parts.append(f'  <polygon points="{pts}" {style}/>')
    elif ast == MSO_SHAPE.HEXAGON:
        pts = (f"{left + w // 4},{top} {left + 3 * w // 4},{top} "
               f"{left + w},{top + h // 2} {left + 3 * w // 4},{top + h} "
               f"{left + w // 4},{top + h} {left},{top + h // 2}")
        parts.append(f'  <polygon points="{pts}" {style}/>')
    else:
        # Any other auto-shape we don't have a hand-mapping for: emit a
        # bounding-box rect so it's visible.
        parts.append(f'  <rect x="{left}" y="{top}" width="{w}" height="{h}" {style}/>')


def _emit_text_inside(shape, parts):
    """Emit `<text>` runs centred inside an AutoShape's bounding box. Each
    paragraph becomes a `<tspan>` line; multi-line text uses dy stacking."""
    left, top, w, h = _bounds(shape)
    if left is None:
        return
    tf = shape.text_frame
    if not tf.text.strip():
        return

    # Collect line metadata. We can't know exact line height without
    # rendering, so we approximate from font size or 0.18 of shape height.
    paragraphs = [p for p in tf.paragraphs if p.text.strip()]
    if not paragraphs:
        return

    # Pick a representative font size from the first run with a size set,
    # else default to ~min(h/4, slide-body-default).
    font_size = None
    font_family = None
    for p in paragraphs:
        for r in p.runs:
            if r.font.size is not None:
                font_size = int(r.font.size)
                break
            if r.font.name:
                font_family = r.font.name
        if font_size:
            break
    if font_size is None:
        # Conservative default: half-of-shape-height capped at 360000 EMU (~24pt)
        font_size = min(h // 3, 360000)
    if font_family is None:
        font_family = "Inter, sans-serif"

    # Colour from the first run's font.color, fall back to dark grey.
    fill = "#1F2937"
    try:
        first_run = next((r for p in paragraphs for r in p.runs), None)
        if first_run is not None:
            c = first_run.font.color
            if c.type is not None and c.rgb is not None:
                fill = f"#{str(c.rgb)}"
    except Exception:  # noqa: BLE001
        pass

    cx = left + w // 2
    cy = top + h // 2
    line_h = int(font_size * 1.25)
    n = len(paragraphs)
    # Vertical center: total block height ≈ n * line_h; first line offset is
    # cy - (n-1)*line_h/2 + (font_size baseline correction).
    start_y = cy - (n - 1) * line_h // 2 + font_size // 3

    # Determine alignment from first paragraph.
    align_map = {1: "start", 2: "middle", 3: "end"}  # PP_ALIGN.LEFT/CENTER/RIGHT
    try:
        first_align = paragraphs[0].alignment
        text_anchor = align_map.get(int(first_align), "middle") if first_align else "middle"
    except Exception:  # noqa: BLE001
        text_anchor = "middle"
    if text_anchor == "start":
        anchor_x = left + int(w * 0.05)
    elif text_anchor == "end":
        anchor_x = left + w - int(w * 0.05)
    else:
        anchor_x = cx

    parts.append(
        f'  <text x="{anchor_x}" y="{start_y}" '
        f'font-family="{_xml_escape(font_family)}" '
        f'font-size="{font_size}" fill="{fill}" '
        f'text-anchor="{text_anchor}">'
    )
    for i, p in enumerate(paragraphs):
        text = _xml_escape(p.text.strip())
        dy = 0 if i == 0 else line_h
        # Bold/italic override per paragraph from first run.
        weight = ""
        style = ""
        try:
            first_r = next((r for r in p.runs), None)
            if first_r is not None:
                if first_r.font.bold:
                    weight = ' font-weight="700"'
                if first_r.font.italic:
                    style = ' font-style="italic"'
        except Exception:  # noqa: BLE001
            pass
        parts.append(
            f'    <tspan x="{anchor_x}" dy="{dy}"{weight}{style}>{text}</tspan>'
        )
    parts.append("  </text>")


def _emit_textbox(shape, parts):
    """Pure TEXT_BOX shape — no geometry to draw, just the text."""
    # Reuse _emit_text_inside which centres text in the bounding box.
    _emit_text_inside(shape, parts)


def _emit_picture(shape, parts, image_lookup):
    """PICTURE shape — base64-embed via data URL so the SVG is self-contained."""
    left, top, w, h = _bounds(shape)
    if left is None:
        return
    href = None
    mime = "image/png"
    try:
        # python-pptx exposes shape.image which gives bytes + content_type.
        img = shape.image
        if img is not None:
            mime = img.content_type or "image/png"
            href = f"data:{mime};base64,{base64.b64encode(img.blob).decode('ascii')}"
    except Exception:  # noqa: BLE001
        # Fallback: walk our pre-built rId lookup.
        try:
            blip = shape.element.find(".//" + qn("a:blip"))
            embed_attr = qn("r:embed")
            rid = blip.get(embed_attr) if blip is not None else None
            if rid and rid in image_lookup:
                blob, mime = image_lookup[rid]
                href = f"data:{mime};base64,{base64.b64encode(blob).decode('ascii')}"
        except Exception:  # noqa: BLE001
            pass
    if href is None:
        # Couldn't get image bytes — emit a placeholder rect so layout stays.
        parts.append(
            f'  <rect x="{left}" y="{top}" width="{w}" height="{h}" '
            f'fill="#e5e7eb" stroke="#9ca3af" stroke-width="1"/>'
        )
        return
    parts.append(
        f'  <image x="{left}" y="{top}" width="{w}" height="{h}" '
        f'href="{href}" preserveAspectRatio="xMidYMid meet"/>'
    )


def _emit_line(shape, parts):
    """LINE / CONNECTOR shape — single straight line from one corner to other."""
    left, top, w, h = _bounds(shape)
    if left is None:
        return
    stroke_col, stroke_w = _resolve_line(shape)
    parts.append(
        f'  <line x1="{left}" y1="{top}" x2="{left + w}" y2="{top + h}" '
        f'stroke="{stroke_col}" stroke-width="{stroke_w}"/>'
    )


def _emit_placeholder_rect(shape, parts, *, hint: str):
    """Generic bounding-box rect — used when the shape kind isn't supported
    yet. Faint grey so the layout shows the right footprint without
    pretending to render."""
    left, top, w, h = _bounds(shape)
    if left is None:
        return
    parts.append(
        f'  <rect x="{left}" y="{top}" width="{w}" height="{h}" '
        f'fill="#f3f4f6" stroke="#d1d5db" stroke-width="3000" '
        f'data-fallback="{_xml_escape(hint)}"/>'
    )


# ---------------------------------------------------------------------------
# Geometry / colour helpers
# ---------------------------------------------------------------------------

def _bounds(shape):
    """Return (left, top, width, height) in source EMU, or (None,)*4 if any
    extent is missing."""
    try:
        left = int(shape.left)
        top = int(shape.top)
        w = int(shape.width)
        h = int(shape.height)
        return left, top, w, h
    except (TypeError, AttributeError):
        return (None, None, None, None)


def _resolve_fill(shape) -> str:
    """Return a `#RRGGBB` colour for the shape's fill, or 'none' if no fill."""
    try:
        fill = shape.fill
        ftype = fill.type
        from pptx.enum.dml import MSO_FILL_TYPE
        if ftype == MSO_FILL_TYPE.SOLID:
            try:
                rgb = fill.fore_color.rgb
                if rgb is not None:
                    return f"#{str(rgb)}"
            except Exception:  # noqa: BLE001
                pass
            # Solid scheme/theme color we can't resolve → neutral grey.
            return "#94a3b8"
        if ftype == MSO_FILL_TYPE.BACKGROUND or ftype is None:
            return "none"
    except Exception:  # noqa: BLE001
        pass
    return "#cbd5e1"


def _resolve_line(shape) -> tuple[str, int]:
    """Return (colour, stroke-width-in-EMU) for the shape's line."""
    default = ("#475569", 6350)  # ~0.5pt in EMU
    try:
        line = shape.line
        try:
            rgb = line.color.rgb
            if rgb is not None:
                col = f"#{str(rgb)}"
            else:
                col = default[0]
        except Exception:  # noqa: BLE001
            col = default[0]
        try:
            w = int(line.width) if line.width is not None else default[1]
        except Exception:  # noqa: BLE001
            w = default[1]
        return col, max(1, w)
    except Exception:  # noqa: BLE001
        return default


def _build_image_lookup(slide) -> dict[str, tuple[bytes, str]]:
    """Build {rId: (blob, content_type)} for all picture relationships."""
    out: dict[str, tuple[bytes, str]] = {}
    try:
        for rel in slide.part.rels.values():
            if "image" in rel.reltype.lower():
                try:
                    out[rel.rId] = (rel.target_part.blob, rel.target_part.content_type)
                except Exception:  # noqa: BLE001
                    continue
    except Exception:  # noqa: BLE001
        pass
    return out


def _xml_escape(s: str) -> str:
    return (
        str(s)
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
        .replace("'", "&apos;")
    )


def _empty_svg(w: int, h: int) -> str:
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}" '
        f'viewBox="0 0 {w} {h}"><rect width="{w}" height="{h}" fill="#fafafa"/></svg>'
    )
