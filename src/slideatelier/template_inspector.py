"""Inspect a master .pptx file: list its layouts and placeholders so designers
can see how to wire layout_map in their template JSON."""
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from rich.console import Console
from rich.table import Table

console = Console()


def inspect_master(master_path: Path) -> None:
    if not master_path.exists():
        console.print(f"[red]File not found:[/red] {master_path}")
        return

    prs = Presentation(str(master_path))

    console.print(f"\n[bold]Master:[/bold] {master_path}")
    console.print(f"[dim]Slide size: {prs.slide_width / 914400:.2f}\" × {prs.slide_height / 914400:.2f}\"[/dim]")
    console.print(f"[dim]{len(prs.slide_layouts)} layouts available[/dim]\n")

    table = Table(show_lines=True)
    table.add_column("Idx", style="cyan", no_wrap=True)
    table.add_column("Layout name", style="bold")
    table.add_column("Placeholders (idx · type · name)", overflow="fold")

    for i, layout in enumerate(prs.slide_layouts):
        phs = []
        for p in layout.placeholders:
            ptype = p.placeholder_format.type
            type_name = ptype.name if hasattr(ptype, "name") else str(ptype)
            # Skip the boilerplate date/footer/page-number placeholders (always there)
            if ptype in (PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER):
                continue
            phs.append(f"  {p.placeholder_format.idx} · {type_name} · {p.name!r}")
        table.add_row(str(i), layout.name, "\n".join(phs) if phs else "[dim](none)[/dim]")

    console.print(table)


def detect_layout_map(master_path: Path) -> dict[str, int]:
    """Try to match our 7 layout types to layouts in the master, by name.
    Falls back to the standard PowerPoint layout indices when no name matches."""
    from .template import DEFAULT_LAYOUT_MAP

    prs = Presentation(str(master_path))
    layout_names = [(i, layout.name.lower()) for i, layout in enumerate(prs.slide_layouts)]

    # Heuristic name patterns — case-insensitive substring match
    patterns = {
        "title": ["title slide", "cover", "title page"],
        "section_divider": ["section header", "section divider", "divider", "section"],
        "content": ["title and content", "content", "1 column", "one column"],
        "bullet_list": ["title and content", "bullet", "list", "content"],
        "two_column": ["two content", "comparison", "2 column", "two column"],
        "key_takeaway": ["title only", "key takeaway", "big idea", "headline", "key insight"],
        "closing": ["closing", "next steps", "thank you", "summary", "title and content"],
    }

    detected = {}
    for our_name, candidates in patterns.items():
        for candidate in candidates:
            for idx, name in layout_names:
                if candidate in name:
                    detected[our_name] = idx
                    break
            if our_name in detected:
                break

    # Fill any gaps with PowerPoint's standard mapping (works for a default master)
    for k, v in DEFAULT_LAYOUT_MAP.items():
        if k not in detected and v < len(prs.slide_layouts):
            detected[k] = v

    return detected
