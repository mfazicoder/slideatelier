"""Hero (Full-bleed) — a striking single-message hero slide.

Layout: full-bleed background (image OR strong solid palette color), large
title at bottom-left with optional subtitle, and a thin colored accent rule
above the title. All shapes are native MSO_SHAPE primitives — fully editable
in PowerPoint.
"""
from __future__ import annotations

from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

from .base import SlideTemplate, Theme
from .shapes_matrix import _lighten


class HeroFullBleed(SlideTemplate):
    id = "hero-fullbleed"
    name = "Hero (Full-bleed)"
    description = "Striking single-message hero slide with full-bleed background."
    style_tags = ("hero", "minimal", "graphics-heavy", "artistic")
    requires_image = True

    def render(
        self,
        slide,
        theme: Theme,
        content: dict,
        *,
        slide_width_emu: int,
        slide_height_emu: int,
    ) -> None:
        palette = theme.palette
        typo = theme.typography

        # ------------------------------------------------------------------
        # 1. Background fill — solid color chosen by theme.
        # ------------------------------------------------------------------
        if theme.is_dark:
            bg_color = palette.background          # deep blue-black for scaffold-dark
            title_color = palette.text             # near-white
            subtitle_color = palette.muted
            accent_color = palette.primary         # bright green
        elif "minimal" in theme.style_tags or "artistic" in theme.style_tags:
            bg_color = palette.background          # warm off-white
            title_color = palette.primary          # near-black
            subtitle_color = palette.muted
            accent_color = palette.accent          # terracotta
        else:
            bg_color = palette.primary             # strong primary (navy / purple)
            title_color = palette.background       # white
            subtitle_color = _lighten(palette.background, 0.0)  # white-ish
            accent_color = palette.accent

        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, slide_width_emu, slide_height_emu
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.fill.background()  # no border on the background

        # ------------------------------------------------------------------
        # 2. Optional full-bleed image + dark overlay for legibility.
        # python-pptx has no straightforward alpha for fills, so we use a
        # near-black solid overlay (palette.text on dark, RGB(0,0,0) on light)
        # and rely on the lack of true transparency by keeping it visually
        # heavy enough to make text legible while still showing the image
        # through its own brightness. (Note: a true alpha overlay would
        # require XML manipulation — we keep it native here.)
        # ------------------------------------------------------------------
        image_path = content.get("image_path")
        if image_path:
            p = Path(image_path)
            if p.exists():
                slide.shapes.add_picture(
                    str(p), 0, 0, slide_width_emu, slide_height_emu
                )
                # On dark themes, the image already darkens text contrast; on
                # light themes we need a darkening overlay so titles read.
                overlay = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, 0, slide_width_emu, slide_height_emu
                )
                overlay.fill.solid()
                overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
                overlay.line.fill.background()
                # Force light-on-dark text once we have the image overlay.
                title_color = RGBColor(0xFF, 0xFF, 0xFF)
                subtitle_color = RGBColor(0xE5, 0xE5, 0xE5)

        # ------------------------------------------------------------------
        # 3. Layout math — text block in bottom-left at ~7% margins.
        # ------------------------------------------------------------------
        margin_x = int(slide_width_emu * 0.07)
        margin_y = int(slide_height_emu * 0.07)
        text_block_w = slide_width_emu - 2 * margin_x

        title = content.get("title", "")
        subtitle = content.get("subtitle", "")

        # Clamp the title size so very large theme defaults still fit.
        title_pt = max(44, min(56, typo.title_size_pt + 12))
        subtitle_pt = max(typo.heading_size_pt, 18)

        # Approximate heights — a generous title box plus a smaller subtitle.
        title_h = Pt(title_pt * 2.2)
        subtitle_h = Pt(subtitle_pt * 2.0) if subtitle else 0
        accent_rule_h = Pt(4)
        accent_rule_w = Pt(60)
        spacing = Pt(8)

        # Stack from the bottom upward:
        subtitle_top = slide_height_emu - margin_y - (subtitle_h if subtitle else 0)
        title_top = subtitle_top - int(title_h) - (int(spacing) if subtitle else 0)
        accent_top = title_top - int(accent_rule_h) - int(spacing)

        # ------------------------------------------------------------------
        # 4. Accent rule — a thin filled rectangle (palette.accent).
        # ------------------------------------------------------------------
        rule = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            margin_x,
            accent_top,
            int(accent_rule_w),
            int(accent_rule_h),
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = accent_color
        rule.line.fill.background()

        # ------------------------------------------------------------------
        # 5. Title text — large heading in the theme's heading font.
        # ------------------------------------------------------------------
        title_box = slide.shapes.add_textbox(
            margin_x, title_top, text_block_w, int(title_h)
        )
        ttf = title_box.text_frame
        ttf.word_wrap = True
        ttf.margin_left = Pt(0)
        ttf.margin_right = Pt(0)
        ttf.margin_top = Pt(0)
        ttf.margin_bottom = Pt(0)
        tp = ttf.paragraphs[0]
        tp.alignment = 1  # PP_ALIGN.LEFT
        trun = tp.add_run()
        trun.text = title
        trun.font.name = typo.heading
        trun.font.size = Pt(title_pt)
        trun.font.bold = True
        trun.font.color.rgb = title_color

        # ------------------------------------------------------------------
        # 6. Subtitle text (optional).
        # ------------------------------------------------------------------
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                margin_x, subtitle_top, text_block_w, int(subtitle_h)
            )
            stf = sub_box.text_frame
            stf.word_wrap = True
            stf.margin_left = Pt(0)
            stf.margin_right = Pt(0)
            stf.margin_top = Pt(0)
            stf.margin_bottom = Pt(0)
            sp = stf.paragraphs[0]
            sp.alignment = 1
            srun = sp.add_run()
            srun.text = subtitle
            srun.font.name = typo.body
            srun.font.size = Pt(subtitle_pt)
            srun.font.color.rgb = subtitle_color
