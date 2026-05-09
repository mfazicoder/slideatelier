"""HubSpoke — central hub + 6 surrounding ovals connected by line connectors.

All native: MSO_SHAPE.OVAL for the hub and spokes, MSO_CONNECTOR.STRAIGHT for
the radial lines (which python-pptx returns as LINE shapes — not FREEFORM).
"""
import math

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Pt

from .base import AssetShape, ShapeRenderContext
from .shapes_matrix import _hex, _lighten


_SPOKE_LABELS = [
    "PEOPLE",
    "PROCESS",
    "DATA",
    "TECH",
    "CULTURE",
    "STRATEGY",
]


class HubSpoke(AssetShape):
    id = "hub-spoke"
    name = "Hub & Spoke"
    description = "Central hub with 6 surrounding nodes connected by straight line connectors."
    style_tags = ("network", "structured", "framework", "consulting")
    aspect_ratio_hint = 1.0

    def render(self, slide, ctx: ShapeRenderContext) -> None:
        theme = ctx.theme
        palette = ctx.palette

        # Geometry: hub diameter ~22% of min(width, height); spokes ~16%.
        bbox_min = min(ctx.width, ctx.height)
        hub_d = int(bbox_min * 0.22)
        spoke_d = int(bbox_min * 0.16)

        cx = ctx.left + ctx.width / 2.0
        cy = ctx.top + ctx.height / 2.0

        # Place spokes on an ellipse so they fill the rect when not square.
        ring_rx = (ctx.width - spoke_d) / 2.0 - int(ctx.width * 0.02)
        ring_ry = (ctx.height - spoke_d) / 2.0 - int(ctx.height * 0.02)

        outline_only = theme.accent_treatment == "outline" or theme.is_dark

        # Compute spoke centres at 6 evenly spaced angles, starting top.
        spoke_centres = []
        for i in range(6):
            theta = -math.pi / 2.0 + i * (2 * math.pi / 6.0)
            sx = cx + ring_rx * math.cos(theta)
            sy = cy + ring_ry * math.sin(theta)
            spoke_centres.append((sx, sy))

        # Draw connectors first (so they sit behind the ovals visually). We
        # draw each connector from the hub's edge toward the spoke's edge by
        # using the ovals' outer points along the radial direction.
        for sx, sy in spoke_centres:
            dx = sx - cx
            dy = sy - cy
            dist = math.hypot(dx, dy) or 1.0
            ux = dx / dist
            uy = dy / dist

            # Start: hub edge along the radial direction.
            start_x = int(round(cx + ux * hub_d / 2.0))
            start_y = int(round(cy + uy * hub_d / 2.0))
            # End: spoke edge along the (reverse) radial direction.
            end_x = int(round(sx - ux * spoke_d / 2.0))
            end_y = int(round(sy - uy * spoke_d / 2.0))

            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y
            )
            conn.line.color.rgb = palette.muted if not outline_only else palette.primary
            conn.line.width = theme.line_weight_emu

        # Draw spoke ovals.
        for i, (sx, sy) in enumerate(spoke_centres):
            left = int(round(sx - spoke_d / 2.0))
            top = int(round(sy - spoke_d / 2.0))
            oval = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, left, top, spoke_d, spoke_d
            )
            oval.fill.solid()
            if outline_only:
                oval.fill.fore_color.rgb = palette.background
                oval.line.color.rgb = palette.primary
                lbl_color = palette.text
            else:
                # Alternate primary/accent tints.
                if i % 2 == 0:
                    oval.fill.fore_color.rgb = _lighten(palette.primary, 0.78)
                else:
                    oval.fill.fore_color.rgb = _lighten(palette.accent, 0.70)
                oval.line.color.rgb = palette.primary
                lbl_color = palette.text
            oval.line.width = theme.line_weight_emu

            tf = oval.text_frame
            tf.margin_left = Pt(2)
            tf.margin_right = Pt(2)
            tf.margin_top = Pt(2)
            tf.margin_bottom = Pt(2)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = 2  # CENTER
            run = p.add_run()
            run.text = _SPOKE_LABELS[i]
            run.font.size = Pt(theme.typography.caption_size_pt)
            run.font.name = theme.typography.heading
            run.font.bold = False
            run.font.color.rgb = lbl_color

        # Hub oval (drawn last so it stacks on top of connectors).
        hub_left = int(round(cx - hub_d / 2.0))
        hub_top = int(round(cy - hub_d / 2.0))
        hub = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, hub_left, hub_top, hub_d, hub_d
        )
        hub.fill.solid()
        if outline_only:
            if theme.is_dark:
                hub.fill.fore_color.rgb = palette.primary
                hub_text_color = palette.background
            else:
                hub.fill.fore_color.rgb = palette.background
                hub_text_color = palette.primary
            hub.line.color.rgb = palette.primary
        else:
            hub.fill.fore_color.rgb = palette.primary
            hub.line.color.rgb = palette.primary
            hub_text_color = palette.background
        hub.line.width = theme.line_weight_emu

        htf = hub.text_frame
        htf.margin_left = Pt(2)
        htf.margin_right = Pt(2)
        htf.margin_top = Pt(2)
        htf.margin_bottom = Pt(2)
        hp = htf.paragraphs[0]
        hp.alignment = 2  # CENTER
        hrun = hp.add_run()
        hrun.text = "CORE"
        hrun.font.size = Pt(theme.typography.body_size_pt)
        hrun.font.name = theme.typography.heading
        hrun.font.bold = True
        hrun.font.color.rgb = hub_text_color


# ---------------------------------------------------------------------------
# SVG rendering (Web Deck publishing — Sprint J.B)
# ---------------------------------------------------------------------------

def _render_svg_hub_spoke(
    self: "HubSpoke", ctx: ShapeRenderContext, width_px: int, height_px: int
) -> str:
    theme = ctx.theme
    palette = ctx.palette

    # Mirror render() geometry but in pixel space.
    bbox_min = min(width_px, height_px)
    hub_d = int(bbox_min * 0.22)
    spoke_d = int(bbox_min * 0.16)

    cx = width_px / 2.0
    cy = height_px / 2.0

    ring_rx = (width_px - spoke_d) / 2.0 - int(width_px * 0.02)
    ring_ry = (height_px - spoke_d) / 2.0 - int(height_px * 0.02)

    outline_only = theme.accent_treatment == "outline" or theme.is_dark

    spoke_centres = []
    for i in range(6):
        theta = -math.pi / 2.0 + i * (2 * math.pi / 6.0)
        sx = cx + ring_rx * math.cos(theta)
        sy = cy + ring_ry * math.sin(theta)
        spoke_centres.append((sx, sy))

    heading_font = theme.typography.heading
    cap_pt = theme.typography.caption_size_pt
    body_pt = theme.typography.body_size_pt

    parts: list[str] = []
    parts.append(
        f'<svg xmlns="http://www.w3.org/2000/svg" '
        f'width="{width_px}" height="{height_px}" '
        f'viewBox="0 0 {width_px} {height_px}">'
    )
    # Background
    parts.append(
        f'  <rect x="0" y="0" width="{width_px}" height="{height_px}" '
        f'fill="{_hex(palette.background)}" />'
    )

    # Connector lines (drawn first so spokes/hub stack on top).
    line_color = (
        _hex(palette.muted) if not outline_only else _hex(palette.primary)
    )
    for sx, sy in spoke_centres:
        dx = sx - cx
        dy = sy - cy
        dist = math.hypot(dx, dy) or 1.0
        ux = dx / dist
        uy = dy / dist
        x1 = cx + ux * hub_d / 2.0
        y1 = cy + uy * hub_d / 2.0
        x2 = sx - ux * spoke_d / 2.0
        y2 = sy - uy * spoke_d / 2.0
        parts.append(
            f'  <line x1="{x1:.2f}" y1="{y1:.2f}" '
            f'x2="{x2:.2f}" y2="{y2:.2f}" '
            f'stroke="{line_color}" stroke-width="2" />'
        )

    # Spoke circles.
    spoke_r = spoke_d / 2.0
    for i, (sx, sy) in enumerate(spoke_centres):
        if outline_only:
            spoke_fill = _hex(palette.background)
            spoke_stroke = _hex(palette.primary)
            lbl_color = _hex(palette.text)
        else:
            if i % 2 == 0:
                spoke_fill = _hex(_lighten(palette.primary, 0.78))
            else:
                spoke_fill = _hex(_lighten(palette.accent, 0.70))
            spoke_stroke = _hex(palette.primary)
            lbl_color = _hex(palette.text)
        parts.append(
            f'  <circle cx="{sx:.2f}" cy="{sy:.2f}" r="{spoke_r:.2f}" '
            f'fill="{spoke_fill}" stroke="{spoke_stroke}" stroke-width="1.5" />'
        )
        parts.append(
            f'  <text x="{sx:.2f}" y="{sy + cap_pt / 3.0:.2f}" '
            f'font-family="{heading_font}" font-size="{cap_pt}" '
            f'fill="{lbl_color}" text-anchor="middle">{_SPOKE_LABELS[i]}</text>'
        )

    # Hub circle (drawn last, on top).
    hub_r = hub_d / 2.0
    if outline_only:
        if theme.is_dark:
            hub_fill = _hex(palette.primary)
            hub_text_color = _hex(palette.background)
        else:
            hub_fill = _hex(palette.background)
            hub_text_color = _hex(palette.primary)
        hub_stroke = _hex(palette.primary)
    else:
        hub_fill = _hex(palette.primary)
        hub_stroke = _hex(palette.primary)
        hub_text_color = _hex(palette.background)
    parts.append(
        f'  <circle cx="{cx:.2f}" cy="{cy:.2f}" r="{hub_r:.2f}" '
        f'fill="{hub_fill}" stroke="{hub_stroke}" stroke-width="1.5" />'
    )
    parts.append(
        f'  <text x="{cx:.2f}" y="{cy + body_pt / 3.0:.2f}" '
        f'font-family="{heading_font}" font-size="{body_pt}" '
        f'font-weight="bold" fill="{hub_text_color}" '
        f'text-anchor="middle">CORE</text>'
    )

    parts.append("</svg>")
    return "\n".join(parts)


HubSpoke.render_svg = _render_svg_hub_spoke  # type: ignore[attr-defined]
