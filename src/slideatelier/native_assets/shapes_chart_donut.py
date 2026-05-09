"""ChartDonut — native PowerPoint doughnut chart.

Uses slide.shapes.add_chart() with XL_CHART_TYPE.DOUGHNUT. Each data point
gets a different palette color so the slices are visually distinct under any
theme.
"""
import math

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from .base import AssetShape, ShapeRenderContext
from .shapes_matrix import _lighten


def _hex(color) -> str:
    """RGBColor → '#RRGGBB' for SVG."""
    return "#{:02X}{:02X}{:02X}".format(color[0], color[1], color[2])


# Categories + values shared between PPT render() and SVG render_svg().
_CATEGORIES = ("A", "B", "C", "D")
_VALUES = (35, 25, 22, 18)


class ChartDonut(AssetShape):
    id = "chart-donut"
    name = "Donut Chart"
    description = "Native doughnut chart with 4 categorical slices, themed colors per slice."
    style_tags = ("data", "chart", "share", "contemporary")
    aspect_ratio_hint = 1.0

    def render(self, slide, ctx: ShapeRenderContext) -> None:
        theme = ctx.theme
        palette = ctx.palette

        cd = CategoryChartData()
        cd.categories = list(_CATEGORIES)
        cd.add_series("Share", _VALUES)

        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT,
            ctx.left,
            ctx.top,
            ctx.width,
            ctx.height,
            cd,
        )
        chart = chart_shape.chart
        chart.has_title = False
        chart.has_legend = True
        try:
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False
        except Exception:
            pass

        # Theme palette per slice. Use 4 distinct shades:
        slice_colors = [
            palette.primary,
            palette.accent,
            _lighten(palette.primary, 0.55),
            _lighten(palette.accent, 0.55),
        ]

        plot = chart.plots[0]
        try:
            series = plot.series[0]
            for i, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                fill.fore_color.rgb = slice_colors[i % len(slice_colors)]
                line = point.format.line
                line.color.rgb = palette.background
        except Exception:
            # Fallback: just color the series uniformly.
            try:
                fill = plot.series[0].format.fill
                fill.solid()
                fill.fore_color.rgb = palette.primary
            except Exception:
                pass


# ---------------------------------------------------------------------------
# SVG rendering (Web Deck publishing — Sprint J.C)
# ---------------------------------------------------------------------------

def _polar(cx: float, cy: float, r: float, angle_rad: float) -> tuple[float, float]:
    return cx + r * math.cos(angle_rad), cy + r * math.sin(angle_rad)


def _donut_slice_path(
    cx: float, cy: float, r_outer: float, r_inner: float,
    start_rad: float, end_rad: float,
) -> str:
    """Build an SVG path for a donut slice using elliptical arc commands.

    The slice is bounded by two arcs (outer + inner) and two radial line
    segments. SVG arc command 'A rx ry x-rot large-arc sweep x y' is used for
    both arcs — the principle "circles stay circles" preserved via native
    arc primitives, no polyline approximation.
    """
    sweep = end_rad - start_rad
    large_arc = 1 if sweep > math.pi else 0

    x_outer_start, y_outer_start = _polar(cx, cy, r_outer, start_rad)
    x_outer_end, y_outer_end = _polar(cx, cy, r_outer, end_rad)
    x_inner_end, y_inner_end = _polar(cx, cy, r_inner, end_rad)
    x_inner_start, y_inner_start = _polar(cx, cy, r_inner, start_rad)

    # M outer-start
    # A r r 0 large 1 outer-end       (outer arc clockwise)
    # L inner-end
    # A r r 0 large 0 inner-start     (inner arc anti-clockwise)
    # Z
    return (
        f"M {x_outer_start:.3f} {y_outer_start:.3f} "
        f"A {r_outer:.3f} {r_outer:.3f} 0 {large_arc} 1 "
        f"{x_outer_end:.3f} {y_outer_end:.3f} "
        f"L {x_inner_end:.3f} {y_inner_end:.3f} "
        f"A {r_inner:.3f} {r_inner:.3f} 0 {large_arc} 0 "
        f"{x_inner_start:.3f} {y_inner_start:.3f} Z"
    )


def _render_svg_chart_donut(
    self: "ChartDonut", ctx: ShapeRenderContext, width_px: int, height_px: int
) -> str:
    """Inline SVG mirror of ChartDonut — emits <path> arc segments.

    Uses SVG elliptical-arc 'A' commands so circular curves stay circular.
    """
    theme = ctx.theme
    palette = ctx.palette

    values = list(_VALUES)
    total = float(sum(values)) or 1.0

    # Geometry — donut centered, radius 40% of the smaller dimension; hole
    # radius is 55% of the outer.
    cx = width_px / 2.0
    cy = height_px / 2.0
    r_outer = min(width_px, height_px) * 0.40
    r_inner = r_outer * 0.55

    slice_colors = [
        palette.primary,
        palette.accent,
        _lighten(palette.primary, 0.55),
        _lighten(palette.accent, 0.55),
    ]
    stroke_color = _hex(palette.background)

    parts: list[str] = []
    parts.append(
        f'<svg xmlns="http://www.w3.org/2000/svg" '
        f'width="{width_px}" height="{height_px}" '
        f'viewBox="0 0 {width_px} {height_px}">'
    )
    # Background
    parts.append(
        f'  <rect x="0" y="0" width="{width_px}" height="{height_px}" '
        f'fill="{_hex(palette.background)}" />'
    )

    # Start at -90° (12 o'clock), sweep clockwise.
    start_angle = -math.pi / 2.0
    for i, v in enumerate(values):
        sweep = (v / total) * 2.0 * math.pi
        end_angle = start_angle + sweep
        d = _donut_slice_path(cx, cy, r_outer, r_inner, start_angle, end_angle)
        fill = _hex(slice_colors[i % len(slice_colors)])
        parts.append(
            f'  <path d="{d}" fill="{fill}" '
            f'stroke="{stroke_color}" stroke-width="1" />'
        )
        start_angle = end_angle

    # Center label — total value.
    body_font = theme.typography.body
    body_pt = theme.typography.body_size_pt
    parts.append(
        f'  <text x="{cx:.2f}" y="{cy + body_pt / 3:.2f}" '
        f'font-family="{body_font}" font-size="{body_pt}" '
        f'fill="{_hex(palette.text)}" text-anchor="middle">'
        f'{int(total)}</text>'
    )

    parts.append("</svg>")
    return "\n".join(parts)


ChartDonut.render_svg = _render_svg_chart_donut  # type: ignore[attr-defined]
