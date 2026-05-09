"""Funnel — a 4-stage narrowing funnel diagram.

Each stage is rendered as a native MSO_SHAPE.TRAPEZOID primitive (or
ROUNDED_RECTANGLE / RECTANGLE fallback when corner_radius_pct demands it).
The four bars are stacked vertically, progressively narrower, centered, with
captions to the right.
"""
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

from .base import AssetShape, ShapeRenderContext


# Stage data: (label, caption)
_STAGES = [
    ("Stage 1", "Awareness"),
    ("Stage 2", "Consideration"),
    ("Stage 3", "Decision"),
    ("Stage 4", "Action"),
]
_WIDTH_PCTS = [1.00, 0.80, 0.60, 0.40]


class Funnel(AssetShape):
    id = "funnel"
    name = "Funnel"
    description = "A 4-stage narrowing funnel for marketing/conversion frameworks (Awareness → Action)."
    style_tags = ("process", "narrowing", "consulting", "marketing")
    aspect_ratio_hint = 1.2

    def render(self, slide, ctx: ShapeRenderContext) -> None:
        theme = ctx.theme
        palette = ctx.palette

        n_stages = len(_STAGES)

        # Reserve ~32% of the bounding rect on the right side for captions.
        caption_w = int(ctx.width * 0.30)
        funnel_area_w = ctx.width - caption_w
        funnel_left = ctx.left
        funnel_top = ctx.top

        # Vertical layout — small gap between stages.
        gap_emu = int(ctx.height * 0.02)
        total_gap = gap_emu * (n_stages - 1)
        stage_h = (ctx.height - total_gap) // n_stages

        # Decide shape kind. TRAPEZOID gives the canonical funnel look but
        # python-pptx's TRAPEZOID renders narrower at top, wider at bottom — for
        # a funnel we want narrowing top→bottom, so we instead use stacked
        # rectangles whose widths shrink. (TRAPEZOID is non-mirrorable without
        # rotation and rotation breaks text orientation.) Rounded variant when
        # the theme asks for it.
        rect_kind = (
            MSO_SHAPE.ROUNDED_RECTANGLE
            if theme.corner_radius_pct > 0
            else MSO_SHAPE.RECTANGLE
        )

        # Choose fills/lines per accent_treatment + dark mode.
        # For "fill": gradient-ish step from primary (top) → accent (bottom),
        # implemented as discrete blended colors per stage.
        # For "outline": background fills, primary outline.
        if theme.accent_treatment == "outline" or theme.is_dark:
            stage_fills = [palette.background] * n_stages
            line_color = palette.primary
            text_on_stage = palette.text
            use_inverted_label = False
        else:
            stage_fills = [
                _blend(palette.primary, palette.accent, i / (n_stages - 1))
                for i in range(n_stages)
            ]
            line_color = palette.muted
            text_on_stage = palette.background  # white-ish on filled stages
            use_inverted_label = True

        for i, (label, caption) in enumerate(_STAGES):
            w = int(funnel_area_w * _WIDTH_PCTS[i])
            # center horizontally within the funnel area
            left = funnel_left + (funnel_area_w - w) // 2
            top = funnel_top + i * (stage_h + gap_emu)

            stage = slide.shapes.add_shape(rect_kind, left, top, w, stage_h)
            stage.fill.solid()
            stage.fill.fore_color.rgb = stage_fills[i]
            stage.line.color.rgb = line_color
            stage.line.width = theme.line_weight_emu

            tf = stage.text_frame
            tf.margin_left = Pt(8)
            tf.margin_right = Pt(8)
            tf.margin_top = Pt(4)
            tf.margin_bottom = Pt(4)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = 2  # PP_ALIGN.CENTER  (2 == CENTER per python-pptx enum int)
            run = p.add_run()
            run.text = label
            run.font.size = Pt(theme.typography.body_size_pt)
            run.font.name = theme.typography.heading
            run.font.bold = True
            run.font.color.rgb = (
                text_on_stage if use_inverted_label else palette.text
            )

            # Caption to the right of the funnel area (one per stage).
            cap = slide.shapes.add_textbox(
                funnel_left + funnel_area_w + int(ctx.width * 0.01),
                top,
                caption_w - int(ctx.width * 0.01),
                stage_h,
            )
            ctf = cap.text_frame
            ctf.word_wrap = True
            ctf.margin_left = Pt(4)
            ctf.margin_right = Pt(4)
            ctf.margin_top = Pt(0)
            ctf.margin_bottom = Pt(0)
            cp = ctf.paragraphs[0]
            cp.alignment = 1  # PP_ALIGN.LEFT
            crun = cp.add_run()
            crun.text = caption
            crun.font.size = Pt(theme.typography.caption_size_pt)
            crun.font.name = theme.typography.body
            crun.font.color.rgb = palette.muted

    def render_svg(self, ctx: ShapeRenderContext, width_px: int, height_px: int) -> str:
        """Inline SVG mirroring render(): 4 trapezoid <polygon>s + step captions.

        Although the PPT version uses RECTANGLEs whose widths shrink (because
        python-pptx's TRAPEZOID is wider at the bottom), we are free in SVG to
        draw true narrowing trapezoids — and we should, because that is what a
        funnel visually IS. Each step is a <polygon> with 4 points.
        """
        theme = ctx.theme
        palette = ctx.palette

        n_stages = len(_STAGES)

        caption_w = int(width_px * 0.30)
        funnel_area_w = width_px - caption_w
        funnel_left = 0
        funnel_top = 0

        gap_px = int(height_px * 0.02)
        total_gap = gap_px * (n_stages - 1)
        stage_h = (height_px - total_gap) // n_stages

        if theme.accent_treatment == "outline" or theme.is_dark:
            stage_fills = [palette.background] * n_stages
            line_color = palette.primary
            text_on_stage = palette.text
            use_inverted_label = False
        else:
            stage_fills = [
                _blend(palette.primary, palette.accent, i / (n_stages - 1))
                for i in range(n_stages)
            ]
            line_color = palette.muted
            text_on_stage = palette.background
            use_inverted_label = True

        body_pt = theme.typography.body_size_pt
        cap_pt = theme.typography.caption_size_pt
        heading_font = theme.typography.heading
        body_font = theme.typography.body

        parts: list[str] = []
        parts.append(
            f'<svg xmlns="http://www.w3.org/2000/svg" '
            f'width="{width_px}" height="{height_px}" '
            f'viewBox="0 0 {width_px} {height_px}">'
        )
        parts.append(
            f'  <rect x="0" y="0" width="{width_px}" height="{height_px}" '
            f'fill="{_hex(palette.background)}" />'
        )

        for i, (label, caption) in enumerate(_STAGES):
            top_w = int(funnel_area_w * _WIDTH_PCTS[i])
            # Bottom width is the next step's width (or 50% of top for last).
            if i < n_stages - 1:
                bottom_w = int(funnel_area_w * _WIDTH_PCTS[i + 1])
            else:
                bottom_w = max(int(top_w * 0.5), 1)

            top_y = funnel_top + i * (stage_h + gap_px)
            bottom_y = top_y + stage_h

            top_left_x = funnel_left + (funnel_area_w - top_w) // 2
            top_right_x = top_left_x + top_w
            bottom_left_x = funnel_left + (funnel_area_w - bottom_w) // 2
            bottom_right_x = bottom_left_x + bottom_w

            points = (
                f"{top_left_x},{top_y} "
                f"{top_right_x},{top_y} "
                f"{bottom_right_x},{bottom_y} "
                f"{bottom_left_x},{bottom_y}"
            )
            fill = _hex(stage_fills[i])
            stroke = _hex(line_color)
            parts.append(
                f'  <polygon points="{points}" '
                f'fill="{fill}" stroke="{stroke}" stroke-width="1" />'
            )

            label_color = _hex(text_on_stage if use_inverted_label else palette.text)
            cx = funnel_left + funnel_area_w // 2
            cy = top_y + stage_h // 2 + body_pt // 3
            parts.append(
                f'  <text x="{cx}" y="{cy}" '
                f'font-family="{heading_font}" font-size="{body_pt}" '
                f'font-weight="bold" fill="{label_color}" '
                f'text-anchor="middle">{label}</text>'
            )

            cap_x = funnel_left + funnel_area_w + int(width_px * 0.01)
            cap_y = top_y + stage_h // 2 + cap_pt // 3
            parts.append(
                f'  <text x="{cap_x}" y="{cap_y}" '
                f'font-family="{body_font}" font-size="{cap_pt}" '
                f'fill="{_hex(palette.muted)}">{caption}</text>'
            )

        parts.append("</svg>")
        return "\n".join(parts)


def _hex(color) -> str:
    """RGBColor → '#RRGGBB' for SVG."""
    return "#{:02X}{:02X}{:02X}".format(color[0], color[1], color[2])


def _blend(a: RGBColor, b: RGBColor, t: float) -> RGBColor:
    """Linear blend of two RGBColors. t=0 → a, t=1 → b."""
    t = max(0.0, min(1.0, t))
    r = int(a[0] + (b[0] - a[0]) * t)
    g = int(a[1] + (b[1] - a[1]) * t)
    bb = int(a[2] + (b[2] - a[2]) * t)
    return RGBColor(r, g, bb)
