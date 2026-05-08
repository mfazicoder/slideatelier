"""Materialize the native asset framework as concrete .pptx files + LibraryAsset
catalog entries.

For every (shape × theme) and (template × theme) combination we:
    1. Create a fresh standalone Presentation.
    2. Render the shape (or template) onto a single blank slide using the theme.
    3. Save it under ``output_dir/native/<id>.pptx``.
    4. Optionally render a thumbnail PNG to ``output_dir/thumbnails/<flat-id>.png``
       using the existing LibreOffice → PDF → PNG pipeline. We skip cleanly
       when LibreOffice ('soffice') isn't on PATH.
    5. Build a LibraryAsset entry with style_tags merged from the shape/template
       and the theme.

The resulting list of LibraryAsset entries is intended to be merged into the
existing 539-asset catalog by the ``atelier library build-natives`` CLI command.
"""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches

from ..library import LibraryAsset, slugify
from .base import AssetShape, SlideTemplate, Theme
from .registry import get_registry


# Standard 16:9 widescreen — matches the rest of slideAtelier.
_SLIDE_WIDTH_IN = 13.333
_SLIDE_HEIGHT_IN = 7.5

_NATIVE_FRAMEWORK_CATEGORY = "Native — Frameworks"
_NATIVE_TEMPLATE_CATEGORY = "Native — Templates"


def _flat_id(asset_id: str) -> str:
    """Same flattening rule used by the rest of the library/web code: replace
    slashes with double underscores so the id maps cleanly onto a filename."""
    return asset_id.replace("/", "__")


def _make_blank_presentation() -> Presentation:
    """Build a fresh 16:9 presentation with one truly blank slide."""
    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_WIDTH_IN)
    prs.slide_height = Inches(_SLIDE_HEIGHT_IN)
    # Layout 6 = blank in python-pptx's default master.
    prs.slides.add_slide(prs.slide_layouts[6])
    return prs


def _template_sample_content(template: SlideTemplate, theme: Theme) -> dict:
    """Minimal placeholder content for each known template, so we can render
    a representative preview without external inputs."""
    if template.id == "hero-fullbleed":
        return {"title": theme.name, "subtitle": "Sample"}
    if template.id == "exec-summary":
        return {
            "title": theme.name,
            "takeaway": "Sample takeaway",
            "columns": [
                {"heading": "A", "bullets": ["x"]},
                {"heading": "B", "bullets": ["y"]},
                {"heading": "C", "bullets": ["z"]},
            ],
        }
    # Generic fallback — at least supply a title so templates that read it
    # don't crash.
    return {"title": theme.name}


def _render_shape_preview(shape: AssetShape, theme: Theme, out_path: Path) -> int:
    """Render a single shape onto a blank slide using the shape's own
    full-slide convenience helper. Returns the resulting shape count."""
    prs = _make_blank_presentation()
    slide = prs.slides[0]
    shape.render_into_full_slide(
        slide,
        theme,
        slide_width_emu=prs.slide_width,
        slide_height_emu=prs.slide_height,
    )
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return len(list(slide.shapes))


def _render_template_preview(template: SlideTemplate, theme: Theme, out_path: Path) -> int:
    """Render a slide template onto a blank slide with placeholder content."""
    prs = _make_blank_presentation()
    slide = prs.slides[0]
    template.render(
        slide,
        theme,
        _template_sample_content(template, theme),
        slide_width_emu=prs.slide_width,
        slide_height_emu=prs.slide_height,
    )
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return len(list(slide.shapes))


def _merge_tags(*sources: tuple[str, ...] | list[str] | None) -> list[str]:
    """De-duplicate while preserving order across heterogenous tag sources."""
    seen: set[str] = set()
    out: list[str] = []
    for src in sources:
        if not src:
            continue
        for t in src:
            if t and t not in seen:
                seen.add(t)
                out.append(t)
    return out


def _generate_thumbnail(
    pptx_path: Path,
    thumbnails_dir: Path,
    flat_id: str,
    *,
    dpi: int = 100,
    max_width_px: int = 800,
) -> bool:
    """Best-effort thumbnail generation for a single .pptx using LibreOffice.
    Returns True if a PNG was written, False otherwise (no soffice, conversion
    failed, etc.). We deliberately swallow exceptions so the catalog build can
    proceed on machines without LibreOffice."""
    if shutil.which("soffice") is None:
        return False
    try:
        # Lazy import: pdf2image / poppler aren't required when soffice is absent.
        from ..library_thumbnails import _convert_pptx_to_pdf, _split_pdf_to_pngs
        import tempfile

        thumbnails_dir.mkdir(parents=True, exist_ok=True)
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            pdf = _convert_pptx_to_pdf(pptx_path, tmp)
            if pdf is None:
                return False
            pages = _split_pdf_to_pngs(pdf, dpi=dpi)
            if not pages:
                return False
            img = pages[0]
            if img.width > max_width_px:
                ratio = max_width_px / img.width
                img = img.resize((max_width_px, int(img.height * ratio)))
            out_png = thumbnails_dir / f"{flat_id}.png"
            img.save(out_png, "PNG", optimize=True)
            return out_png.exists()
    except Exception:  # noqa: BLE001 — thumbnails are optional, never fatal.
        return False


def build_native_catalog(
    output_dir: Path,
    *,
    generate_thumbnails: bool = True,
) -> list[LibraryAsset]:
    """Build .pptx files + LibraryAsset entries for every (shape × theme) and
    (template × theme) combination. See module docstring for full pipeline."""
    output_dir = Path(output_dir)
    native_dir = output_dir / "native"
    thumbnails_dir = output_dir / "thumbnails"

    reg = get_registry()
    themes = reg.themes()
    shapes = reg.shapes()
    templates = reg.templates()

    entries: list[LibraryAsset] = []
    framework_slug = slugify(_NATIVE_FRAMEWORK_CATEGORY)
    template_slug = slugify(_NATIVE_TEMPLATE_CATEGORY)

    # ---------- shapes × themes ----------
    for shape in shapes:
        for theme in themes:
            asset_id = f"native/{shape.id}/{theme.id}"
            file_stem = f"{shape.id}__{theme.id}"
            pptx_path = native_dir / f"{file_stem}.pptx"
            shape_count = _render_shape_preview(shape, theme, pptx_path)

            has_thumb = False
            if generate_thumbnails:
                has_thumb = _generate_thumbnail(
                    pptx_path, thumbnails_dir, _flat_id(asset_id)
                )

            entries.append(
                LibraryAsset(
                    id=asset_id,
                    file_path=str(pptx_path.resolve()),
                    slide_index=1,
                    category=_NATIVE_FRAMEWORK_CATEGORY,
                    category_slug=framework_slug,
                    file_stem=file_stem,
                    shape_count=shape_count,
                    text_snippets=[],
                    slide_width_inches=_SLIDE_WIDTH_IN,
                    slide_height_inches=_SLIDE_HEIGHT_IN,
                    has_thumbnail=has_thumb,
                    notes=f"Native shape · {theme.name}",
                    style_tags=_merge_tags(
                        shape.style_tags, theme.style_tags, ("native",)
                    ),
                )
            )

    # ---------- templates × themes ----------
    for template in templates:
        for theme in themes:
            asset_id = f"native/{template.id}/{theme.id}"
            file_stem = f"template__{template.id}__{theme.id}"
            pptx_path = native_dir / f"{file_stem}.pptx"
            shape_count = _render_template_preview(template, theme, pptx_path)

            has_thumb = False
            if generate_thumbnails:
                has_thumb = _generate_thumbnail(
                    pptx_path, thumbnails_dir, _flat_id(asset_id)
                )

            entries.append(
                LibraryAsset(
                    id=asset_id,
                    file_path=str(pptx_path.resolve()),
                    slide_index=1,
                    category=_NATIVE_TEMPLATE_CATEGORY,
                    category_slug=template_slug,
                    file_stem=file_stem,
                    shape_count=shape_count,
                    text_snippets=[],
                    slide_width_inches=_SLIDE_WIDTH_IN,
                    slide_height_inches=_SLIDE_HEIGHT_IN,
                    has_thumbnail=has_thumb,
                    notes=f"Native template · {theme.name}",
                    style_tags=_merge_tags(
                        template.style_tags, theme.style_tags, ("native",)
                    ),
                )
            )

    return entries
