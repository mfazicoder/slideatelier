"""TimelineDots — horizontal line with 5 numbered milestones, alternating
above/below callouts for the labels.

Native: MSO_CONNECTOR.STRAIGHT for the spine, MSO_SHAPE.OVAL for milestones,
textboxes for the callouts.
"""
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Pt

from .base import AssetShape, ShapeRenderContext
from .shapes_matrix import _lighten


def _hex(color) -> str:
    """RGBColor → '#RRGGBB' for SVG."""
    return "#{:02X}{:02X}{:02X}".format(color[0], color[1], color[2])


_MILESTONES = [
    ("Q1 '25", "Discovery"),
    ("Q2 '25", "Pilot"),
    ("Q3 '25", "Scale"),
    ("Q4 '25", "Optimize"),
    ("Q1 '26", "Renew"),
]


class TimelineDots(AssetShape):
    id = "timeline-dots"
    name = "Timeline Dots"
    description = "Horizontal timeline with 5 numbered milestones and alternating callouts."
    style_tags = ("process", "sequential", "modern", "contemporary")
    aspect_ratio_hint = 2.4

    def render(self, slide, ctx: ShapeRenderContext) -> None:
        theme = ctx.theme
        palette = ctx.palette

        outline_only = theme.accent_treatment == "outline" or theme.is_dark

        n = len(_MILESTONES)

        # Vertical layout: top callout band, spine line, bottom callout band.
        # Spine sits at vertical center.
        spine_y = ctx.top + ctx.height // 2

        # Dot diameter ~ 9% of bbox height (or 6% of width, whichever fits).
        dot_d = int(min(ctx.height * 0.18, ctx.width * 0.06))
        dot_d = max(dot_d, 200000)  # min ~22pt for legibility

        # Compute x positions for the n dots, evenly spaced inside a margin.
        margin_x = int(ctx.width * 0.05)
        usable_w = ctx.width - 2 * margin_x
        if n > 1:
            stride = usable_w // (n - 1)
        else:
            stride = 0
        first_cx = ctx.left + margin_x

        dot_centres = [first_cx + i * stride for i in range(n)]

        # Spine connector — extend slightly past first/last dots for visual
        # finish.
        spine_left = ctx.left + int(margin_x * 0.5)
        spine_right = ctx.left + ctx.width - int(margin_x * 0.5)
        spine = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            spine_left,
            spine_y,
            spine_right,
            spine_y,
        )
        spine.line.color.rgb = palette.muted if not outline_only else palette.primary
        spine.line.width = max(theme.line_weight_emu, 19050)  # at least ~1.5pt

        # Callout band heights (above and below the spine).
        callout_h = (ctx.height - dot_d) // 2 - int(ctx.height * 0.02)

        for i, (date, label) in enumerate(_MILESTONES):
            cx = dot_centres[i]
            # Dot oval centred on the spine.
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                cx - dot_d // 2,
                spine_y - dot_d // 2,
                dot_d,
                dot_d,
            )
            dot.fill.solid()
            if outline_only:
                if theme.is_dark:
                    dot.fill.fore_color.rgb = palette.primary
                    dot_text_color = palette.background
                else:
                    dot.fill.fore_color.rgb = palette.background
                    dot_text_color = palette.primary
                dot.line.color.rgb = palette.primary
            else:
                dot.fill.fore_color.rgb = palette.primary
                dot.line.color.rgb = palette.primary
                dot_text_color = palette.background
            dot.line.width = theme.line_weight_emu

            dtf = dot.text_frame
            dtf.margin_left = Pt(0)
            dtf.margin_right = Pt(0)
            dtf.margin_top = Pt(0)
            dtf.margin_bottom = Pt(0)
            dp = dtf.paragraphs[0]
            dp.alignment = 2  # CENTER
            drun = dp.add_run()
            drun.text = str(i + 1)
            drun.font.size = Pt(theme.typography.caption_size_pt)
            drun.font.name = theme.typography.heading
            drun.font.bold = True
            drun.font.color.rgb = dot_text_color

            # Callout — alternate above/below.
            callout_w = stride - int(stride * 0.15) if n > 1 else int(ctx.width * 0.18)
            callout_w = max(callout_w, dot_d * 3)
            cleft = cx - callout_w // 2
            if i % 2 == 0:
                # Above the spine
                ctop = ctx.top
                ch = callout_h
            else:
                # Below the spine
                ctop = spine_y + dot_d // 2 + int(ctx.height * 0.02)
                ch = callout_h

            cbox = slide.shapes.add_textbox(cleft, ctop, callout_w, ch)
            ctf = cbox.text_frame
            ctf.word_wrap = True
            ctf.margin_left = Pt(2)
            ctf.margin_right = Pt(2)
            ctf.margin_top = Pt(2)
            ctf.margin_bottom = Pt(2)

            # First paragraph: date
            cp = ctf.paragraphs[0]
            cp.alignment = 2  # CENTER
            crun = cp.add_run()
            crun.text = date
            crun.font.size = Pt(theme.typography.body_size_pt)
            crun.font.name = theme.typography.heading
            crun.font.bold = True
            crun.font.color.rgb = palette.text

            # Second paragraph: label
            lp = ctf.add_paragraph()
            lp.alignment = 2  # CENTER
            lrun = lp.add_run()
            lrun.text = label
            lrun.font.size = Pt(theme.typography.caption_size_pt)
            lrun.font.name = theme.typography.body
            lrun.font.color.rgb = palette.muted


# ---------------------------------------------------------------------------
# SVG rendering (Web Deck publishing — Sprint J.C)
# ---------------------------------------------------------------------------

def _render_svg_timeline_dots(
    self: "TimelineDots", ctx: ShapeRenderContext, width_px: int, height_px: int
) -> str:
    """Inline SVG mirror of TimelineDots — baseline <line> + N <circle>s +
    alternating <text> callouts above/below.
    """
    theme = ctx.theme
    palette = ctx.palette

    outline_only = theme.accent_treatment == "outline" or theme.is_dark

    milestones = list(_MILESTONES)
    n = len(milestones)

    # Layout — spine at vertical centre.
    spine_y = height_px // 2
    dot_d = int(min(height_px * 0.18, width_px * 0.06))
    dot_d = max(dot_d, 24)
    dot_r = dot_d / 2.0

    margin_x = int(width_px * 0.05)
    usable_w = width_px - 2 * margin_x
    stride = (usable_w / (n - 1)) if n > 1 else 0
    first_cx = margin_x
    dot_centres = [first_cx + i * stride for i in range(n)]

    spine_left = int(margin_x * 0.5)
    spine_right = width_px - int(margin_x * 0.5)

    if outline_only:
        spine_color = _hex(palette.primary)
        dot_fill = _hex(palette.background) if not theme.is_dark else _hex(palette.primary)
        dot_text_color = _hex(palette.primary) if not theme.is_dark else _hex(palette.background)
        dot_stroke = _hex(palette.primary)
    else:
        spine_color = _hex(palette.muted)
        dot_fill = _hex(palette.primary)
        dot_text_color = _hex(palette.background)
        dot_stroke = _hex(palette.primary)

    heading_font = theme.typography.heading
    body_font = theme.typography.body
    body_pt = theme.typography.body_size_pt
    cap_pt = theme.typography.caption_size_pt

    parts: list[str] = []
    parts.append(
        f'<svg xmlns="http://www.w3.org/2000/svg" '
        f'width="{width_px}" height="{height_px}" '
        f'viewBox="0 0 {width_px} {height_px}">'
    )
    # Background
    parts.append(
        f'  <rect x="0" y="0" width="{width_px}" height="{height_px}" '
        f'fill="{_hex(palette.background)}" />'
    )

    # Spine baseline.
    parts.append(
        f'  <line x1="{spine_left}" y1="{spine_y}" '
        f'x2="{spine_right}" y2="{spine_y}" '
        f'stroke="{spine_color}" stroke-width="2" />'
    )

    # Dots + numbered labels + alternating callouts.
    for i, (date, label) in enumerate(milestones):
        cx = dot_centres[i]
        # Circle marker.
        parts.append(
            f'  <circle cx="{cx:.2f}" cy="{spine_y}" r="{dot_r:.2f}" '
            f'fill="{dot_fill}" stroke="{dot_stroke}" stroke-width="1" />'
        )
        # Number inside the dot.
        parts.append(
            f'  <text x="{cx:.2f}" y="{spine_y + cap_pt / 3:.2f}" '
            f'font-family="{heading_font}" font-size="{cap_pt}" '
            f'font-weight="bold" fill="{dot_text_color}" '
            f'text-anchor="middle">{i + 1}</text>'
        )

        # Callout placement — alternate above / below.
        if i % 2 == 0:
            # Above the spine.
            date_y = spine_y - dot_r - 12 - body_pt
            label_y = date_y + body_pt + 4
        else:
            # Below the spine.
            date_y = spine_y + dot_r + 12 + body_pt
            label_y = date_y + body_pt + 4

        parts.append(
            f'  <text x="{cx:.2f}" y="{date_y:.2f}" '
            f'font-family="{heading_font}" font-size="{body_pt}" '
            f'font-weight="bold" fill="{_hex(palette.text)}" '
            f'text-anchor="middle">{date}</text>'
        )
        parts.append(
            f'  <text x="{cx:.2f}" y="{label_y:.2f}" '
            f'font-family="{body_font}" font-size="{cap_pt}" '
            f'fill="{_hex(palette.muted)}" '
            f'text-anchor="middle">{label}</text>'
        )

    parts.append("</svg>")
    return "\n".join(parts)


TimelineDots.render_svg = _render_svg_timeline_dots  # type: ignore[attr-defined]
