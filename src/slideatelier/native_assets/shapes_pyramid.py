"""Pyramid — 4-level hierarchy (Foundation → Vision).

Implemented as a stack of native MSO_SHAPE.TRAPEZOID primitives, progressively
narrower from bottom to top. python-pptx's TRAPEZOID renders narrower at the
top by default (which is what we want), so each tier is just a trapezoid with
shrinking width.
"""
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

from .base import AssetShape, ShapeRenderContext
from .shapes_matrix import _lighten


# Bottom → top
_TIERS = [
    "Foundation",
    "Structure",
    "Strategy",
    "Vision",
]
_WIDTH_PCTS = [1.00, 0.78, 0.55, 0.32]


class Pyramid(AssetShape):
    id = "pyramid"
    name = "Pyramid"
    description = "4-level hierarchy pyramid (Foundation → Vision)."
    style_tags = ("hierarchy", "structured", "consulting", "framework")
    aspect_ratio_hint = 1.4

    def render(self, slide, ctx: ShapeRenderContext) -> None:
        theme = ctx.theme
        palette = ctx.palette

        n = len(_TIERS)
        # 30% on the right reserved for tier captions.
        caption_w = int(ctx.width * 0.28)
        pyramid_area_w = ctx.width - caption_w
        gap_emu = int(ctx.height * 0.015)
        total_gap = gap_emu * (n - 1)
        tier_h = (ctx.height - total_gap) // n

        outline_only = theme.accent_treatment == "outline" or theme.is_dark

        for i in range(n):
            # i=0 is bottom, i=3 is top (Vision).
            tier_idx_top_first = n - 1 - i  # 3,2,1,0 for top→bottom drawing
            label = _TIERS[i]
            w_pct = _WIDTH_PCTS[i]

            top = ctx.top + tier_idx_top_first * (tier_h + gap_emu)
            w = int(pyramid_area_w * w_pct)
            left = ctx.left + (pyramid_area_w - w) // 2

            # Top tier uses TRAPEZOID for the pyramidal apex; lower tiers
            # use trapezoids too so the silhouette stays pyramidal.
            shape_kind = MSO_SHAPE.TRAPEZOID

            tier = slide.shapes.add_shape(shape_kind, left, top, w, tier_h)
            tier.fill.solid()

            if outline_only:
                tier.fill.fore_color.rgb = palette.background
                tier.line.color.rgb = palette.primary
                label_color = palette.text
            else:
                # Bottom = lightest, top = solid primary. Index i goes 0..3 where
                # i=3 is top → most saturated.
                tint = 0.78 - i * 0.20  # 0.78, 0.58, 0.38, 0.18
                tint = max(0.0, tint)
                tier.fill.fore_color.rgb = _lighten(palette.primary, tint)
                tier.line.color.rgb = palette.primary
                # White-ish text for top 2 tiers (dark fill), dark text for bottom.
                if i >= 2:
                    label_color = palette.background
                else:
                    label_color = palette.text

            tier.line.width = theme.line_weight_emu

            tf = tier.text_frame
            tf.margin_left = Pt(4)
            tf.margin_right = Pt(4)
            tf.margin_top = Pt(4)
            tf.margin_bottom = Pt(4)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = 2  # CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(theme.typography.body_size_pt)
            run.font.name = theme.typography.heading
            run.font.bold = True
            run.font.color.rgb = label_color

            # Caption to the right of each tier.
            cap = slide.shapes.add_textbox(
                ctx.left + pyramid_area_w + int(ctx.width * 0.01),
                top,
                caption_w - int(ctx.width * 0.01),
                tier_h,
            )
            ctf = cap.text_frame
            ctf.word_wrap = True
            ctf.margin_left = Pt(4)
            ctf.margin_right = Pt(4)
            cp = ctf.paragraphs[0]
            cp.alignment = 1  # LEFT
            crun = cp.add_run()
            crun.text = f"Tier {n - i}"
            crun.font.size = Pt(theme.typography.caption_size_pt)
            crun.font.name = theme.typography.body
            crun.font.color.rgb = palette.muted

    def render_svg(self, ctx: ShapeRenderContext, width_px: int, height_px: int) -> str:
        """Inline SVG mirroring render(): 4 stacked trapezoid <polygon>s
        forming a pyramid (Foundation at bottom → Vision at top), with
        captions to the right of each tier.
        """
        theme = ctx.theme
        palette = ctx.palette

        n = len(_TIERS)
        caption_w = int(width_px * 0.28)
        pyramid_area_w = width_px - caption_w
        gap_px = int(height_px * 0.015)
        total_gap = gap_px * (n - 1)
        tier_h = (height_px - total_gap) // n

        outline_only = theme.accent_treatment == "outline" or theme.is_dark

        body_pt = theme.typography.body_size_pt
        cap_pt = theme.typography.caption_size_pt
        heading_font = theme.typography.heading
        body_font = theme.typography.body

        parts: list[str] = []
        parts.append(
            f'<svg xmlns="http://www.w3.org/2000/svg" '
            f'width="{width_px}" height="{height_px}" '
            f'viewBox="0 0 {width_px} {height_px}">'
        )
        parts.append(
            f'  <rect x="0" y="0" width="{width_px}" height="{height_px}" '
            f'fill="{_hex_color(palette.background)}" />'
        )

        for i in range(n):
            # i=0 is bottom (Foundation), i=3 is top (Vision)
            tier_idx_top_first = n - 1 - i
            label = _TIERS[i]
            w_pct = _WIDTH_PCTS[i]
            # Bottom width = own row pct, top width = next-up row pct (or 70%
            # for the apex). This produces true narrowing trapezoids that
            # stack into a pyramid silhouette — matching MSO_SHAPE.TRAPEZOID's
            # narrower-at-top default.
            bottom_w = int(pyramid_area_w * w_pct)
            if i < n - 1:
                top_w = int(pyramid_area_w * _WIDTH_PCTS[i + 1])
            else:
                top_w = max(int(bottom_w * 0.5), 1)

            top_y = tier_idx_top_first * (tier_h + gap_px)
            bottom_y = top_y + tier_h

            top_left_x = (pyramid_area_w - top_w) // 2
            top_right_x = top_left_x + top_w
            bottom_left_x = (pyramid_area_w - bottom_w) // 2
            bottom_right_x = bottom_left_x + bottom_w

            points = (
                f"{top_left_x},{top_y} "
                f"{top_right_x},{top_y} "
                f"{bottom_right_x},{bottom_y} "
                f"{bottom_left_x},{bottom_y}"
            )

            if outline_only:
                fill = palette.background
                stroke = palette.primary
                label_color = palette.text
            else:
                tint = max(0.0, 0.78 - i * 0.20)
                fill = _lighten(palette.primary, tint)
                stroke = palette.primary
                if i >= 2:
                    label_color = palette.background
                else:
                    label_color = palette.text

            parts.append(
                f'  <polygon points="{points}" '
                f'fill="{_hex_color(fill)}" stroke="{_hex_color(stroke)}" '
                f'stroke-width="1" />'
            )
            tx = pyramid_area_w // 2
            ty = top_y + tier_h // 2 + body_pt // 3
            parts.append(
                f'  <text x="{tx}" y="{ty}" '
                f'font-family="{heading_font}" font-size="{body_pt}" '
                f'font-weight="bold" fill="{_hex_color(label_color)}" '
                f'text-anchor="middle">{label}</text>'
            )

            cap_x = pyramid_area_w + int(width_px * 0.01)
            cap_y = top_y + tier_h // 2 + cap_pt // 3
            parts.append(
                f'  <text x="{cap_x}" y="{cap_y}" '
                f'font-family="{body_font}" font-size="{cap_pt}" '
                f'fill="{_hex_color(palette.muted)}">Tier {n - i}</text>'
            )

        parts.append("</svg>")
        return "\n".join(parts)


def _hex_color(color) -> str:
    """RGBColor → '#RRGGBB' for SVG."""
    return "#{:02X}{:02X}{:02X}".format(color[0], color[1], color[2])
