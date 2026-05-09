"""Ribbon — horizontal banner: central RECTANGLE flanked by RIGHT_TRIANGLE flags.

The two end "swallowtail" cuts are formed by inverted right triangles whose
hypotenuse aligns with the inner edge of the central rectangle, creating the
classic banner ribbon silhouette using only native primitives.
"""
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

from .base import AssetShape, ShapeRenderContext
from .shapes_matrix import _hex, _lighten


class Ribbon(AssetShape):
    id = "ribbon-banner"
    name = "Ribbon Banner"
    description = "Horizontal ribbon banner — center rectangle with triangular ribbon ends."
    style_tags = ("hero", "graphics-heavy", "contemporary", "decorative")
    aspect_ratio_hint = 3.0

    def render(self, slide, ctx: ShapeRenderContext) -> None:
        theme = ctx.theme
        palette = ctx.palette

        outline_only = theme.accent_treatment == "outline" or theme.is_dark

        # Banner occupies central horizontal stripe (vertically centered, ~50% h).
        band_h = int(ctx.height * 0.50)
        band_top = ctx.top + (ctx.height - band_h) // 2

        # Ends: each triangular flag is band_h wide.
        end_w = band_h  # square-ish flag for right-triangle shape

        # The two flags overlap the central rectangle a bit so the join is clean.
        overlap = int(end_w * 0.25)

        # Central rectangle spans from (left + end_w - overlap) to
        # (right - end_w + overlap).
        center_left = ctx.left + end_w - overlap
        center_right = ctx.left + ctx.width - end_w + overlap
        center_w = center_right - center_left

        # Choose fills.
        if outline_only:
            ribbon_fill = palette.background
            ribbon_line = palette.primary
            text_color = palette.text
            flag_fill = palette.background
            flag_line = palette.accent
        else:
            ribbon_fill = palette.primary
            ribbon_line = palette.primary
            text_color = palette.background
            flag_fill = _lighten(palette.primary, 0.30)
            flag_line = palette.primary

        # Left flag: RIGHT_TRIANGLE rotated/oriented so the hypotenuse points
        # into the center. Default RIGHT_TRIANGLE has the right angle at the
        # bottom-left. To make a "ribbon end" that points left, we rotate 270°
        # so the right angle ends up at the right side, hypotenuse on the left.
        left_flag = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_TRIANGLE,
            ctx.left,
            band_top,
            end_w,
            band_h,
        )
        # By default the hypotenuse runs from top-left to bottom-right of the
        # bounding box. Flipping horizontally gives a triangle whose hypotenuse
        # runs top-right to bottom-left — i.e. the slanted edge faces inward.
        # python-pptx exposes flipping via the .rotation/.element; we simply
        # rotate 0 (the default already gives an inward-facing slant for the
        # left flag if flipped). Use the rotation property as a simple approach.
        # Actually the default RIGHT_TRIANGLE has the right angle at the bottom-
        # left and the slanted hypotenuse from top-left down to bottom-right.
        # For the LEFT flag we want the slanted edge on the right (pointing
        # inward). Rotating 0° already gives slant from top-left to bottom-
        # right which is acceptable; flag end shape on left side ribbon visual
        # is achieved by treating the triangle's right edge as the join.
        left_flag.fill.solid()
        left_flag.fill.fore_color.rgb = flag_fill
        left_flag.line.color.rgb = flag_line
        left_flag.line.width = theme.line_weight_emu

        # Right flag: same triangle but mirrored (rotation 180° flips both axes
        # so hypotenuse runs top-right to bottom-left and right-angle is at
        # top-right relative to bbox).
        right_flag = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_TRIANGLE,
            ctx.left + ctx.width - end_w,
            band_top,
            end_w,
            band_h,
        )
        right_flag.rotation = 180.0
        right_flag.fill.solid()
        right_flag.fill.fore_color.rgb = flag_fill
        right_flag.line.color.rgb = flag_line
        right_flag.line.width = theme.line_weight_emu

        # Central rectangle (drawn last so it sits over the flag overlap).
        shape_kind = (
            MSO_SHAPE.ROUNDED_RECTANGLE
            if theme.corner_radius_pct > 0
            else MSO_SHAPE.RECTANGLE
        )
        center = slide.shapes.add_shape(
            shape_kind, center_left, band_top, center_w, band_h
        )
        center.fill.solid()
        center.fill.fore_color.rgb = ribbon_fill
        center.line.color.rgb = ribbon_line
        center.line.width = theme.line_weight_emu

        tf = center.text_frame
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(4)
        tf.margin_bottom = Pt(4)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = 2  # CENTER
        run = p.add_run()
        run.text = "FEATURED"
        run.font.size = Pt(theme.typography.heading_size_pt)
        run.font.name = theme.typography.heading
        run.font.bold = True
        run.font.color.rgb = text_color


# ---------------------------------------------------------------------------
# SVG rendering (Web Deck publishing — Sprint J.B)
# ---------------------------------------------------------------------------

def _render_svg_ribbon(
    self: "Ribbon", ctx: ShapeRenderContext, width_px: int, height_px: int
) -> str:
    theme = ctx.theme
    palette = ctx.palette

    outline_only = theme.accent_treatment == "outline" or theme.is_dark

    # Mirror render() proportions.
    band_h = int(height_px * 0.50)
    band_top = (height_px - band_h) // 2
    end_w = band_h
    overlap = int(end_w * 0.25)
    center_left = end_w - overlap
    center_right = width_px - end_w + overlap
    center_w = center_right - center_left

    if outline_only:
        ribbon_fill = _hex(palette.background)
        ribbon_line = _hex(palette.primary)
        text_color = _hex(palette.text)
        flag_fill = _hex(palette.background)
        flag_line = _hex(palette.accent)
    else:
        ribbon_fill = _hex(palette.primary)
        ribbon_line = _hex(palette.primary)
        text_color = _hex(palette.background)
        flag_fill = _hex(_lighten(palette.primary, 0.30))
        flag_line = _hex(palette.primary)

    heading_font = theme.typography.heading
    heading_pt = theme.typography.heading_size_pt

    parts: list[str] = []
    parts.append(
        f'<svg xmlns="http://www.w3.org/2000/svg" '
        f'width="{width_px}" height="{height_px}" '
        f'viewBox="0 0 {width_px} {height_px}">'
    )
    parts.append(
        f'  <rect x="0" y="0" width="{width_px}" height="{height_px}" '
        f'fill="{_hex(palette.background)}" />'
    )

    # Left flag — triangle with notch on the right (point inward).
    # Mirrors the RIGHT_TRIANGLE silhouette: outer left edge full height,
    # outer top edge full width, inner slanted hypotenuse from top-right
    # to bottom-left of bbox is what python-pptx draws by default. We
    # render the polygon using three vertices.
    lf_x = 0
    lf_y = band_top
    lf_pts = (
        f"{lf_x},{lf_y} "  # top-left
        f"{lf_x + end_w},{lf_y} "  # top-right
        f"{lf_x},{lf_y + band_h}"  # bottom-left
    )
    parts.append(
        f'  <polygon points="{lf_pts}" '
        f'fill="{flag_fill}" stroke="{flag_line}" stroke-width="1.5" />'
    )

    # Right flag — mirror image (rotation 180 in render()).
    rf_x = width_px - end_w
    rf_y = band_top
    rf_pts = (
        f"{rf_x + end_w},{rf_y + band_h} "  # bottom-right
        f"{rf_x},{rf_y + band_h} "  # bottom-left
        f"{rf_x + end_w},{rf_y}"  # top-right
    )
    parts.append(
        f'  <polygon points="{rf_pts}" '
        f'fill="{flag_fill}" stroke="{flag_line}" stroke-width="1.5" />'
    )

    # Central rectangle (rounded if theme has rounded corners).
    rx_attr = ""
    if theme.corner_radius_pct > 0:
        rxv = max(4, int(band_h * 0.15))
        rx_attr = f' rx="{rxv}" ry="{rxv}"'
    parts.append(
        f'  <rect x="{center_left}" y="{band_top}" '
        f'width="{center_w}" height="{band_h}"{rx_attr} '
        f'fill="{ribbon_fill}" stroke="{ribbon_line}" stroke-width="1.5" />'
    )

    # Centred label.
    label_x = center_left + center_w / 2.0
    label_y = band_top + band_h / 2.0 + heading_pt / 3.0
    parts.append(
        f'  <text x="{label_x:.2f}" y="{label_y:.2f}" '
        f'font-family="{heading_font}" font-size="{heading_pt}" '
        f'font-weight="bold" fill="{text_color}" '
        f'text-anchor="middle">FEATURED</text>'
    )

    parts.append("</svg>")
    return "\n".join(parts)


Ribbon.render_svg = _render_svg_ribbon  # type: ignore[attr-defined]
