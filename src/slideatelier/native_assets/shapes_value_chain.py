"""Value Chain — Porter's value chain framework as native shapes.

Five primary-activity chevrons flow left → right (top 2/3 of the bounding rect),
with four wider support-activity rectangles below (bottom 1/3). A small
right-arrow marked "MARGIN" caps the far-right edge.

All shapes are native python-pptx primitives (CHEVRON, RECTANGLE,
ROUNDED_RECTANGLE, RIGHT_ARROW) so the deck remains fully editable.
"""
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

from .base import AssetShape, ShapeRenderContext
from .shapes_matrix import _lighten


PRIMARY_LABELS = (
    "Inbound logistics",
    "Operations",
    "Outbound logistics",
    "Marketing & sales",
    "Service",
)

SUPPORT_LABELS = (
    "Firm infrastructure",
    "HR management",
    "Tech development",
    "Procurement",
)


class ValueChain(AssetShape):
    id = "value-chain"
    name = "Value Chain"
    description = "Porter's value chain — 5 primary activities + 4 support activities."
    style_tags = ("process", "consulting", "porter", "framework")
    aspect_ratio_hint = 2.0  # wider than tall

    def render(self, slide, ctx: ShapeRenderContext) -> None:
        theme = ctx.theme
        palette = ctx.palette

        # Vertical split: top 55% primary chevrons, gap 5%, bottom 30% support
        # Reserve ~10% on the right for the MARGIN arrow.
        margin_w = int(ctx.width * 0.08)
        chain_left = ctx.left
        chain_w = ctx.width - margin_w - int(ctx.width * 0.01)  # tiny gap before MARGIN

        primary_h = int(ctx.height * 0.55)
        gap_h = int(ctx.height * 0.05)
        support_h = ctx.height - primary_h - gap_h

        # ----- Primary chevrons (top row) -----
        n_primary = 5
        # 5 chevrons each 18% of chain width with 2% overlap (so total ~ 5*18 - 4*2 = 82%).
        # Compute by allocating equal stride with overlap.
        # Effective stride = (chain_w - chevron_w) / (n - 1) for full coverage.
        chevron_w = int(chain_w * 0.22)
        if n_primary > 1:
            stride = (chain_w - chevron_w) // (n_primary - 1)
        else:
            stride = chain_w

        # Treatment-driven coloring
        outline_only = theme.accent_treatment == "outline"
        is_dark = theme.is_dark

        for i, label in enumerate(PRIMARY_LABELS):
            cleft = chain_left + i * stride
            ctop = ctx.top
            chev = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                cleft,
                ctop,
                chevron_w,
                primary_h,
            )
            chev.fill.solid()

            # Last chevron uses accent so the chain "graduates" toward margin.
            base_color = palette.accent if i == n_primary - 1 else palette.primary

            if outline_only:
                chev.fill.fore_color.rgb = palette.background
                chev.line.color.rgb = base_color
                label_color = base_color
            elif is_dark:
                chev.fill.fore_color.rgb = base_color
                chev.line.color.rgb = base_color
                label_color = palette.background if _is_lightish(palette.background) else palette.text
                # palette.background on dark theme is dark; we want light text on dark fill
                label_color = palette.text
            else:
                chev.fill.fore_color.rgb = base_color
                chev.line.color.rgb = base_color
                label_color = palette.background  # white-ish text on filled chevron

            chev.line.width = theme.line_weight_emu

            tf = chev.text_frame
            tf.margin_left = Pt(4)
            tf.margin_right = Pt(4)
            tf.margin_top = Pt(4)
            tf.margin_bottom = Pt(4)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = 2  # CENTER (PP_ALIGN.CENTER == 2 in python-pptx? actually CENTER=2)
            # Note: PP_ALIGN.CENTER is 2 in python-pptx enum
            run = p.add_run()
            run.text = label
            run.font.size = Pt(theme.typography.body_size_pt)
            run.font.name = theme.typography.heading
            run.font.bold = True
            run.font.color.rgb = label_color

        # ----- Support rectangles (bottom row) -----
        n_support = 4
        support_top = ctx.top + primary_h + gap_h
        support_total_w = chain_w
        support_gap = int(support_total_w * 0.015)
        support_w = (support_total_w - support_gap * (n_support - 1)) // n_support

        # Support fill/treatment
        if outline_only:
            support_fill = palette.background
            support_line = palette.muted
            support_text_color = palette.text
        elif is_dark:
            support_fill = _darken(palette.background, 0.15)
            support_line = palette.muted
            support_text_color = palette.text
        else:
            support_fill = _lighten(palette.muted, 0.85)
            support_line = palette.muted
            support_text_color = palette.text

        shape_kind = (
            MSO_SHAPE.ROUNDED_RECTANGLE
            if theme.corner_radius_pct > 0
            else MSO_SHAPE.RECTANGLE
        )

        for i, label in enumerate(SUPPORT_LABELS):
            sleft = chain_left + i * (support_w + support_gap)
            rect = slide.shapes.add_shape(
                shape_kind,
                sleft,
                support_top,
                support_w,
                support_h,
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = support_fill
            rect.line.color.rgb = support_line
            rect.line.width = theme.line_weight_emu

            tf = rect.text_frame
            tf.margin_left = Pt(6)
            tf.margin_right = Pt(6)
            tf.margin_top = Pt(4)
            tf.margin_bottom = Pt(4)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = 2  # CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(theme.typography.caption_size_pt)
            run.font.name = theme.typography.body
            run.font.bold = False
            run.font.color.rgb = support_text_color

        # ----- MARGIN right-arrow accent -----
        margin_left = chain_left + chain_w + int(ctx.width * 0.005)
        margin_top = ctx.top + int(primary_h * 0.15)
        margin_height = int(primary_h * 0.7)
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            margin_left,
            margin_top,
            margin_w,
            margin_height,
        )
        arrow.fill.solid()
        if outline_only:
            arrow.fill.fore_color.rgb = palette.background
            arrow.line.color.rgb = palette.accent
            margin_text_color = palette.accent
        else:
            arrow.fill.fore_color.rgb = palette.accent
            arrow.line.color.rgb = palette.accent
            margin_text_color = palette.background
        arrow.line.width = theme.line_weight_emu

        atf = arrow.text_frame
        atf.margin_left = Pt(2)
        atf.margin_right = Pt(2)
        atf.margin_top = Pt(2)
        atf.margin_bottom = Pt(2)
        ap = atf.paragraphs[0]
        ap.alignment = 2  # CENTER
        arun = ap.add_run()
        arun.text = "MARGIN"
        arun.font.size = Pt(theme.typography.caption_size_pt)
        arun.font.name = theme.typography.heading
        arun.font.bold = True
        arun.font.color.rgb = margin_text_color


    def render_svg(self, ctx: ShapeRenderContext, width_px: int, height_px: int) -> str:
        """Inline SVG mirroring render(): 5 chevron <polygon>s on top, 4 support
        <rect>s on the bottom, and 1 right-arrow <polygon> for MARGIN.
        """
        theme = ctx.theme
        palette = ctx.palette

        margin_w = int(width_px * 0.08)
        chain_left = 0
        chain_w = width_px - margin_w - int(width_px * 0.01)

        primary_h = int(height_px * 0.55)
        gap_h = int(height_px * 0.05)
        support_h = height_px - primary_h - gap_h

        n_primary = 5
        chevron_w = int(chain_w * 0.22)
        if n_primary > 1:
            stride = (chain_w - chevron_w) // (n_primary - 1)
        else:
            stride = chain_w

        outline_only = theme.accent_treatment == "outline"
        is_dark = theme.is_dark

        body_pt = theme.typography.body_size_pt
        cap_pt = theme.typography.caption_size_pt
        heading_font = theme.typography.heading
        body_font = theme.typography.body

        parts: list[str] = []
        parts.append(
            f'<svg xmlns="http://www.w3.org/2000/svg" '
            f'width="{width_px}" height="{height_px}" '
            f'viewBox="0 0 {width_px} {height_px}">'
        )
        parts.append(
            f'  <rect x="0" y="0" width="{width_px}" height="{height_px}" '
            f'fill="{_hex_color(palette.background)}" />'
        )

        # Chevron polygon: 5 points — left flat side, indent for arrow notch on
        # the left, point on the right. Mirrors MSO_SHAPE.CHEVRON which is a
        # right-pointing arrow with a left V-notch (so chevrons can tile).
        # Points (clockwise from top-left):
        #   (x0, y0)            → top-left
        #   (x0 + (w - notch), y0)        → top-right of body
        #   (x0 + w, y0 + h/2)            → tip
        #   (x0 + (w - notch), y0 + h)    → bottom-right of body
        #   (x0, y0 + h)                  → bottom-left
        #   (x0 + notch, y0 + h/2)        → notch in
        def _chevron_points(x0: int, y0: int, w: int, h: int) -> str:
            notch = h // 2
            verts = [
                (x0, y0),
                (x0 + w - notch, y0),
                (x0 + w, y0 + h // 2),
                (x0 + w - notch, y0 + h),
                (x0, y0 + h),
                (x0 + notch, y0 + h // 2),
            ]
            return " ".join(f"{vx},{vy}" for vx, vy in verts)

        for i, label in enumerate(PRIMARY_LABELS):
            cleft = chain_left + i * stride
            ctop = 0
            base_color = palette.accent if i == n_primary - 1 else palette.primary

            if outline_only:
                fill = palette.background
                stroke = base_color
                label_color = base_color
            elif is_dark:
                fill = base_color
                stroke = base_color
                label_color = palette.text
            else:
                fill = base_color
                stroke = base_color
                label_color = palette.background

            parts.append(
                f'  <polygon points="{_chevron_points(cleft, ctop, chevron_w, primary_h)}" '
                f'fill="{_hex_color(fill)}" stroke="{_hex_color(stroke)}" '
                f'stroke-width="1" />'
            )
            tx = cleft + chevron_w // 2
            ty = ctop + primary_h // 2 + body_pt // 3
            parts.append(
                f'  <text x="{tx}" y="{ty}" '
                f'font-family="{heading_font}" font-size="{body_pt}" '
                f'font-weight="bold" fill="{_hex_color(label_color)}" '
                f'text-anchor="middle">{_xml_escape(label)}</text>'
            )

        # Support rectangles
        n_support = 4
        support_top = primary_h + gap_h
        support_total_w = chain_w
        support_gap = int(support_total_w * 0.015)
        support_w = (support_total_w - support_gap * (n_support - 1)) // n_support

        if outline_only:
            support_fill = palette.background
            support_line = palette.muted
            support_text_color = palette.text
        elif is_dark:
            support_fill = _darken(palette.background, 0.15)
            support_line = palette.muted
            support_text_color = palette.text
        else:
            support_fill = _lighten(palette.muted, 0.85)
            support_line = palette.muted
            support_text_color = palette.text

        for i, label in enumerate(SUPPORT_LABELS):
            sleft = chain_left + i * (support_w + support_gap)
            parts.append(
                f'  <rect x="{sleft}" y="{support_top}" '
                f'width="{support_w}" height="{support_h}" '
                f'fill="{_hex_color(support_fill)}" '
                f'stroke="{_hex_color(support_line)}" stroke-width="1" />'
            )
            tx = sleft + support_w // 2
            ty = support_top + support_h // 2 + cap_pt // 3
            parts.append(
                f'  <text x="{tx}" y="{ty}" '
                f'font-family="{body_font}" font-size="{cap_pt}" '
                f'fill="{_hex_color(support_text_color)}" '
                f'text-anchor="middle">{_xml_escape(label)}</text>'
            )

        # MARGIN right-arrow polygon (RIGHT_ARROW = 7-point arrow shape).
        margin_left = chain_left + chain_w + int(width_px * 0.005)
        margin_top = int(primary_h * 0.15)
        margin_height = int(primary_h * 0.7)

        if outline_only:
            arrow_fill = palette.background
            arrow_stroke = palette.accent
            margin_text_color = palette.accent
        else:
            arrow_fill = palette.accent
            arrow_stroke = palette.accent
            margin_text_color = palette.background

        # Right-arrow polygon: rectangle body with triangle head on the right.
        # Body height = 60% of overall, head occupies right ~40% width.
        ax = margin_left
        ay = margin_top
        aw = margin_w
        ah = margin_height
        body_top = ay + ah // 4
        body_bot = ay + 3 * ah // 4
        head_start_x = ax + aw // 2
        arrow_pts = (
            f"{ax},{body_top} "
            f"{head_start_x},{body_top} "
            f"{head_start_x},{ay} "
            f"{ax + aw},{ay + ah // 2} "
            f"{head_start_x},{ay + ah} "
            f"{head_start_x},{body_bot} "
            f"{ax},{body_bot}"
        )
        parts.append(
            f'  <polygon points="{arrow_pts}" '
            f'fill="{_hex_color(arrow_fill)}" '
            f'stroke="{_hex_color(arrow_stroke)}" stroke-width="1" />'
        )
        tx = ax + aw // 4
        ty = ay + ah // 2 + cap_pt // 3
        parts.append(
            f'  <text x="{tx}" y="{ty}" '
            f'font-family="{heading_font}" font-size="{cap_pt}" '
            f'font-weight="bold" fill="{_hex_color(margin_text_color)}" '
            f'text-anchor="middle">MARGIN</text>'
        )

        parts.append("</svg>")
        return "\n".join(parts)


def _hex_color(color) -> str:
    """RGBColor → '#RRGGBB' for SVG."""
    return "#{:02X}{:02X}{:02X}".format(color[0], color[1], color[2])


def _xml_escape(s: str) -> str:
    """Escape XML special chars so labels like 'Marketing & sales' don't break
    SVG parsing. Apply to every text label embedded in the rendered SVG."""
    return (
        str(s)
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
        .replace("'", "&apos;")
    )


def _is_lightish(color) -> bool:
    """Rough luminance check — True if color is closer to white than black."""
    return (color[0] + color[1] + color[2]) / 3 > 128


def _darken(color, factor: float):
    """Blend an RGBColor toward black. factor=0 → original, factor=1 → black."""
    from pptx.dml.color import RGBColor
    r = max(0, int(color[0] * (1 - factor)))
    g = max(0, int(color[1] * (1 - factor)))
    b = max(0, int(color[2] * (1 - factor)))
    return RGBColor(r, g, b)
