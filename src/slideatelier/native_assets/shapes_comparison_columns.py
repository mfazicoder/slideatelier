"""ComparisonColumns — 3 vertical option columns side by side.

Each column = header rectangle (colored) on top + body rectangle below for
bullet text. Three differentiated colors (primary, accent, muted) make the
options scan at a glance.
"""
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

from .base import AssetShape, ShapeRenderContext
from .shapes_matrix import _lighten


def _hex(color) -> str:
    """RGBColor → '#RRGGBB' for SVG."""
    return "#{:02X}{:02X}{:02X}".format(color[0], color[1], color[2])


_OPTIONS = ["Option A", "Option B", "Option C"]
_BULLETS = (
    "First differentiator\nSecond benefit\nKey constraint",
    "First differentiator\nSecond benefit\nKey constraint",
    "First differentiator\nSecond benefit\nKey constraint",
)


class ComparisonColumns(AssetShape):
    id = "comparison-columns"
    name = "Comparison Columns"
    description = "3 side-by-side option columns with header + body for fast comparison."
    style_tags = ("comparison", "structured", "framework", "consulting")
    aspect_ratio_hint = 1.5

    def render(self, slide, ctx: ShapeRenderContext) -> None:
        theme = ctx.theme
        palette = ctx.palette

        outline_only = theme.accent_treatment == "outline" or theme.is_dark

        n = 3
        gap = int(ctx.width * 0.02)
        col_w = (ctx.width - gap * (n - 1)) // n
        header_h = int(ctx.height * 0.18)
        body_h = ctx.height - header_h

        shape_kind = (
            MSO_SHAPE.ROUNDED_RECTANGLE
            if theme.corner_radius_pct > 0
            else MSO_SHAPE.RECTANGLE
        )

        # Distinct color per column
        header_colors = [palette.primary, palette.accent, palette.muted]

        for i, label in enumerate(_OPTIONS):
            cleft = ctx.left + i * (col_w + gap)
            ctop = ctx.top

            base_color = header_colors[i]

            # Header
            hdr = slide.shapes.add_shape(
                shape_kind, cleft, ctop, col_w, header_h
            )
            hdr.fill.solid()
            if outline_only:
                hdr.fill.fore_color.rgb = palette.background
                hdr.line.color.rgb = base_color
                hdr_text_color = base_color
            else:
                hdr.fill.fore_color.rgb = base_color
                hdr.line.color.rgb = base_color
                hdr_text_color = palette.background
            hdr.line.width = theme.line_weight_emu

            htf = hdr.text_frame
            htf.margin_left = Pt(6)
            htf.margin_right = Pt(6)
            htf.margin_top = Pt(4)
            htf.margin_bottom = Pt(4)
            htf.word_wrap = True
            hp = htf.paragraphs[0]
            hp.alignment = 2  # CENTER
            hrun = hp.add_run()
            hrun.text = label
            hrun.font.size = Pt(theme.typography.heading_size_pt)
            hrun.font.name = theme.typography.heading
            hrun.font.bold = True
            hrun.font.color.rgb = hdr_text_color

            # Body
            body = slide.shapes.add_shape(
                shape_kind, cleft, ctop + header_h, col_w, body_h
            )
            body.fill.solid()
            if outline_only:
                body.fill.fore_color.rgb = palette.background
                body.line.color.rgb = palette.muted
                body_text_color = palette.text
            else:
                body.fill.fore_color.rgb = _lighten(base_color, 0.88)
                body.line.color.rgb = palette.muted
                body_text_color = palette.text
            body.line.width = theme.line_weight_emu

            btf = body.text_frame
            btf.margin_left = Pt(8)
            btf.margin_right = Pt(8)
            btf.margin_top = Pt(8)
            btf.margin_bottom = Pt(8)
            btf.word_wrap = True

            bullets = _BULLETS[i].split("\n")
            for j, bullet in enumerate(bullets):
                if j == 0:
                    bp = btf.paragraphs[0]
                else:
                    bp = btf.add_paragraph()
                bp.alignment = 1  # LEFT
                brun = bp.add_run()
                brun.text = f"• {bullet}"
                brun.font.size = Pt(theme.typography.body_size_pt)
                brun.font.name = theme.typography.body
                brun.font.color.rgb = body_text_color


# ---------------------------------------------------------------------------
# SVG rendering (Web Deck publishing — Sprint J.C)
# ---------------------------------------------------------------------------

def _render_svg_comparison_columns(
    self: "ComparisonColumns", ctx: ShapeRenderContext, width_px: int, height_px: int
) -> str:
    """Inline SVG mirror of ComparisonColumns — header + body <rect> per column,
    plus thin <line> dividers between columns.
    """
    theme = ctx.theme
    palette = ctx.palette

    outline_only = theme.accent_treatment == "outline" or theme.is_dark

    options = list(_OPTIONS)
    bullets_per_col = [b.split("\n") for b in _BULLETS]
    n = len(options)

    gap = int(width_px * 0.02)
    col_w = (width_px - gap * (n - 1)) // n
    header_h = int(height_px * 0.18)
    body_h = height_px - header_h

    header_colors = [palette.primary, palette.accent, palette.muted]

    heading_font = theme.typography.heading
    body_font = theme.typography.body
    heading_pt = theme.typography.heading_size_pt
    body_pt = theme.typography.body_size_pt

    parts: list[str] = []
    parts.append(
        f'<svg xmlns="http://www.w3.org/2000/svg" '
        f'width="{width_px}" height="{height_px}" '
        f'viewBox="0 0 {width_px} {height_px}">'
    )
    # Background
    parts.append(
        f'  <rect x="0" y="0" width="{width_px}" height="{height_px}" '
        f'fill="{_hex(palette.background)}" />'
    )

    for i, label in enumerate(options):
        cleft = i * (col_w + gap)
        base_color = header_colors[i]

        if outline_only:
            hdr_fill = _hex(palette.background)
            hdr_stroke = _hex(base_color)
            hdr_text = _hex(base_color)
            body_fill = _hex(palette.background)
            body_stroke = _hex(palette.muted)
            body_text = _hex(palette.text)
        else:
            hdr_fill = _hex(base_color)
            hdr_stroke = _hex(base_color)
            hdr_text = _hex(palette.background)
            body_fill = _hex(_lighten(base_color, 0.88))
            body_stroke = _hex(palette.muted)
            body_text = _hex(palette.text)

        # Header rect
        parts.append(
            f'  <rect x="{cleft}" y="0" width="{col_w}" height="{header_h}" '
            f'fill="{hdr_fill}" stroke="{hdr_stroke}" stroke-width="1" />'
        )
        # Header label
        parts.append(
            f'  <text x="{cleft + col_w / 2:.2f}" '
            f'y="{header_h / 2 + heading_pt / 3:.2f}" '
            f'font-family="{heading_font}" font-size="{heading_pt}" '
            f'font-weight="bold" fill="{hdr_text}" '
            f'text-anchor="middle">{label}</text>'
        )

        # Body rect
        parts.append(
            f'  <rect x="{cleft}" y="{header_h}" '
            f'width="{col_w}" height="{body_h}" '
            f'fill="{body_fill}" stroke="{body_stroke}" stroke-width="1" />'
        )

        # Bullets — one <text> per bullet, vertically stacked.
        line_height = body_pt + 6
        bullets = bullets_per_col[i]
        bx = cleft + 12
        by_start = header_h + 16 + body_pt
        for j, bullet in enumerate(bullets):
            by = by_start + j * line_height
            parts.append(
                f'  <text x="{bx}" y="{by:.2f}" '
                f'font-family="{body_font}" font-size="{body_pt}" '
                f'fill="{body_text}">• {bullet}</text>'
            )

    # Dividers — vertical <line>s between columns to make the gap explicit.
    divider_color = _hex(palette.muted)
    for i in range(1, n):
        x = i * col_w + (i - 1) * gap + gap / 2
        parts.append(
            f'  <line x1="{x:.2f}" y1="0" x2="{x:.2f}" y2="{height_px}" '
            f'stroke="{divider_color}" stroke-width="1" '
            f'stroke-dasharray="4,4" />'
        )

    parts.append("</svg>")
    return "\n".join(parts)


ComparisonColumns.render_svg = _render_svg_comparison_columns  # type: ignore[attr-defined]
