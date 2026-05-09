"""Iceberg — visible tip + submerged depths model.

Above the waterline: an MSO_SHAPE.ISOSCELES_TRIANGLE pointing up (the visible
tip, "Visible"). Below: an MSO_SHAPE.PENTAGON. python-pptx's PENTAGON points
upward by default; we render the lower mass as an inverted pentagon by
rotating it 180° via the rotation property. A horizontal RECTANGLE acts as
the waterline. Sub-categories label the underwater region.
"""
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

from .base import AssetShape, ShapeRenderContext
from .shapes_matrix import _hex, _lighten


_UNDERWATER_LABELS = ["Beliefs", "Values", "Assumptions", "Behaviors"]


class Iceberg(AssetShape):
    id = "iceberg"
    name = "Iceberg Model"
    description = "Visible tip above the waterline, hidden depths below — beliefs/values/assumptions/behaviors."
    style_tags = ("framework", "depth", "consulting", "graphics-heavy")
    aspect_ratio_hint = 0.9

    def render(self, slide, ctx: ShapeRenderContext) -> None:
        theme = ctx.theme
        palette = ctx.palette

        outline_only = theme.accent_treatment == "outline" or theme.is_dark

        # Vertical layout: top 30% above water, 6% waterline, 64% below water.
        above_h = int(ctx.height * 0.30)
        line_h = int(ctx.height * 0.04)
        below_h = ctx.height - above_h - line_h

        # Above-water tip — isosceles triangle, narrower than the underwater mass.
        tip_w = int(ctx.width * 0.42)
        tip_left = ctx.left + (ctx.width - tip_w) // 2
        tip_top = ctx.top
        tip = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            tip_left,
            tip_top,
            tip_w,
            above_h,
        )
        tip.fill.solid()
        if outline_only:
            tip.fill.fore_color.rgb = palette.background
            tip.line.color.rgb = palette.primary
            tip_text = palette.text
        else:
            tip.fill.fore_color.rgb = palette.primary
            tip.line.color.rgb = palette.primary
            tip_text = palette.background
        tip.line.width = theme.line_weight_emu

        ttf = tip.text_frame
        ttf.margin_left = Pt(2)
        ttf.margin_right = Pt(2)
        ttf.margin_top = Pt(2)
        ttf.margin_bottom = Pt(2)
        tp = ttf.paragraphs[0]
        tp.alignment = 2  # CENTER
        trun = tp.add_run()
        trun.text = "Visible"
        trun.font.size = Pt(theme.typography.caption_size_pt)
        trun.font.name = theme.typography.heading
        trun.font.bold = True
        trun.font.color.rgb = tip_text

        # Waterline rectangle — full width, accent color.
        line_top = ctx.top + above_h
        waterline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            ctx.left,
            line_top,
            ctx.width,
            line_h,
        )
        waterline.fill.solid()
        if outline_only:
            waterline.fill.fore_color.rgb = palette.background
            waterline.line.color.rgb = palette.accent
        else:
            waterline.fill.fore_color.rgb = _lighten(palette.accent, 0.55)
            waterline.line.color.rgb = palette.accent
        waterline.line.width = theme.line_weight_emu

        wtf = waterline.text_frame
        wtf.margin_left = Pt(4)
        wtf.margin_right = Pt(4)
        wtf.margin_top = Pt(0)
        wtf.margin_bottom = Pt(0)
        wp = wtf.paragraphs[0]
        wp.alignment = 1  # LEFT
        wrun = wp.add_run()
        wrun.text = "Waterline"
        wrun.font.size = Pt(theme.typography.caption_size_pt)
        wrun.font.name = theme.typography.body
        wrun.font.color.rgb = palette.muted

        # Below-water pentagon (inverted). PENTAGON's default orientation is
        # pointed up; rotating 180° flips it to point down — the classic
        # iceberg-under-water silhouette.
        below_w = int(ctx.width * 0.78)
        below_left = ctx.left + (ctx.width - below_w) // 2
        below_top = line_top + line_h
        below = slide.shapes.add_shape(
            MSO_SHAPE.PENTAGON,
            below_left,
            below_top,
            below_w,
            below_h,
        )
        below.rotation = 180.0
        below.fill.solid()
        if outline_only:
            below.fill.fore_color.rgb = palette.background
            below.line.color.rgb = palette.primary
            below_text = palette.text
        else:
            below.fill.fore_color.rgb = _lighten(palette.primary, 0.55)
            below.line.color.rgb = palette.primary
            below_text = palette.text
        below.line.width = theme.line_weight_emu

        # Don't put text directly inside the rotated pentagon (it would also
        # rotate). Instead use 4 textboxes stacked vertically over the pentagon
        # area, each labelled with one underwater concept.
        labels_top = below_top + int(below_h * 0.10)
        labels_h = int(below_h * 0.80)
        each_h = labels_h // len(_UNDERWATER_LABELS)
        for i, lbl in enumerate(_UNDERWATER_LABELS):
            tb = slide.shapes.add_textbox(
                below_left,
                labels_top + i * each_h,
                below_w,
                each_h,
            )
            ltf = tb.text_frame
            ltf.word_wrap = True
            ltf.margin_left = Pt(4)
            ltf.margin_right = Pt(4)
            ltf.margin_top = Pt(2)
            ltf.margin_bottom = Pt(2)
            lp = ltf.paragraphs[0]
            lp.alignment = 2  # CENTER
            lrun = lp.add_run()
            lrun.text = lbl
            lrun.font.size = Pt(theme.typography.caption_size_pt)
            lrun.font.name = theme.typography.heading
            lrun.font.color.rgb = below_text


# ---------------------------------------------------------------------------
# SVG rendering (Web Deck publishing — Sprint J.B)
# ---------------------------------------------------------------------------

def _render_svg_iceberg(
    self: "Iceberg", ctx: ShapeRenderContext, width_px: int, height_px: int
) -> str:
    theme = ctx.theme
    palette = ctx.palette

    outline_only = theme.accent_treatment == "outline" or theme.is_dark

    # Match render(): top 30% above water, 4% waterline, rest below water.
    above_h = int(height_px * 0.30)
    line_h = int(height_px * 0.04)
    below_h = height_px - above_h - line_h

    # Tip triangle
    tip_w = int(width_px * 0.42)
    tip_left = (width_px - tip_w) // 2
    tip_top = 0
    tip_apex_x = tip_left + tip_w / 2.0
    tip_apex_y = tip_top
    tip_bl_x = tip_left
    tip_bl_y = tip_top + above_h
    tip_br_x = tip_left + tip_w
    tip_br_y = tip_top + above_h

    if outline_only:
        tip_fill = _hex(palette.background)
        tip_stroke = _hex(palette.primary)
        tip_text = _hex(palette.text)
    else:
        tip_fill = _hex(palette.primary)
        tip_stroke = _hex(palette.primary)
        tip_text = _hex(palette.background)

    # Waterline rectangle
    line_top = above_h
    if outline_only:
        water_fill = _hex(palette.background)
        water_stroke = _hex(palette.accent)
    else:
        water_fill = _hex(_lighten(palette.accent, 0.55))
        water_stroke = _hex(palette.accent)

    # Below-water inverted pentagon. We mirror the python-pptx PENTAGON
    # (a regular pentagon point-up, then rotated 180 → point-down). Build
    # it as an explicit polygon. Pentagon vertices for a point-up
    # pentagon inscribed in a unit-square bbox approximate to:
    # top apex at (0.5, 0), upper-left (0, 0.38), upper-right (1, 0.38),
    # lower-left (0.18, 1), lower-right (0.82, 1).
    # Rotated 180° about the bbox centre flips top↔bottom and left↔right:
    # bottom apex at (0.5, 1), lower-left (1, 0.62), lower-right (0, 0.62),
    # upper-right (0.82, 0), upper-left (0.18, 0).
    below_w = int(width_px * 0.78)
    below_left = (width_px - below_w) // 2
    below_top = line_top + line_h
    pent_pts = [
        (below_left + 0.18 * below_w, below_top + 0.0 * below_h),
        (below_left + 0.82 * below_w, below_top + 0.0 * below_h),
        (below_left + 1.00 * below_w, below_top + 0.62 * below_h),
        (below_left + 0.50 * below_w, below_top + 1.00 * below_h),
        (below_left + 0.00 * below_w, below_top + 0.62 * below_h),
    ]

    if outline_only:
        below_fill = _hex(palette.background)
        below_stroke = _hex(palette.primary)
        below_text = _hex(palette.text)
    else:
        below_fill = _hex(_lighten(palette.primary, 0.55))
        below_stroke = _hex(palette.primary)
        below_text = _hex(palette.text)

    heading_font = theme.typography.heading
    body_font = theme.typography.body
    cap_pt = theme.typography.caption_size_pt

    parts: list[str] = []
    parts.append(
        f'<svg xmlns="http://www.w3.org/2000/svg" '
        f'width="{width_px}" height="{height_px}" '
        f'viewBox="0 0 {width_px} {height_px}">'
    )
    # Background
    parts.append(
        f'  <rect x="0" y="0" width="{width_px}" height="{height_px}" '
        f'fill="{_hex(palette.background)}" />'
    )

    # Tip triangle polygon
    tip_pts_str = (
        f"{tip_apex_x:.2f},{tip_apex_y:.2f} "
        f"{tip_br_x:.2f},{tip_br_y:.2f} "
        f"{tip_bl_x:.2f},{tip_bl_y:.2f}"
    )
    parts.append(
        f'  <polygon points="{tip_pts_str}" '
        f'fill="{tip_fill}" stroke="{tip_stroke}" stroke-width="1.5" />'
    )
    # Tip label
    tip_label_x = tip_left + tip_w / 2.0
    tip_label_y = tip_top + above_h * 0.65
    parts.append(
        f'  <text x="{tip_label_x:.2f}" y="{tip_label_y:.2f}" '
        f'font-family="{heading_font}" font-size="{cap_pt}" '
        f'font-weight="bold" fill="{tip_text}" '
        f'text-anchor="middle">Visible</text>'
    )

    # Waterline rectangle (acts as horizontal divider)
    parts.append(
        f'  <rect x="0" y="{line_top}" width="{width_px}" height="{line_h}" '
        f'fill="{water_fill}" stroke="{water_stroke}" stroke-width="1.5" />'
    )
    # Waterline label
    parts.append(
        f'  <text x="6" y="{line_top + line_h - 2}" '
        f'font-family="{body_font}" font-size="{cap_pt}" '
        f'fill="{_hex(palette.muted)}">Waterline</text>'
    )

    # Below-water polygon
    pent_pts_str = " ".join(f"{x:.2f},{y:.2f}" for (x, y) in pent_pts)
    parts.append(
        f'  <polygon points="{pent_pts_str}" '
        f'fill="{below_fill}" stroke="{below_stroke}" stroke-width="1.5" />'
    )

    # Underwater labels (stacked vertically over the pentagon).
    labels_top = below_top + int(below_h * 0.10)
    labels_h = int(below_h * 0.80)
    each_h = labels_h // len(_UNDERWATER_LABELS)
    for i, lbl in enumerate(_UNDERWATER_LABELS):
        lx = below_left + below_w / 2.0
        ly = labels_top + i * each_h + each_h / 2.0 + cap_pt / 3.0
        parts.append(
            f'  <text x="{lx:.2f}" y="{ly:.2f}" '
            f'font-family="{heading_font}" font-size="{cap_pt}" '
            f'fill="{below_text}" text-anchor="middle">{lbl}</text>'
        )

    parts.append("</svg>")
    return "\n".join(parts)


Iceberg.render_svg = _render_svg_iceberg  # type: ignore[attr-defined]
