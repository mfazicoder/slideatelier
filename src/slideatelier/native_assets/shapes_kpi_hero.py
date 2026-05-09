"""KpiHero — oversized number callout with a small caption + accent anchor.

Designed to feel hero-like: one giant percentage in the primary color, one
small caption underneath, plus a thin accent rectangle/chevron as a visual
anchor. All native primitives.
"""
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt

from .base import AssetShape, ShapeRenderContext
from .shapes_matrix import _hex


class KpiHero(AssetShape):
    id = "kpi-hero"
    name = "KPI Hero"
    description = "Oversized hero metric (e.g. 42%) with caption and accent anchor."
    style_tags = ("metric", "graphics-heavy", "hero", "contemporary")
    aspect_ratio_hint = 1.6

    def render(self, slide, ctx: ShapeRenderContext) -> None:
        theme = ctx.theme
        palette = ctx.palette

        # Vertical stack: accent anchor (~6%), big number (~62%), caption (~22%).
        anchor_h = int(ctx.height * 0.06)
        number_h = int(ctx.height * 0.62)
        caption_h = ctx.height - anchor_h - number_h - int(ctx.height * 0.02)

        # 1) Accent anchor (chevron when accent_treatment="fill", thin
        # rectangle for outline themes). Centered horizontally, ~22% width.
        anchor_w = int(ctx.width * 0.22)
        anchor_left = ctx.left + (ctx.width - anchor_w) // 2
        anchor_top = ctx.top

        outline_only = theme.accent_treatment == "outline" or theme.is_dark
        if outline_only:
            anchor_kind = MSO_SHAPE.RECTANGLE
        else:
            anchor_kind = MSO_SHAPE.CHEVRON if theme.corner_radius_pct == 0 else MSO_SHAPE.ROUNDED_RECTANGLE

        anchor = slide.shapes.add_shape(
            anchor_kind, anchor_left, anchor_top, anchor_w, anchor_h
        )
        anchor.fill.solid()
        if outline_only:
            anchor.fill.fore_color.rgb = palette.accent
            anchor.line.color.rgb = palette.accent
        else:
            anchor.fill.fore_color.rgb = palette.accent
            anchor.line.color.rgb = palette.accent
        anchor.line.width = theme.line_weight_emu

        # 2) Big number — placed in a wide textbox.
        number_top = ctx.top + anchor_h + int(ctx.height * 0.01)
        num_box = slide.shapes.add_textbox(
            ctx.left,
            number_top,
            ctx.width,
            number_h,
        )
        ntf = num_box.text_frame
        ntf.word_wrap = True
        ntf.margin_left = Pt(4)
        ntf.margin_right = Pt(4)
        ntf.margin_top = Pt(0)
        ntf.margin_bottom = Pt(0)
        np_ = ntf.paragraphs[0]
        np_.alignment = 2  # CENTER
        nrun = np_.add_run()
        nrun.text = "42%"
        # Hero font size: scale by number_h.
        # Approx 1pt = 12700 EMU. Take ~70% of number_h converted to pt.
        hero_pt = max(48, int((number_h / 12700.0) * 0.62))
        nrun.font.size = Pt(hero_pt)
        nrun.font.name = theme.typography.heading
        nrun.font.bold = True
        nrun.font.color.rgb = palette.primary

        # 3) Small caption.
        cap_top = number_top + number_h + int(ctx.height * 0.01)
        cap_box = slide.shapes.add_textbox(
            ctx.left,
            cap_top,
            ctx.width,
            caption_h,
        )
        ctf = cap_box.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Pt(4)
        ctf.margin_right = Pt(4)
        cp = ctf.paragraphs[0]
        cp.alignment = 2  # CENTER
        crun = cp.add_run()
        crun.text = "Conversion lift YoY"
        crun.font.size = Pt(theme.typography.heading_size_pt)
        crun.font.name = theme.typography.body
        crun.font.color.rgb = palette.muted


# ---------------------------------------------------------------------------
# SVG rendering (Web Deck publishing — Sprint J.B)
# ---------------------------------------------------------------------------

def _render_svg_kpi_hero(
    self: "KpiHero", ctx: ShapeRenderContext, width_px: int, height_px: int
) -> str:
    theme = ctx.theme
    palette = ctx.palette

    outline_only = theme.accent_treatment == "outline" or theme.is_dark

    # Mirror render() layout proportions.
    anchor_h = int(height_px * 0.06)
    number_h = int(height_px * 0.62)
    caption_h = height_px - anchor_h - number_h - int(height_px * 0.02)

    anchor_w = int(width_px * 0.22)
    anchor_left = (width_px - anchor_w) // 2
    anchor_top = 0

    # Hero font size: matches the EMU calc in render() (number_h is now px).
    # Approx 1pt = 1.333 px — scale ~ 0.62 of the box height in pt for a
    # large readable hero number.
    hero_pt = max(48, int(number_h * 0.62))
    body_pt = theme.typography.body_size_pt
    heading_pt = theme.typography.heading_size_pt
    heading_font = theme.typography.heading
    body_font = theme.typography.body

    parts: list[str] = []
    parts.append(
        f'<svg xmlns="http://www.w3.org/2000/svg" '
        f'width="{width_px}" height="{height_px}" '
        f'viewBox="0 0 {width_px} {height_px}">'
    )
    parts.append(
        f'  <rect x="0" y="0" width="{width_px}" height="{height_px}" '
        f'fill="{_hex(palette.background)}" />'
    )

    # 1) Accent anchor — chevron polygon (filled themes) or rectangle
    #    (outline themes). Matches the python-pptx CHEVRON / RECTANGLE choice.
    accent_hex = _hex(palette.accent)
    if outline_only:
        # Thin rectangle anchor.
        parts.append(
            f'  <rect x="{anchor_left}" y="{anchor_top}" '
            f'width="{anchor_w}" height="{anchor_h}" '
            f'fill="{accent_hex}" stroke="{accent_hex}" stroke-width="1" />'
        )
    else:
        # Chevron silhouette: pointed-right pentagon.
        ax = anchor_left
        ay = anchor_top
        aw = anchor_w
        ah = anchor_h
        notch = int(ah * 0.5)
        chevron_pts = (
            f"{ax},{ay} "
            f"{ax + aw - notch},{ay} "
            f"{ax + aw},{ay + ah / 2.0:.2f} "
            f"{ax + aw - notch},{ay + ah} "
            f"{ax},{ay + ah}"
        )
        parts.append(
            f'  <polygon points="{chevron_pts}" '
            f'fill="{accent_hex}" stroke="{accent_hex}" stroke-width="1" />'
        )

    # 2) Hero number text. Centred horizontally; baseline placed within
    # the number band.
    number_top = anchor_h + int(height_px * 0.01)
    num_x = width_px / 2.0
    num_y = number_top + number_h * 0.78
    parts.append(
        f'  <text x="{num_x:.2f}" y="{num_y:.2f}" '
        f'font-family="{heading_font}" font-size="{hero_pt}" '
        f'font-weight="bold" fill="{_hex(palette.primary)}" '
        f'text-anchor="middle">42%</text>'
    )

    # 3) Caption text below.
    cap_top = number_top + number_h + int(height_px * 0.01)
    cap_x = width_px / 2.0
    cap_y = cap_top + caption_h / 2.0 + heading_pt / 3.0
    parts.append(
        f'  <text x="{cap_x:.2f}" y="{cap_y:.2f}" '
        f'font-family="{body_font}" font-size="{heading_pt}" '
        f'fill="{_hex(palette.muted)}" '
        f'text-anchor="middle">Conversion lift YoY</text>'
    )

    # Underline accent below caption — keeps the "accent anchor" visual
    # tie even when the top chevron is small.
    underline_y = cap_y + 6
    underline_w = max(60, int(width_px * 0.18))
    underline_x1 = (width_px - underline_w) / 2.0
    underline_x2 = underline_x1 + underline_w
    parts.append(
        f'  <line x1="{underline_x1:.2f}" y1="{underline_y:.2f}" '
        f'x2="{underline_x2:.2f}" y2="{underline_y:.2f}" '
        f'stroke="{accent_hex}" stroke-width="2" />'
    )

    parts.append("</svg>")
    return "\n".join(parts)


KpiHero.render_svg = _render_svg_kpi_hero  # type: ignore[attr-defined]
