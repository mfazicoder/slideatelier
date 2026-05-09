"""ChartBar — native PowerPoint clustered column chart.

Uses slide.shapes.add_chart() with XL_CHART_TYPE.COLUMN_CLUSTERED. Theme
palette is applied to the data series fill so the chart picks up the active
theme. Output is a real chart (shape_type=CHART, integer 3) — fully editable
in PowerPoint, not a freeform.
"""
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from .base import AssetShape, ShapeRenderContext


def _hex(color) -> str:
    """RGBColor → '#RRGGBB' for SVG."""
    return "#{:02X}{:02X}{:02X}".format(color[0], color[1], color[2])


# Categories + values shared between PPT render() and SVG render_svg().
_CATEGORIES = ("Q1", "Q2", "Q3", "Q4")
_VALUES = (12, 18, 24, 31)


class ChartBar(AssetShape):
    id = "chart-bar"
    name = "Bar Chart"
    description = "Native clustered column chart with one series across 4 categories."
    style_tags = ("data", "chart", "comparison", "framework")
    aspect_ratio_hint = 1.6

    def render(self, slide, ctx: ShapeRenderContext) -> None:
        theme = ctx.theme
        palette = ctx.palette

        cd = CategoryChartData()
        cd.categories = list(_CATEGORIES)
        cd.add_series("Revenue", _VALUES)

        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            ctx.left,
            ctx.top,
            ctx.width,
            ctx.height,
            cd,
        )
        chart = chart_shape.chart
        chart.has_title = False
        chart.has_legend = False  # single series; legend redundant

        # Apply theme color to the series.
        plot = chart.plots[0]
        try:
            series = plot.series[0]
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = palette.primary
            line = series.format.line
            line.color.rgb = palette.primary
        except Exception:
            # If python-pptx version doesn't support series.format.line, skip.
            pass


# ---------------------------------------------------------------------------
# SVG rendering (Web Deck publishing — Sprint J.C)
# ---------------------------------------------------------------------------

def _render_svg_chart_bar(
    self: "ChartBar", ctx: ShapeRenderContext, width_px: int, height_px: int
) -> str:
    """Inline SVG mirror of ChartBar — emits <rect> bars + <text> x-axis labels.

    No external chart library; bars are scaled to placeholder _VALUES.
    """
    theme = ctx.theme
    palette = ctx.palette

    categories = list(_CATEGORIES)
    values = list(_VALUES)
    n = len(values)

    # Layout: reserve bottom strip for category labels.
    label_h = max(16, int(height_px * 0.08))
    plot_left = int(width_px * 0.06)
    plot_right = width_px - int(width_px * 0.04)
    plot_top = int(height_px * 0.06)
    plot_bottom = height_px - label_h
    plot_w = plot_right - plot_left
    plot_h = plot_bottom - plot_top

    # Bar geometry — gap = ~30% of slot width.
    slot_w = plot_w / n
    bar_w = slot_w * 0.7
    max_v = max(values) if values else 1

    body_font = theme.typography.body
    cap_pt = theme.typography.caption_size_pt

    parts: list[str] = []
    parts.append(
        f'<svg xmlns="http://www.w3.org/2000/svg" '
        f'width="{width_px}" height="{height_px}" '
        f'viewBox="0 0 {width_px} {height_px}">'
    )
    # Background
    parts.append(
        f'  <rect x="0" y="0" width="{width_px}" height="{height_px}" '
        f'fill="{_hex(palette.background)}" />'
    )

    # Baseline axis line.
    axis_color = _hex(palette.muted)
    parts.append(
        f'  <line x1="{plot_left}" y1="{plot_bottom}" '
        f'x2="{plot_right}" y2="{plot_bottom}" '
        f'stroke="{axis_color}" stroke-width="1" />'
    )

    # Highlight the tallest bar with the accent color; rest use primary.
    primary = _hex(palette.primary)
    accent = _hex(palette.accent)
    highlight_idx = values.index(max_v)

    for i, v in enumerate(values):
        bar_h = (v / max_v) * plot_h if max_v else 0
        slot_x = plot_left + i * slot_w
        bx = slot_x + (slot_w - bar_w) / 2
        by = plot_bottom - bar_h
        fill = accent if i == highlight_idx else primary
        parts.append(
            f'  <rect x="{bx:.2f}" y="{by:.2f}" '
            f'width="{bar_w:.2f}" height="{bar_h:.2f}" '
            f'fill="{fill}" />'
        )

        # Category label below the bar.
        lx = slot_x + slot_w / 2
        ly = plot_bottom + label_h * 0.65
        parts.append(
            f'  <text x="{lx:.2f}" y="{ly:.2f}" '
            f'font-family="{body_font}" font-size="{cap_pt}" '
            f'fill="{_hex(palette.text)}" text-anchor="middle">{categories[i]}</text>'
        )

    parts.append("</svg>")
    return "\n".join(parts)


ChartBar.render_svg = _render_svg_chart_bar  # type: ignore[attr-defined]
