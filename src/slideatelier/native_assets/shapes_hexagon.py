"""HexagonGrid — a 7-cell honeycomb (1 center + 6 surrounding) using native
MSO_SHAPE.HEXAGON primitives. Geometry assumes python-pptx's flat-top
hexagon orientation (the default).

Math (flat-top hexagon, "width" W = flat-to-flat horizontal distance,
"height" H = vertex-to-vertex vertical distance):
    H = W * 2 / sqrt(3)               (so W = H * sqrt(3) / 2)
For two hexagons sharing an edge (touching), the centre-to-centre distances
in the six honeycomb directions are:
    right / left                    : dx = W,            dy = 0
    upper-right / lower-right       : dx =  W * 0.5,     dy = ± H * 0.75
    upper-left  / lower-left        : dx = -W * 0.5,     dy = ± H * 0.75

(Equivalently, neighbour centres lie at distance W from the centre at angles
0°, 60°, 120°, 180°, 240°, 300° measured from the +x axis but with the
vertical axis scaled by 2/sqrt(3) — which is exactly what the formula above
encodes.)
"""
import math

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

from .base import AssetShape, ShapeRenderContext


# Six surrounding labels — 0° is "right", going counter-clockwise.
_SURROUND_LABELS = [
    "PEOPLE",     # right
    "PROCESS",    # upper-right
    "DATA",       # upper-left
    "TECH",       # left
    "CULTURE",    # lower-left
    "STRATEGY",   # lower-right
]


def _lighten(color, factor: float):
    """Blend an RGBColor toward white. factor=0 → original, factor=1 → white."""
    from pptx.dml.color import RGBColor
    r = min(255, int(color[0] + (255 - color[0]) * factor))
    g = min(255, int(color[1] + (255 - color[1]) * factor))
    b = min(255, int(color[2] + (255 - color[2]) * factor))
    return RGBColor(r, g, b)


class HexagonGrid(AssetShape):
    id = "hexagon-grid"
    name = "Hexagon Grid"
    description = "A 7-cell honeycomb (1 core + 6 surrounding pillars) — modern framework diagram."
    style_tags = ("structured", "honeycomb", "contemporary", "framework")
    aspect_ratio_hint = 1.0

    def render(self, slide, ctx: ShapeRenderContext) -> None:
        theme = ctx.theme
        palette = ctx.palette

        # Hexagon "width" W is flat-to-flat distance. We need 3 hex-widths
        # horizontally (left-neighbour + center + right-neighbour) plus a
        # small margin, so W is bounded by ctx.width / 3.
        # Vertically we need: center hex height + 2 * 0.75 * H (the up/down
        # offset of the 60°/120°/240°/300° neighbours), i.e. 2.5 * H.
        #     => H bounded by ctx.height / 2.5
        # Pick the smaller of the two so all 7 fit. Apply a small inset so
        # nothing touches the bbox edges.
        sqrt3 = math.sqrt(3.0)

        # Two candidate sizes (one each from horizontal/vertical constraint):
        w_from_x = ctx.width / 3.0
        h_from_y = ctx.height / 2.5
        # Convert each to the other dimension and pick the limiting one.
        # Given W: H = W * 2 / sqrt(3)
        # Given H: W = H * sqrt(3) / 2
        w_candidate_from_h = h_from_y * sqrt3 / 2.0
        W = min(w_from_x, w_candidate_from_h) * 0.95  # 5% inset
        H = W * 2.0 / sqrt3

        # Center of the bbox
        cx = ctx.left + ctx.width / 2.0
        cy = ctx.top + ctx.height / 2.0

        # Six neighbour centres in honeycomb (flat-top) layout.
        # Order: right, upper-right, upper-left, left, lower-left, lower-right
        neighbour_offsets = [
            ( W,          0.0),
            ( W * 0.5,   -H * 0.75),
            (-W * 0.5,   -H * 0.75),
            (-W,          0.0),
            (-W * 0.5,    H * 0.75),
            ( W * 0.5,    H * 0.75),
        ]

        # Treatment per theme.
        if theme.accent_treatment == "outline" or theme.is_dark:
            outline_only = True
            line_color = palette.primary
        else:
            outline_only = False
            line_color = palette.muted

        # Helper: add a hexagon centred at (px, py) with width W and height H.
        def _add_hex(px: float, py: float, fill_rgb, line_rgb, label: str,
                     label_color, bold: bool = False) -> None:
            left = int(round(px - W / 2.0))
            top = int(round(py - H / 2.0))
            width = int(round(W))
            height = int(round(H))
            shape = slide.shapes.add_shape(
                MSO_SHAPE.HEXAGON, left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_rgb
            shape.line.color.rgb = line_rgb
            shape.line.width = theme.line_weight_emu
            tf = shape.text_frame
            tf.margin_left = Pt(2)
            tf.margin_right = Pt(2)
            tf.margin_top = Pt(2)
            tf.margin_bottom = Pt(2)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = 2  # PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(theme.typography.caption_size_pt)
            run.font.name = theme.typography.heading
            run.font.bold = bold
            run.font.color.rgb = label_color

        # Surrounding hexagons (draw first so the centre sits visually on top).
        for i, (dx, dy) in enumerate(neighbour_offsets):
            if outline_only:
                fill = palette.background
                lbl_color = palette.text
            else:
                # Alternate between accent-tinted and primary-tinted.
                if i % 2 == 0:
                    fill = _lighten(palette.accent, 0.65)
                else:
                    fill = _lighten(palette.primary, 0.80)
                lbl_color = palette.text

            _add_hex(
                cx + dx,
                cy + dy,
                fill_rgb=fill,
                line_rgb=line_color,
                label=_SURROUND_LABELS[i],
                label_color=lbl_color,
                bold=False,
            )

        # Centre hexagon — always solid primary, "CORE" in background colour.
        if outline_only and theme.is_dark:
            # On dark themes use primary as a strong solid fill against the
            # dark background, with the dark background colour for the text.
            centre_fill = palette.primary
            centre_text = palette.background
        elif outline_only:
            # Pure outline theme (light) — keep it outline too for consistency.
            centre_fill = palette.background
            centre_text = palette.primary
        else:
            centre_fill = palette.primary
            centre_text = palette.background

        _add_hex(
            cx,
            cy,
            fill_rgb=centre_fill,
            line_rgb=palette.primary,
            label="CORE",
            label_color=centre_text,
            bold=True,
        )

    def render_svg(self, ctx: ShapeRenderContext, width_px: int, height_px: int) -> str:
        """Inline SVG mirroring render(): 7 flat-top hexagon <polygon>s
        (1 centre + 6 honeycomb neighbours), each with 6 points.
        """
        theme = ctx.theme
        palette = ctx.palette

        sqrt3 = math.sqrt(3.0)

        w_from_x = width_px / 3.0
        h_from_y = height_px / 2.5
        w_candidate_from_h = h_from_y * sqrt3 / 2.0
        W = min(w_from_x, w_candidate_from_h) * 0.95
        H = W * 2.0 / sqrt3

        cx = width_px / 2.0
        cy = height_px / 2.0

        neighbour_offsets = [
            (W, 0.0),
            (W * 0.5, -H * 0.75),
            (-W * 0.5, -H * 0.75),
            (-W, 0.0),
            (-W * 0.5, H * 0.75),
            (W * 0.5, H * 0.75),
        ]

        if theme.accent_treatment == "outline" or theme.is_dark:
            outline_only = True
            line_color = palette.primary
        else:
            outline_only = False
            line_color = palette.muted

        cap_pt = theme.typography.caption_size_pt
        heading_font = theme.typography.heading

        def _hex_points(px: float, py: float) -> str:
            """Flat-top hexagon: vertices at top, top-right, bottom-right,
            bottom, bottom-left, top-left. W = flat-to-flat horizontal,
            H = vertex-to-vertex vertical. Quarter-height steps for the slanted
            corners match python-pptx's MSO_SHAPE.HEXAGON geometry."""
            half_w = W / 2.0
            half_h = H / 2.0
            quarter_h = H / 4.0
            verts = [
                (px,          py - half_h),       # top vertex
                (px + half_w, py - quarter_h),    # top-right
                (px + half_w, py + quarter_h),    # bottom-right
                (px,          py + half_h),       # bottom vertex
                (px - half_w, py + quarter_h),    # bottom-left
                (px - half_w, py - quarter_h),    # top-left
            ]
            return " ".join(f"{x:.2f},{y:.2f}" for x, y in verts)

        parts: list[str] = []
        parts.append(
            f'<svg xmlns="http://www.w3.org/2000/svg" '
            f'width="{width_px}" height="{height_px}" '
            f'viewBox="0 0 {width_px} {height_px}">'
        )
        parts.append(
            f'  <rect x="0" y="0" width="{width_px}" height="{height_px}" '
            f'fill="{_hex_color(palette.background)}" />'
        )

        # Six surrounding hexagons.
        for i, (dx, dy) in enumerate(neighbour_offsets):
            if outline_only:
                fill = palette.background
                lbl_color = palette.text
            else:
                if i % 2 == 0:
                    fill = _lighten(palette.accent, 0.65)
                else:
                    fill = _lighten(palette.primary, 0.80)
                lbl_color = palette.text

            px = cx + dx
            py = cy + dy
            parts.append(
                f'  <polygon points="{_hex_points(px, py)}" '
                f'fill="{_hex_color(fill)}" stroke="{_hex_color(line_color)}" '
                f'stroke-width="1" />'
            )
            label = _SURROUND_LABELS[i]
            ty = py + cap_pt // 3
            parts.append(
                f'  <text x="{px:.2f}" y="{ty:.2f}" '
                f'font-family="{heading_font}" font-size="{cap_pt}" '
                f'fill="{_hex_color(lbl_color)}" text-anchor="middle">{label}</text>'
            )

        # Centre hexagon.
        if outline_only and theme.is_dark:
            centre_fill = palette.primary
            centre_text = palette.background
        elif outline_only:
            centre_fill = palette.background
            centre_text = palette.primary
        else:
            centre_fill = palette.primary
            centre_text = palette.background

        parts.append(
            f'  <polygon points="{_hex_points(cx, cy)}" '
            f'fill="{_hex_color(centre_fill)}" '
            f'stroke="{_hex_color(palette.primary)}" stroke-width="1" />'
        )
        ty = cy + cap_pt // 3
        parts.append(
            f'  <text x="{cx:.2f}" y="{ty:.2f}" '
            f'font-family="{heading_font}" font-size="{cap_pt}" '
            f'font-weight="bold" '
            f'fill="{_hex_color(centre_text)}" text-anchor="middle">CORE</text>'
        )

        parts.append("</svg>")
        return "\n".join(parts)


def _hex_color(color) -> str:
    """RGBColor → '#RRGGBB' for SVG."""
    return "#{:02X}{:02X}{:02X}".format(color[0], color[1], color[2])
