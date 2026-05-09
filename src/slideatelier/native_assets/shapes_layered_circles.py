"""LayeredCircles — 3 overlapping ovals forming a venn-style diagram.

Native MSO_SHAPE.OVAL primitives. To make overlap regions visible, all three
circles use semi-transparent-feeling tints (we approximate transparency by
using lightened palette colors) and an outline-only treatment falls back to
just outlined circles.
"""
import math

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

from .base import AssetShape, ShapeRenderContext
from .shapes_matrix import _hex, _lighten


_LABELS = ["People", "Process", "Product"]


class LayeredCircles(AssetShape):
    id = "layered-circles"
    name = "Layered Circles"
    description = "3 overlapping ovals (venn-style) for tri-factor frameworks."
    style_tags = ("comparison", "venn", "structured", "contemporary")
    aspect_ratio_hint = 1.2

    def render(self, slide, ctx: ShapeRenderContext) -> None:
        theme = ctx.theme
        palette = ctx.palette

        outline_only = theme.accent_treatment == "outline" or theme.is_dark

        # Three circles in a triangular arrangement: 2 on top, 1 on bottom.
        # Diameter ~52% of bbox width so they overlap meaningfully.
        d = int(min(ctx.width * 0.55, ctx.height * 0.65))
        cx = ctx.left + ctx.width / 2.0
        cy = ctx.top + ctx.height / 2.0

        # Centers: equilateral triangle with side ~ 0.55 * d so adjacent
        # circles overlap by ~45%.
        side = d * 0.62
        # Circle 0 (top-left), 1 (top-right), 2 (bottom-center)
        centers = [
            (cx - side / 2.0, cy - side * math.sqrt(3.0) / 6.0),
            (cx + side / 2.0, cy - side * math.sqrt(3.0) / 6.0),
            (cx, cy + side * math.sqrt(3.0) / 3.0),
        ]

        # Per-circle base color: primary, accent, muted (or theme variants).
        base_colors = [palette.primary, palette.accent, palette.muted]

        for i, ((px, py), base, label) in enumerate(zip(centers, base_colors, _LABELS)):
            left = int(round(px - d / 2.0))
            top = int(round(py - d / 2.0))
            oval = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, left, top, d, d
            )
            oval.fill.solid()

            if outline_only:
                oval.fill.fore_color.rgb = palette.background
                oval.line.color.rgb = base
                lbl_color = base
            else:
                # Use a heavily lightened tint so overlaps are visually
                # distinct (approximating transparency without freeform).
                oval.fill.fore_color.rgb = _lighten(base, 0.65)
                oval.line.color.rgb = base
                lbl_color = palette.text
            oval.line.width = theme.line_weight_emu

            tf = oval.text_frame
            tf.margin_left = Pt(8)
            tf.margin_right = Pt(8)
            tf.margin_top = Pt(8)
            tf.margin_bottom = Pt(8)
            p = tf.paragraphs[0]
            p.alignment = 2  # CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(theme.typography.body_size_pt)
            run.font.name = theme.typography.heading
            run.font.bold = True
            run.font.color.rgb = lbl_color


# ---------------------------------------------------------------------------
# SVG rendering (Web Deck publishing — Sprint J.B)
# ---------------------------------------------------------------------------

def _render_svg_layered_circles(
    self: "LayeredCircles", ctx: ShapeRenderContext, width_px: int, height_px: int
) -> str:
    theme = ctx.theme
    palette = ctx.palette

    outline_only = theme.accent_treatment == "outline" or theme.is_dark

    # Mirror render(): three circles in equilateral triangle, diameter
    # ~52% of the bbox.
    d = int(min(width_px * 0.55, height_px * 0.65))
    cx = width_px / 2.0
    cy = height_px / 2.0
    side = d * 0.62
    centers = [
        (cx - side / 2.0, cy - side * math.sqrt(3.0) / 6.0),
        (cx + side / 2.0, cy - side * math.sqrt(3.0) / 6.0),
        (cx, cy + side * math.sqrt(3.0) / 3.0),
    ]
    base_colors = [palette.primary, palette.accent, palette.muted]

    heading_font = theme.typography.heading
    body_pt = theme.typography.body_size_pt

    parts: list[str] = []
    parts.append(
        f'<svg xmlns="http://www.w3.org/2000/svg" '
        f'width="{width_px}" height="{height_px}" '
        f'viewBox="0 0 {width_px} {height_px}">'
    )
    parts.append(
        f'  <rect x="0" y="0" width="{width_px}" height="{height_px}" '
        f'fill="{_hex(palette.background)}" />'
    )

    r = d / 2.0
    for i, ((px, py), base, label) in enumerate(zip(centers, base_colors, _LABELS)):
        if outline_only:
            fill = _hex(palette.background)
            stroke = _hex(base)
            lbl_color = _hex(base)
        else:
            fill = _hex(_lighten(base, 0.65))
            stroke = _hex(base)
            lbl_color = _hex(palette.text)
        parts.append(
            f'  <circle cx="{px:.2f}" cy="{py:.2f}" r="{r:.2f}" '
            f'fill="{fill}" stroke="{stroke}" stroke-width="1.5" '
            f'fill-opacity="0.65" />'
        )
        parts.append(
            f'  <text x="{px:.2f}" y="{py + body_pt / 3.0:.2f}" '
            f'font-family="{heading_font}" font-size="{body_pt}" '
            f'font-weight="bold" fill="{lbl_color}" '
            f'text-anchor="middle">{label}</text>'
        )

    parts.append("</svg>")
    return "\n".join(parts)


LayeredCircles.render_svg = _render_svg_layered_circles  # type: ignore[attr-defined]
