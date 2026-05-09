"""Matrix 2x2 — a four-quadrant framework with axis labels.

Renders four equal rectangles in a 2x2 grid with axis labels. The treatment
adapts to the theme: filled tints for consulting-corporate, outline-only for
scaffold-dark, etc.
"""
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

from .base import AssetShape, ShapeRenderContext


class Matrix2x2(AssetShape):
    id = "matrix-2x2"
    name = "2×2 Matrix"
    description = "Four-quadrant framework with axis labels (BCG-style, prioritization, importance × urgency)."
    style_tags = ("framework", "structured", "comparison", "consulting")
    aspect_ratio_hint = 1.0  # square-ish

    def render(self, slide, ctx: ShapeRenderContext) -> None:
        theme = ctx.theme
        palette = ctx.palette

        # Reserve 12% on left for y-axis label, 12% on bottom for x-axis label
        axis_label_emu = int(min(ctx.width, ctx.height) * 0.10)
        grid_left = ctx.left + axis_label_emu
        grid_top = ctx.top
        grid_w = ctx.width - axis_label_emu
        grid_h = ctx.height - axis_label_emu

        cell_w = grid_w // 2
        cell_h = grid_h // 2

        # Quadrant fill colors per theme treatment
        if theme.accent_treatment == "outline" or theme.is_dark:
            fill_colors = [palette.background] * 4
            line_color = palette.primary
        else:
            # subtle tinted fills for consulting-corporate
            fill_colors = [
                palette.background,
                _lighten(palette.primary, 0.92),
                _lighten(palette.primary, 0.92),
                palette.primary,
            ]
            line_color = palette.muted

        for r in range(2):
            for c in range(2):
                idx = r * 2 + c
                shape_kind = MSO_SHAPE.ROUNDED_RECTANGLE if theme.corner_radius_pct > 0 else MSO_SHAPE.RECTANGLE
                rect = slide.shapes.add_shape(
                    shape_kind,
                    grid_left + c * cell_w,
                    grid_top + r * cell_h,
                    cell_w,
                    cell_h,
                )
                rect.fill.solid()
                rect.fill.fore_color.rgb = fill_colors[idx]
                rect.line.color.rgb = line_color
                rect.line.width = theme.line_weight_emu
                # Keep label area visible — small text in top-left of each cell
                tf = rect.text_frame
                tf.margin_left = Pt(8)
                tf.margin_top = Pt(8)
                tf.margin_right = Pt(8)
                tf.margin_bottom = Pt(8)
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = ["Low / Low", "High / Low", "Low / High", "High / High"][idx]
                run.font.size = Pt(theme.typography.body_size_pt)
                run.font.name = theme.typography.heading
                # Adapt label color for dark fills
                if idx == 3 and not theme.is_dark and theme.accent_treatment != "outline":
                    run.font.color.rgb = palette.background
                else:
                    run.font.color.rgb = palette.text

        # Y-axis label (rotated text — but python-pptx doesn't easily rotate
        # text frames cleanly, so we use a vertical column of text on the left)
        y_axis = slide.shapes.add_textbox(
            ctx.left,
            ctx.top,
            axis_label_emu,
            grid_h,
        )
        ytf = y_axis.text_frame
        ytf.word_wrap = True
        yp = ytf.paragraphs[0]
        yp.alignment = 2  # PP_ALIGN.RIGHT
        yrun = yp.add_run()
        yrun.text = "Y axis →"
        yrun.font.size = Pt(theme.typography.caption_size_pt)
        yrun.font.color.rgb = palette.muted
        yrun.font.name = theme.typography.body

        # X-axis label (along the bottom)
        x_axis = slide.shapes.add_textbox(
            grid_left,
            grid_top + grid_h,
            grid_w,
            axis_label_emu,
        )
        xtf = x_axis.text_frame
        xp = xtf.paragraphs[0]
        xrun = xp.add_run()
        xrun.text = "X axis →"
        xrun.font.size = Pt(theme.typography.caption_size_pt)
        xrun.font.color.rgb = palette.muted
        xrun.font.name = theme.typography.body

    def render_svg(self, ctx: ShapeRenderContext, width_px: int, height_px: int) -> str:
        """Inline SVG mirroring render(): 4 <rect>s in a 2x2 grid + 2 axis labels."""
        theme = ctx.theme
        palette = ctx.palette

        # Mirror the EMU layout: reserve ~10% on left/bottom for axis labels.
        axis_label = int(min(width_px, height_px) * 0.10)
        grid_left = axis_label
        grid_top = 0
        grid_w = width_px - axis_label
        grid_h = height_px - axis_label

        cell_w = grid_w // 2
        cell_h = grid_h // 2

        if theme.accent_treatment == "outline" or theme.is_dark:
            fill_colors = [palette.background] * 4
            line_color = palette.primary
        else:
            fill_colors = [
                palette.background,
                _lighten(palette.primary, 0.92),
                _lighten(palette.primary, 0.92),
                palette.primary,
            ]
            line_color = palette.muted

        cell_labels = ["Quadrant 1", "Quadrant 2", "Quadrant 3", "Quadrant 4"]

        body_font = theme.typography.body
        heading_font = theme.typography.heading
        body_pt = theme.typography.body_size_pt
        cap_pt = theme.typography.caption_size_pt

        parts: list[str] = []
        parts.append(
            f'<svg xmlns="http://www.w3.org/2000/svg" '
            f'width="{width_px}" height="{height_px}" '
            f'viewBox="0 0 {width_px} {height_px}">'
        )
        parts.append(
            f'  <rect x="0" y="0" width="{width_px}" height="{height_px}" '
            f'fill="{_hex(palette.background)}" />'
        )

        for r in range(2):
            for c in range(2):
                idx = r * 2 + c
                x = grid_left + c * cell_w
                y = grid_top + r * cell_h
                fill = _hex(fill_colors[idx])
                stroke = _hex(line_color)
                parts.append(
                    f'  <rect x="{x}" y="{y}" width="{cell_w}" height="{cell_h}" '
                    f'fill="{fill}" stroke="{stroke}" stroke-width="1" />'
                )
                label = cell_labels[idx]
                if idx == 3 and not theme.is_dark and theme.accent_treatment != "outline":
                    txt_color = _hex(palette.background)
                else:
                    txt_color = _hex(palette.text)
                tx = x + 8
                ty = y + 8 + body_pt
                parts.append(
                    f'  <text x="{tx}" y="{ty}" '
                    f'font-family="{heading_font}" font-size="{body_pt}" '
                    f'fill="{txt_color}">{label}</text>'
                )

        # Y-axis label rotated -90 around its anchor.
        y_axis_x = axis_label - 4
        y_axis_y = grid_top + grid_h // 2
        parts.append(
            f'  <text x="{y_axis_x}" y="{y_axis_y}" '
            f'font-family="{body_font}" font-size="{cap_pt}" '
            f'fill="{_hex(palette.muted)}" text-anchor="middle" '
            f'transform="rotate(-90 {y_axis_x} {y_axis_y})">Y axis →</text>'
        )

        # X-axis label along the bottom.
        x_axis_x = grid_left + grid_w // 2
        x_axis_y = grid_top + grid_h + axis_label // 2 + cap_pt // 2
        parts.append(
            f'  <text x="{x_axis_x}" y="{x_axis_y}" '
            f'font-family="{body_font}" font-size="{cap_pt}" '
            f'fill="{_hex(palette.muted)}" text-anchor="middle">X axis →</text>'
        )

        parts.append("</svg>")
        return "\n".join(parts)


def _lighten(color, factor: float):
    """Blend an RGBColor toward white. factor=0 → original, factor=1 → white."""
    from pptx.dml.color import RGBColor
    r = min(255, int(color[0] + (255 - color[0]) * factor))
    g = min(255, int(color[1] + (255 - color[1]) * factor))
    b = min(255, int(color[2] + (255 - color[2]) * factor))
    return RGBColor(r, g, b)


def _hex(color) -> str:
    """RGBColor → '#RRGGBB' for SVG."""
    return "#{:02X}{:02X}{:02X}".format(color[0], color[1], color[2])
