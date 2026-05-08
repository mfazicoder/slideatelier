"""Executive Summary — the classic consulting summary slide.

Layout: title at top, a key-takeaway callout band in the middle, and 2-3
columns of supporting bullets below. All shapes are native MSO_SHAPE
primitives — fully editable in PowerPoint.
"""
from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

from .base import SlideTemplate, Theme
from .shapes_matrix import _lighten


class ExecSummary(SlideTemplate):
    id = "exec-summary"
    name = "Executive Summary"
    description = "Classic consulting summary: title + key takeaway callout + 2-3 columns of supporting bullets."
    style_tags = ("structured", "consulting", "summary", "professional")
    requires_image = False

    def render(
        self,
        slide,
        theme: Theme,
        content: dict,
        *,
        slide_width_emu: int,
        slide_height_emu: int,
    ) -> None:
        palette = theme.palette
        typo = theme.typography

        title = content.get("title", "")
        takeaway = content.get("takeaway", "")
        columns = content.get("columns", []) or []
        n_cols = max(1, len(columns))

        # ------------------------------------------------------------------
        # Vertical layout: 18% title bar, 22% callout, 60% columns.
        # ------------------------------------------------------------------
        margin_x = int(slide_width_emu * 0.06)
        content_w = slide_width_emu - 2 * margin_x

        title_band_top = int(slide_height_emu * 0.05)
        title_band_h = int(slide_height_emu * 0.13)

        callout_top = int(slide_height_emu * 0.20)
        callout_h = int(slide_height_emu * 0.18)

        cols_top = int(slide_height_emu * 0.42)
        cols_h = int(slide_height_emu * 0.52)

        # ------------------------------------------------------------------
        # 1. Title bar.
        # ------------------------------------------------------------------
        title_box = slide.shapes.add_textbox(
            margin_x, title_band_top, content_w, title_band_h
        )
        ttf = title_box.text_frame
        ttf.word_wrap = True
        ttf.margin_left = Pt(0)
        ttf.margin_right = Pt(0)
        ttf.margin_top = Pt(0)
        ttf.margin_bottom = Pt(0)
        tp = ttf.paragraphs[0]
        tp.alignment = 1  # LEFT
        trun = tp.add_run()
        trun.text = title
        trun.font.name = typo.heading
        trun.font.size = Pt(typo.title_size_pt)
        trun.font.bold = True
        trun.font.color.rgb = palette.text if theme.is_dark else palette.primary

        # Thin underline rule in accent color, sitting just below the title.
        rule_h = Pt(3)
        rule_w = int(content_w * 0.12)
        rule = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            margin_x,
            title_band_top + title_band_h + int(Pt(2)),
            rule_w,
            int(rule_h),
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = palette.accent
        rule.line.fill.background()

        # ------------------------------------------------------------------
        # 2. Key-takeaway callout band.
        # ------------------------------------------------------------------
        callout_kind = (
            MSO_SHAPE.ROUNDED_RECTANGLE
            if theme.corner_radius_pct > 0
            else MSO_SHAPE.RECTANGLE
        )
        callout = slide.shapes.add_shape(
            callout_kind, margin_x, callout_top, content_w, callout_h
        )

        if theme.is_dark or theme.accent_treatment == "outline":
            # Outline-only treatment: transparent background look — use the
            # theme background as the fill so it blends with the slide.
            callout.fill.solid()
            callout.fill.fore_color.rgb = palette.background
            callout.line.color.rgb = palette.accent
            callout.line.width = theme.line_weight_emu
            callout_text_color = palette.text
        elif theme.accent_treatment == "fill":
            # Light tinted fill from primary, no border.
            callout.fill.solid()
            callout.fill.fore_color.rgb = _lighten(palette.primary, 0.85)
            callout.line.fill.background()
            callout_text_color = palette.primary
        else:
            # Underline / shadow treatments: gentle accent tint.
            callout.fill.solid()
            callout.fill.fore_color.rgb = _lighten(palette.accent, 0.80)
            callout.line.fill.background()
            callout_text_color = palette.text

        # Inset text inside the callout.
        ctf = callout.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Pt(18)
        ctf.margin_right = Pt(18)
        ctf.margin_top = Pt(12)
        ctf.margin_bottom = Pt(12)
        cp = ctf.paragraphs[0]
        cp.alignment = 1  # LEFT
        crun = cp.add_run()
        crun.text = takeaway
        crun.font.name = typo.body
        crun.font.size = Pt(typo.body_size_pt + 4)
        crun.font.bold = True
        crun.font.color.rgb = callout_text_color

        # ------------------------------------------------------------------
        # 3. Columns — equal-width split across content_w with vertical
        #    separator rules between columns.
        # ------------------------------------------------------------------
        gutter = int(content_w * 0.025)
        total_gutter = gutter * (n_cols - 1)
        col_w = (content_w - total_gutter) // n_cols

        for i, col in enumerate(columns):
            col_left = margin_x + i * (col_w + gutter)

            # Vertical separator before columns 2+ (sits in the gutter).
            if i > 0:
                sep_left = col_left - gutter // 2
                sep = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    sep_left,
                    cols_top + int(cols_h * 0.05),
                    int(Pt(0.75)),
                    int(cols_h * 0.90),
                )
                sep.fill.solid()
                sep.fill.fore_color.rgb = palette.muted
                sep.line.fill.background()

            heading_text = col.get("heading", "")
            bullets = col.get("bullets", []) or []

            # Column heading.
            heading_h = int(cols_h * 0.14)
            head_box = slide.shapes.add_textbox(
                col_left, cols_top, col_w, heading_h
            )
            htf = head_box.text_frame
            htf.word_wrap = True
            htf.margin_left = Pt(0)
            htf.margin_right = Pt(0)
            htf.margin_top = Pt(0)
            htf.margin_bottom = Pt(0)
            hp = htf.paragraphs[0]
            hp.alignment = 1
            hrun = hp.add_run()
            hrun.text = heading_text
            hrun.font.name = typo.heading
            hrun.font.size = Pt(typo.heading_size_pt)
            hrun.font.bold = True
            hrun.font.color.rgb = palette.primary if not theme.is_dark else palette.primary

            # Bullets list.
            bullets_top = cols_top + heading_h + int(Pt(4))
            bullets_h = cols_h - heading_h - int(Pt(4))
            bul_box = slide.shapes.add_textbox(
                col_left, bullets_top, col_w, bullets_h
            )
            btf = bul_box.text_frame
            btf.word_wrap = True
            btf.margin_left = Pt(0)
            btf.margin_right = Pt(0)
            btf.margin_top = Pt(0)
            btf.margin_bottom = Pt(0)

            for bi, bullet in enumerate(bullets):
                if bi == 0:
                    bp = btf.paragraphs[0]
                else:
                    bp = btf.add_paragraph()
                bp.alignment = 1
                bp.space_after = Pt(4)
                brun = bp.add_run()
                brun.text = f"•  {bullet}"
                brun.font.name = typo.body
                brun.font.size = Pt(typo.body_size_pt)
                brun.font.color.rgb = palette.text
