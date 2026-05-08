from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .asset_copier import copy_slide_into, copy_slide_shapes_onto
from .library import load_catalog
from .models import Slide, SlideDeck
from .template import Template, hex_to_rgb


# ============================================================================
# Public API
# ============================================================================

class DeckRenderer:
    def __init__(
        self,
        template: Template,
        templates_dir: Path | None = None,
        library_catalog_path: Path | None = None,
    ):
        self.tpl = template
        self.templates_dir = templates_dir or Path("./templates")
        self.library_catalog_path = library_catalog_path or Path("./library/catalog.json")
        self._catalog_cache = None

    def _catalog(self):
        """Lazy-load the catalog only if a slide in the deck references an asset."""
        if self._catalog_cache is None:
            if self.library_catalog_path.exists():
                self._catalog_cache = load_catalog(self.library_catalog_path)
            else:
                self._catalog_cache = False  # sentinel — catalog missing
        return self._catalog_cache or None

    def render(self, deck: SlideDeck, output_path: Path) -> None:
        # Load catalog if any slide refers to library assets — either as a
        # full-slide replacement (asset_ref) or as a growable extra.
        needs_catalog = any(s.asset_ref for s in deck.slides) or any(
            e.type == "library_asset"
            for s in deck.slides
            for e in (s.extras or [])
        )
        catalog = self._catalog() if needs_catalog else None

        if self.tpl.master_pptx:
            MasterRenderer(self.tpl, self.templates_dir, catalog).render(deck, output_path)
        else:
            JSONRenderer(self.tpl, catalog).render(deck, output_path)


def render_deck(
    deck: SlideDeck,
    output_path: Path,
    template: Template | None = None,
    templates_dir: Path | None = None,
) -> None:
    """Convenience wrapper. Uses a built-in default Template if none provided."""
    if template is None:
        template = Template()
    DeckRenderer(template, templates_dir).render(deck, output_path)


# ============================================================================
# MASTER MODE — renders into a designer-provided master.pptx
# ============================================================================

class MasterRenderer:
    """Opens a master .pptx and adds generated slides using its master layouts.
    All visual decisions (colors, fonts, logos, footers, theme) are inherited
    from the master file."""

    def __init__(self, template: Template, templates_dir: Path, catalog=None):
        self.tpl = template
        self.templates_dir = templates_dir
        self.catalog = catalog
        self.master_path = self._resolve_master_path(template.master_pptx)

    def _resolve_master_path(self, raw: str) -> Path:
        p = Path(raw)
        if p.is_absolute():
            return p
        # Try relative to templates dir first, then to CWD
        rel_to_templates = self.templates_dir / p
        if rel_to_templates.exists():
            return rel_to_templates
        return p

    def render(self, deck: SlideDeck, output_path: Path) -> None:
        if not self.master_path.exists():
            raise FileNotFoundError(
                f"Master .pptx not found: {self.master_path}. "
                f"Check the template's master_pptx path."
            )
        prs = Presentation(str(self.master_path))

        for slide_data in deck.slides:
            self._render_slide(prs, slide_data, deck)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

    def _render_slide(self, prs: Presentation, slide_data: Slide, deck: SlideDeck) -> None:
        # Asset reference — copy the library slide directly (preserves native shapes).
        if slide_data.asset_ref and self.catalog:
            asset = self.catalog.find(slide_data.asset_ref)
            if asset is not None:
                copy_slide_into(Path(asset.file_path), asset.slide_index, prs, match_dimensions=False)
                # Note: text substitution into the asset is a v0.7 feature (phase 5).
                # For now the asset's placeholder text remains as-is.
                if slide_data.speaker_notes:
                    new_slide = prs.slides[-1]
                    new_slide.notes_slide.notes_text_frame.text = slide_data.speaker_notes
                return

        layout_idx = self.tpl.layout_map.get(slide_data.layout, 1)  # fallback to "Title and Content"
        if layout_idx >= len(prs.slide_layouts):
            raise IndexError(
                f"Layout index {layout_idx} for {slide_data.layout!r} is out of range. "
                f"Master has {len(prs.slide_layouts)} layouts. "
                f"Run `atelier template inspect` to see what's available."
            )
        layout = prs.slide_layouts[layout_idx]
        s = prs.slides.add_slide(layout)

        if slide_data.layout == "title":
            self._fill_title_slide(s, deck)
        elif slide_data.layout == "two_column":
            self._fill_two_column(s, slide_data)
        else:
            # Default: title + body content (covers content, bullet_list, section_divider, key_takeaway, closing)
            self._fill_title_and_body(s, slide_data)

        # Render any growable extras (library_asset, etc.) AFTER layout content.
        self._render_extras(prs, s, slide_data)

        if slide_data.speaker_notes:
            s.notes_slide.notes_text_frame.text = slide_data.speaker_notes

    def _render_extras(self, prs: Presentation, slide, slide_data: Slide) -> None:
        """Stamp library-asset extras onto the slide, honoring `bbox` and
        `color_override` from extra.config when present."""
        slide_w = prs.slide_width
        slide_h = prs.slide_height
        for extra in (slide_data.extras or []):
            if extra.type != "library_asset":
                continue
            asset_ref = (extra.config or {}).get("asset_ref")
            if not asset_ref or self.catalog is None:
                continue
            asset = self.catalog.find(asset_ref)
            if asset is None:
                continue
            target = _resolve_extra_target(extra, slide_w, slide_h)
            recolor = _resolve_color_override(extra)
            try:
                copy_slide_shapes_onto(
                    Path(asset.file_path),
                    asset.slide_index,
                    slide,
                    target_left_emu=int(target[0]),
                    target_top_emu=int(target[1]),
                    target_width_emu=int(target[2]),
                    target_height_emu=int(target[3]),
                    recolor_to=recolor,
                )
            except Exception as e:  # noqa: BLE001
                import sys
                print(f"[renderer] failed to stamp extra {asset_ref}: {e}", file=sys.stderr)
                continue

    def _fill_title_slide(self, slide, deck: SlideDeck):
        title_ph = _find_placeholder(slide, [PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.TITLE])
        if title_ph is not None:
            title_ph.text = deck.title
        if deck.subtitle:
            sub_ph = _find_placeholder(slide, [PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.BODY])
            if sub_ph is not None:
                sub_ph.text = deck.subtitle

    def _fill_title_and_body(self, slide, slide_data: Slide):
        title_ph = _find_placeholder(slide, [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE])
        if title_ph is not None:
            title_ph.text = slide_data.title
        body_ph = _find_placeholder(
            slide, [PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.BODY], skip=[title_ph] if title_ph else []
        )
        if body_ph is not None and slide_data.body:
            _fill_text_frame_with_bullets(body_ph.text_frame, slide_data.body)

    def _fill_two_column(self, slide, slide_data: Slide):
        title_ph = _find_placeholder(slide, [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE])
        if title_ph is not None:
            title_ph.text = slide_data.title
        # Find two content placeholders by sorted (idx) order, skipping the title
        content_phs = sorted(
            (
                p for p in slide.placeholders
                if p.placeholder_format.type in (PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.BODY)
                and p is not title_ph
            ),
            key=lambda p: p.placeholder_format.idx,
        )
        if len(content_phs) >= 1 and slide_data.body_left:
            _fill_text_frame_with_bullets(content_phs[0].text_frame, slide_data.body_left)
        if len(content_phs) >= 2 and slide_data.body_right:
            _fill_text_frame_with_bullets(content_phs[1].text_frame, slide_data.body_right)


def _resolve_extra_target(extra, slide_width_emu: int, slide_height_emu: int) -> tuple:
    """Compute (left, top, width, height) in EMU for an extra.

    Priority: bbox (normalized 0..1 rect) → position-based default fraction.
    The bbox is interpreted relative to the slide canvas. position values fall
    back to the same fractional defaults the WYSIWYG overlay uses (right=40%
    column, below=34% strip, left=40% column, inline=centered 50% box)."""
    cfg = extra.config or {} if hasattr(extra, "config") else {}
    bbox = cfg.get("bbox") if isinstance(cfg, dict) else None
    if isinstance(bbox, dict):
        try:
            left = float(bbox.get("left", 0.0))
            top = float(bbox.get("top", 0.0))
            width = float(bbox.get("width", 0.0))
            height = float(bbox.get("height", 0.0))
        except (TypeError, ValueError):
            bbox = None
        else:
            # Clamp to [0,1] so a misbehaving client can't push shapes off-canvas.
            left = max(0.0, min(1.0, left))
            top = max(0.0, min(1.0, top))
            width = max(0.01, min(1.0 - left, width))
            height = max(0.01, min(1.0 - top, height))
            return (
                int(slide_width_emu * left),
                int(slide_height_emu * top),
                int(slide_width_emu * width),
                int(slide_height_emu * height),
            )

    pos = getattr(extra, "position", "right")
    if pos == "right":
        # right column: starts at 50%, takes 48% width, vertically inset 18%-95%
        frac = (0.50, 0.18, 0.48, 0.77)
    elif pos == "left":
        frac = (0.02, 0.18, 0.45, 0.77)
    elif pos == "below":
        frac = (0.04, 0.55, 0.92, 0.40)
    else:  # inline
        frac = (0.15, 0.27, 0.70, 0.60)
    return (
        int(slide_width_emu * frac[0]),
        int(slide_height_emu * frac[1]),
        int(slide_width_emu * frac[2]),
        int(slide_height_emu * frac[3]),
    )


def _resolve_color_override(extra):
    """Parse extra.config.color_override (#RRGGBB) into an RGBColor or None."""
    cfg = extra.config or {} if hasattr(extra, "config") else {}
    if not isinstance(cfg, dict):
        return None
    raw = (cfg.get("color_override") or "").strip()
    if not raw:
        return None
    s = raw.lstrip("#")
    if len(s) != 6:
        return None
    try:
        from pptx.dml.color import RGBColor
        return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except (ValueError, ImportError):
        return None


def _find_placeholder(slide, types_in_priority_order, skip=None):
    """Return the first placeholder matching any type in priority order, optionally skipping shapes."""
    skip = skip or []
    for ptype in types_in_priority_order:
        for p in slide.placeholders:
            if p in skip:
                continue
            if p.placeholder_format.type == ptype:
                return p
    return None


def _fill_text_frame_with_bullets(tf, items: list[str]):
    if not items:
        return
    tf.text = items[0]
    for item in items[1:]:
        p = tf.add_paragraph()
        p.text = item


# ============================================================================
# JSON MODE — renders blank slides with absolute positioning. Fallback when
# no master.pptx is provided.
# ============================================================================

class JSONRenderer:
    def __init__(self, template: Template, catalog=None):
        self.tpl = template
        self.catalog = catalog
        c = template.colors
        self.color_primary = hex_to_rgb(c.primary)
        self.color_accent = hex_to_rgb(c.accent)
        self.color_text = hex_to_rgb(c.text)
        self.color_muted = hex_to_rgb(c.muted)

    def render(self, deck: SlideDeck, output_path: Path) -> None:
        prs = Presentation()
        prs.slide_width = Inches(self.tpl.slide_width_inches)
        prs.slide_height = Inches(self.tpl.slide_height_inches)

        for slide in deck.slides:
            self._render_slide(prs, slide, deck)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

    def _has_right_extra(self, slide: Slide) -> bool:
        return any(
            e.position == "right" and e.type in ("library_asset", "chart_bar", "chart_donut", "matrix_2x2")
            for e in (slide.extras or [])
        )

    def _block_rect(self, slide_data: Slide | None, name: str):
        """Return (left, top, width, height) in EMU for a block whose
        `block_bbox[name]` is set, else None. The wireframe stage stores
        bboxes as {left,top,width,height} normalised 0..1 over the slide's
        16:9 area; we project that onto the actual slide width/height in EMU.
        """
        if slide_data is None:
            return None
        bb_map = getattr(slide_data, "block_bbox", None) or {}
        bb = bb_map.get(name)
        if not bb:
            return None
        slide_w = int(self.tpl.slide_width_inches * 914400)
        slide_h = int(self.tpl.slide_height_inches * 914400)
        try:
            return (
                int(float(bb["left"])   * slide_w),
                int(float(bb["top"])    * slide_h),
                int(float(bb["width"])  * slide_w),
                int(float(bb["height"]) * slide_h),
            )
        except (KeyError, TypeError, ValueError):
            return None

    def _render_slide(self, prs: Presentation, slide: Slide, deck: SlideDeck) -> None:
        # Asset reference — copy the library slide directly (preserves native shapes).
        if slide.asset_ref and self.catalog:
            asset = self.catalog.find(slide.asset_ref)
            if asset is not None:
                copy_slide_into(Path(asset.file_path), asset.slide_index, prs, match_dimensions=False)
                if slide.speaker_notes:
                    new_slide = prs.slides[-1]
                    new_slide.notes_slide.notes_text_frame.text = slide.speaker_notes
                return

        blank = prs.slide_layouts[6]
        s = prs.slides.add_slide(blank)

        if slide.layout == "title":
            self._render_title_slide(s, deck, slide)
        elif slide.layout == "section_divider":
            self._render_section_divider(s, slide)
        elif slide.layout == "two_column":
            self._render_two_column(s, slide)
        elif slide.layout == "key_takeaway":
            self._render_key_takeaway(s, slide)
        else:
            self._render_content(s, slide)

        # Render any growable extras AFTER layout content.
        self._render_extras(s, slide)

        if slide.speaker_notes:
            s.notes_slide.notes_text_frame.text = slide.speaker_notes

    def _render_extras(self, slide, slide_data: Slide) -> None:
        """Stamp each extra onto the existing slide. For now: library_asset only;
        charts/matrices coming in v1.5+. Position dictates where on the slide."""
        slide_w = self.tpl.slide_width_inches * 914400
        slide_h = self.tpl.slide_height_inches * 914400
        for extra in (slide_data.extras or []):
            if extra.type != "library_asset":
                # chart_bar / chart_donut / matrix_2x2 land in next sprint
                continue
            asset_ref = (extra.config or {}).get("asset_ref")
            if not asset_ref or self.catalog is None:
                continue
            asset = self.catalog.find(asset_ref)
            if asset is None:
                continue
            target = _resolve_extra_target(extra, slide_w, slide_h)
            recolor = _resolve_color_override(extra)
            try:
                copy_slide_shapes_onto(
                    Path(asset.file_path),
                    asset.slide_index,
                    slide,
                    target_left_emu=int(target[0]),
                    target_top_emu=int(target[1]),
                    target_width_emu=int(target[2]),
                    target_height_emu=int(target[3]),
                    recolor_to=recolor,
                )
            except Exception as e:  # noqa: BLE001
                # Don't crash the whole render on one bad extra
                import sys
                print(f"[renderer] failed to stamp extra {asset_ref}: {e}", file=sys.stderr)
                continue

    # Map the model's align string to python-pptx alignment enum.
    _ALIGN_MAP = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }

    def _block_style(self, slide_data: Slide | None, name: str) -> dict:
        """Sparse typography overrides for `name` from `slide.block_style`."""
        if slide_data is None:
            return {}
        bs = getattr(slide_data, "block_style", None) or {}
        return dict(bs.get(name) or {})

    def _apply_block_style_textbox(self, *, font_size, bold, color, align, font_family,
                                    italic=False, style_overrides: dict | None = None):
        """Layer the per-block style overrides on top of layout defaults.
        Returns the resolved kwargs for `_add_textbox`."""
        s = style_overrides or {}
        if s.get("font_family"):
            font_family = s["font_family"]
        if s.get("font_size"):
            try:
                font_size = int(s["font_size"])
            except (TypeError, ValueError):
                pass
        if s.get("color"):
            try:
                color = hex_to_rgb(s["color"])
            except Exception:  # noqa: BLE001
                pass
        if "bold" in s:
            bold = bool(s["bold"])
        if "italic" in s:
            italic = bool(s["italic"])
        if s.get("align") in self._ALIGN_MAP:
            align = self._ALIGN_MAP[s["align"]]
        return {
            "font_size": font_size, "bold": bold, "color": color, "align": align,
            "font_family": font_family, "italic": italic,
        }

    def _add_textbox(self, slide, left, top, width, height, text, *, font_size, bold=False,
                     color=None, align=PP_ALIGN.LEFT, font_family=None, italic=False):
        if color is None:
            color = self.color_text
        if font_family is None:
            font_family = self.tpl.fonts.body
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font_family
        return box

    def _add_bullets(self, slide, left, top, width, height, items, *, marker="•",
                     font_size=None, color=None, font_family=None, bold=False, italic=False,
                     align=PP_ALIGN.LEFT):
        if font_size is None:
            font_size = self.tpl.fonts.sizes.body
        if color is None:
            color = self.color_text
        if font_family is None:
            font_family = self.tpl.fonts.body
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"{marker} {item}"
            p.alignment = align
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.italic = italic
            p.font.color.rgb = color
            p.font.name = font_family
            p.space_after = Pt(8)

    # ----- per-layout renderers ------------------------------------------
    # Each block (title / strap / body / body_left / body_right) is placed at
    # `slide.block_bbox[name]` if set, else at the layout's default Inches()
    # rect. This means the wireframe's freeform positioning round-trips into
    # the .pptx export exactly. Auto-promoted bboxes (set by the wireframe
    # save endpoint when switching to a no-body layout) cause body content to
    # appear in hi-fi where the layout would otherwise drop it.

    def _render_title_slide(self, slide, deck, slide_data=None):
        sizes = self.tpl.fonts.sizes
        title_rect = self._block_rect(slide_data, "title") or (
            Inches(1), Inches(2.8), Inches(11.3), Inches(1.5)
        )
        title_kw = self._apply_block_style_textbox(
            font_size=sizes.title_slide_title, bold=True, color=self.color_primary,
            align=PP_ALIGN.LEFT, font_family=self.tpl.fonts.heading,
            style_overrides=self._block_style(slide_data, "title"),
        )
        self._add_textbox(slide, *title_rect, deck.title, **title_kw)
        if deck.subtitle:
            strap_rect = self._block_rect(slide_data, "strap") or (
                Inches(1), Inches(4.4), Inches(11.3), Inches(0.8)
            )
            strap_kw = self._apply_block_style_textbox(
                font_size=sizes.title_slide_subtitle, bold=False, color=self.color_muted,
                align=PP_ALIGN.LEFT, font_family=self.tpl.fonts.body,
                style_overrides=self._block_style(slide_data, "strap"),
            )
            self._add_textbox(slide, *strap_rect, deck.subtitle, **strap_kw)
        body_rect = self._block_rect(slide_data, "body")
        if body_rect and slide_data is not None and slide_data.body:
            body_kw = self._apply_block_style_textbox(
                font_size=self.tpl.fonts.sizes.body, bold=False, color=self.color_text,
                align=PP_ALIGN.LEFT, font_family=self.tpl.fonts.body,
                style_overrides=self._block_style(slide_data, "body"),
            )
            self._add_bullets(slide, *body_rect, slide_data.body, **body_kw)

    def _render_section_divider(self, slide, slide_data):
        title_rect = self._block_rect(slide_data, "title") or (
            Inches(1), Inches(3.2), Inches(11.3), Inches(1.2)
        )
        title_kw = self._apply_block_style_textbox(
            font_size=self.tpl.fonts.sizes.section_divider, bold=True, color=self.color_primary,
            align=PP_ALIGN.LEFT, font_family=self.tpl.fonts.heading,
            style_overrides=self._block_style(slide_data, "title"),
        )
        self._add_textbox(slide, *title_rect, slide_data.title, **title_kw)
        if slide_data.strap:
            strap_rect = self._block_rect(slide_data, "strap") or (
                Inches(1), Inches(4.6), Inches(11.3), Inches(0.7)
            )
            strap_kw = self._apply_block_style_textbox(
                font_size=14, bold=False, color=self.color_muted,
                align=PP_ALIGN.LEFT, font_family=self.tpl.fonts.body,
                style_overrides=self._block_style(slide_data, "strap"),
            )
            self._add_textbox(slide, *strap_rect, slide_data.strap, **strap_kw)
        body_rect = self._block_rect(slide_data, "body")
        if body_rect and slide_data.body:
            body_kw = self._apply_block_style_textbox(
                font_size=self.tpl.fonts.sizes.body, bold=False, color=self.color_text,
                align=PP_ALIGN.LEFT, font_family=self.tpl.fonts.body,
                style_overrides=self._block_style(slide_data, "body"),
            )
            self._add_bullets(slide, *body_rect, slide_data.body, **body_kw)

    def _render_content(self, slide, slide_data):
        title_rect = self._block_rect(slide_data, "title") or (
            Inches(0.6), Inches(0.5), Inches(12.1), Inches(0.8)
        )
        title_kw = self._apply_block_style_textbox(
            font_size=self.tpl.fonts.sizes.slide_title, bold=True, color=self.color_primary,
            align=PP_ALIGN.LEFT, font_family=self.tpl.fonts.heading,
            style_overrides=self._block_style(slide_data, "title"),
        )
        self._add_textbox(slide, *title_rect, slide_data.title, **title_kw)
        body_top = Inches(1.7)
        if slide_data.strap:
            strap_rect = self._block_rect(slide_data, "strap") or (
                Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.5)
            )
            strap_kw = self._apply_block_style_textbox(
                font_size=14, bold=False, color=self.color_muted,
                align=PP_ALIGN.LEFT, font_family=self.tpl.fonts.body,
                style_overrides=self._block_style(slide_data, "strap"),
            )
            self._add_textbox(slide, *strap_rect, slide_data.strap, **strap_kw)
            if self._block_rect(slide_data, "strap") is None:
                body_top = Inches(2.0)
        if slide_data.body:
            marker = "→" if slide_data.layout == "closing" else "•"
            body_rect = self._block_rect(slide_data, "body")
            body_kw = self._apply_block_style_textbox(
                font_size=self.tpl.fonts.sizes.body, bold=False, color=self.color_text,
                align=PP_ALIGN.LEFT, font_family=self.tpl.fonts.body,
                style_overrides=self._block_style(slide_data, "body"),
            )
            if body_rect:
                self._add_bullets(slide, *body_rect, slide_data.body, marker=marker, **body_kw)
            else:
                body_width = Inches(5.7) if self._has_right_extra(slide_data) else Inches(12.1)
                self._add_bullets(slide, Inches(0.6), body_top, body_width, Inches(5.5),
                                  slide_data.body, marker=marker, **body_kw)

    def _render_two_column(self, slide, slide_data):
        title_rect = self._block_rect(slide_data, "title") or (
            Inches(0.6), Inches(0.5), Inches(12.1), Inches(0.8)
        )
        title_kw = self._apply_block_style_textbox(
            font_size=self.tpl.fonts.sizes.slide_title, bold=True, color=self.color_primary,
            align=PP_ALIGN.LEFT, font_family=self.tpl.fonts.heading,
            style_overrides=self._block_style(slide_data, "title"),
        )
        self._add_textbox(slide, *title_rect, slide_data.title, **title_kw)
        body_top = Inches(1.7)
        if slide_data.strap:
            strap_rect = self._block_rect(slide_data, "strap") or (
                Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.5)
            )
            strap_kw = self._apply_block_style_textbox(
                font_size=14, bold=False, color=self.color_muted,
                align=PP_ALIGN.LEFT, font_family=self.tpl.fonts.body,
                style_overrides=self._block_style(slide_data, "strap"),
            )
            self._add_textbox(slide, *strap_rect, slide_data.strap, **strap_kw)
            if self._block_rect(slide_data, "strap") is None:
                body_top = Inches(2.0)
        if slide_data.body_left:
            left_rect = self._block_rect(slide_data, "body_left") or (
                Inches(0.6), body_top, Inches(5.9), Inches(5.5)
            )
            left_kw = self._apply_block_style_textbox(
                font_size=self.tpl.fonts.sizes.body, bold=False, color=self.color_text,
                align=PP_ALIGN.LEFT, font_family=self.tpl.fonts.body,
                style_overrides=self._block_style(slide_data, "body_left"),
            )
            self._add_bullets(slide, *left_rect, slide_data.body_left, **left_kw)
        if slide_data.body_right:
            right_rect = self._block_rect(slide_data, "body_right") or (
                Inches(6.85), body_top, Inches(5.9), Inches(5.5)
            )
            right_kw = self._apply_block_style_textbox(
                font_size=self.tpl.fonts.sizes.body, bold=False, color=self.color_text,
                align=PP_ALIGN.LEFT, font_family=self.tpl.fonts.body,
                style_overrides=self._block_style(slide_data, "body_right"),
            )
            self._add_bullets(slide, *right_rect, slide_data.body_right, **right_kw)
        body_rect = self._block_rect(slide_data, "body")
        if body_rect and slide_data.body:
            body_kw = self._apply_block_style_textbox(
                font_size=self.tpl.fonts.sizes.body, bold=False, color=self.color_text,
                align=PP_ALIGN.LEFT, font_family=self.tpl.fonts.body,
                style_overrides=self._block_style(slide_data, "body"),
            )
            self._add_bullets(slide, *body_rect, slide_data.body, **body_kw)

    def _render_key_takeaway(self, slide, slide_data):
        sizes = self.tpl.fonts.sizes
        title_rect = self._block_rect(slide_data, "title") or (
            Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.5)
        )
        title_kw = self._apply_block_style_textbox(
            font_size=sizes.key_takeaway_title, bold=True, color=self.color_primary,
            align=PP_ALIGN.CENTER, font_family=self.tpl.fonts.heading,
            style_overrides=self._block_style(slide_data, "title"),
        )
        self._add_textbox(slide, *title_rect, slide_data.title, **title_kw)
        if slide_data.strap:
            strap_rect = self._block_rect(slide_data, "strap") or (
                Inches(1.5), Inches(4.7), Inches(10.3), Inches(0.6)
            )
            strap_kw = self._apply_block_style_textbox(
                font_size=14, bold=False, color=self.color_muted,
                align=PP_ALIGN.CENTER, font_family=self.tpl.fonts.body,
                style_overrides=self._block_style(slide_data, "strap"),
            )
            self._add_textbox(slide, *strap_rect, slide_data.strap, **strap_kw)
        if slide_data.body:
            body_rect = self._block_rect(slide_data, "body") or (
                Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.5)
            )
            body_kw = self._apply_block_style_textbox(
                font_size=sizes.key_takeaway_body, bold=False, color=self.color_muted,
                align=PP_ALIGN.CENTER, font_family=self.tpl.fonts.body,
                style_overrides=self._block_style(slide_data, "body"),
            )
            self._add_textbox(slide, *body_rect, " — ".join(slide_data.body), **body_kw)
